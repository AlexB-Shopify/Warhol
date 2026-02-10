"""Tests for the PPTX engine operations."""


import pytest
from pptx import Presentation
from pptx.util import Inches


class TestSlideOperations:
    def test_create_presentation(self):
        from src.pptx_engine.slide_operations import create_presentation

        prs = create_presentation()
        assert prs.slide_width == Inches(13.333)
        assert prs.slide_height == Inches(7.5)

    def test_add_blank_slide(self):
        from src.pptx_engine.slide_operations import add_blank_slide, create_presentation

        prs = create_presentation()
        add_blank_slide(prs)
        assert len(prs.slides) == 1

    def test_clone_from_template(self, tmp_path):
        from src.pptx_engine.slide_operations import (
            clone_slide_from_template,
            create_presentation,
        )

        # Create a source template with content
        template_prs = Presentation()
        template_prs.slide_width = Inches(13.333)
        template_prs.slide_height = Inches(7.5)
        slide = template_prs.slides.add_slide(template_prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        txBox.text_frame.paragraphs[0].text = "Template Content"

        template_path = tmp_path / "template.pptx"
        template_prs.save(str(template_path))

        # Clone into new presentation
        target = create_presentation()
        clone_slide_from_template(target, template_path, 0)
        assert len(target.slides) == 1

    def test_clone_invalid_index(self, tmp_path):
        from src.pptx_engine.slide_operations import (
            clone_slide_from_template,
            create_presentation,
        )

        template_prs = Presentation()
        template_prs.slides.add_slide(template_prs.slide_layouts[6])
        template_path = tmp_path / "template.pptx"
        template_prs.save(str(template_path))

        target = create_presentation()
        with pytest.raises(IndexError):
            clone_slide_from_template(target, template_path, 5)


class TestTextOperations:
    def test_add_textbox(self):
        from src.pptx_engine.slide_operations import add_blank_slide, create_presentation
        from src.pptx_engine.text_operations import add_textbox

        prs = create_presentation()
        slide = add_blank_slide(prs)

        txBox = add_textbox(
            slide, "Hello World",
            left=1.0, top=1.0, width=8.0, height=1.0,
            font_name="Arial", font_size=24,
        )
        assert txBox is not None
        assert txBox.text_frame.paragraphs[0].text == "Hello World"

    def test_add_bullet_list(self):
        from src.pptx_engine.slide_operations import add_blank_slide, create_presentation
        from src.pptx_engine.text_operations import add_bullet_list

        prs = create_presentation()
        slide = add_blank_slide(prs)

        txBox = add_bullet_list(
            slide, ["Item 1", "Item 2", "Item 3"],
            left=1.0, top=2.0, width=10.0, height=4.0,
        )
        tf = txBox.text_frame
        assert len(tf.paragraphs) == 3
        assert tf.paragraphs[0].text == "Item 1"
        assert tf.paragraphs[2].text == "Item 3"

    def test_add_textbox_with_color(self):
        from src.pptx_engine.slide_operations import add_blank_slide, create_presentation
        from src.pptx_engine.text_operations import add_textbox

        prs = create_presentation()
        slide = add_blank_slide(prs)

        txBox = add_textbox(
            slide, "Colored Text",
            left=1.0, top=1.0, width=8.0, height=1.0,
            font_color="#ff0000", bold=True,
        )
        assert txBox is not None


class TestShapeOperations:
    def test_add_rectangle(self):
        from src.pptx_engine.slide_operations import add_blank_slide, create_presentation
        from src.pptx_engine.shape_operations import add_rectangle

        prs = create_presentation()
        slide = add_blank_slide(prs)

        shape = add_rectangle(
            slide, 1.0, 1.0, 4.0, 2.0,
            fill_color="#1a73e8",
        )
        assert shape is not None

    def test_add_rectangle_no_fill(self):
        from src.pptx_engine.slide_operations import add_blank_slide, create_presentation
        from src.pptx_engine.shape_operations import add_rectangle

        prs = create_presentation()
        slide = add_blank_slide(prs)

        shape = add_rectangle(
            slide, 1.0, 1.0, 4.0, 2.0,
            border_color="#333333",
        )
        assert shape is not None

    def test_add_oval(self):
        from src.pptx_engine.slide_operations import add_blank_slide, create_presentation
        from src.pptx_engine.shape_operations import add_oval

        prs = create_presentation()
        slide = add_blank_slide(prs)

        shape = add_oval(
            slide, 2.0, 2.0, 3.0, 3.0,
            fill_color="#34a853",
        )
        assert shape is not None


class TestEndToEnd:
    def test_build_and_save(self, tmp_path):
        """Test building a complete presentation and saving it."""
        from src.pptx_engine.slide_operations import add_blank_slide, create_presentation
        from src.pptx_engine.text_operations import add_bullet_list, add_textbox
        from src.pptx_engine.shape_operations import add_rectangle

        prs = create_presentation()

        # Title slide
        slide = add_blank_slide(prs)
        add_rectangle(slide, 0, 0, 13.333, 0.15, fill_color="#1a73e8")
        add_textbox(
            slide, "Test Presentation",
            left=1.5, top=2.0, width=10.333, height=1.5,
            font_size=44, bold=True, alignment="center",
        )

        # Content slide
        slide = add_blank_slide(prs)
        add_textbox(
            slide, "Key Points",
            left=0.8, top=0.4, width=11.733, height=0.8,
            font_size=32, bold=True,
        )
        add_bullet_list(
            slide, ["Point One", "Point Two", "Point Three"],
            left=1.0, top=1.6, width=11.333, height=5.0,
        )

        output = tmp_path / "test_output.pptx"
        prs.save(str(output))

        assert output.exists()
        assert output.stat().st_size > 0

        # Verify we can reopen it
        reopened = Presentation(str(output))
        assert len(reopened.slides) == 2
