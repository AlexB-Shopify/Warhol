"""Tests for the design system extractor."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


class TestDesignSystemExtractor:
    def test_extract_from_basic_template(self):
        """Test extracting a design system from the built-in basic template."""
        from src.pptx_engine.design_system_extractor import extract_design_system

        template_dir = Path(__file__).parent.parent / "templates" / "minimal"
        if not (template_dir / "basic.pptx").exists():
            import pytest
            pytest.skip("Basic template not generated yet")

        ds = extract_design_system(template_dir, name="Test")
        assert ds.name == "Test"
        # The basic template uses Arial everywhere
        assert ds.fonts.title_font == "Arial"
        assert ds.fonts.body_font == "Arial"
        # Title size in basic template is 44pt
        assert ds.fonts.title_size == 44

    def test_extract_from_generated_pptx(self, tmp_path):
        """Test extraction from a custom-built .pptx with known fonts/colors."""
        from src.pptx_engine.design_system_extractor import extract_design_system

        # Create a presentation with known design values
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for _ in range(3):  # Multiple slides to build frequency
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # Title with Helvetica 36pt
            title_box = slide.shapes.add_textbox(
                Inches(1), Inches(0.5), Inches(10), Inches(1.2)
            )
            title_box.name = "Title 1"
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = "Title Text"
            p.font.name = "Helvetica"
            p.font.size = Pt(36)
            p.font.color.rgb = RGBColor(0x11, 0x11, 0x11)

            # Body with Georgia 16pt
            body_box = slide.shapes.add_textbox(
                Inches(1), Inches(2), Inches(10), Inches(4)
            )
            tf = body_box.text_frame
            p = tf.paragraphs[0]
            p.text = "Body text here"
            p.font.name = "Georgia"
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

            # Accent shape with brand color
            from pptx.enum.shapes import MSO_SHAPE
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0), Inches(13.333), Inches(0.1)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0xE9, 0x1E, 0x63)  # Pink

        prs.save(str(tmp_path / "brand.pptx"))

        ds = extract_design_system(tmp_path, name="Brand")

        assert ds.name == "Brand"
        assert ds.fonts.title_font == "Helvetica"
        assert ds.fonts.body_font == "Georgia"
        assert ds.fonts.title_size == 36
        assert ds.fonts.body_size == 16
        # The pink accent should be picked up as primary
        assert ds.colors.primary.upper() == "#E91E63"

    def test_extract_empty_directory(self, tmp_path):
        """Test graceful handling of empty directory."""
        from src.pptx_engine.design_system_extractor import extract_design_system

        ds = extract_design_system(tmp_path)
        # Should return defaults
        assert ds.fonts.title_font == "Arial"
        assert ds.colors.primary == "#1a73e8"

    def test_extract_from_single_file(self, tmp_path):
        """Test extracting a design system from a single .pptx file."""
        from src.pptx_engine.design_system_extractor import extract_design_system_from_file

        # Create a minimal pptx with known fonts
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(1))
        tb.name = "Title 1"
        p = tb.text_frame.paragraphs[0]
        p.text = "Title"
        p.font.name = "Courier New"
        p.font.size = Pt(40)

        single_file = tmp_path / "single.pptx"
        prs.save(str(single_file))

        ds = extract_design_system_from_file(single_file, name="Single")
        assert ds.name == "Single"
        assert ds.fonts.title_font == "Courier New"
        assert ds.fonts.title_size == 40

    def test_extract_from_single_file_not_found(self, tmp_path):
        """Test that missing file raises FileNotFoundError."""
        from src.pptx_engine.design_system_extractor import extract_design_system_from_file
        import pytest

        with pytest.raises(FileNotFoundError):
            extract_design_system_from_file(tmp_path / "nonexistent.pptx")

    def test_extract_from_single_file_wrong_extension(self, tmp_path):
        """Test that non-.pptx file raises ValueError."""
        from src.pptx_engine.design_system_extractor import extract_design_system_from_file
        import pytest

        bad_file = tmp_path / "notes.txt"
        bad_file.write_text("not a pptx")

        with pytest.raises(ValueError, match="Expected a .pptx file"):
            extract_design_system_from_file(bad_file)

    def test_extract_and_save_yaml(self, tmp_path):
        """Test the full round-trip: extract -> save -> load."""
        from src.pptx_engine.design_system_extractor import extract_design_system
        from src.schemas.design_system import DesignSystem

        template_dir = Path(__file__).parent.parent / "templates" / "minimal"
        if not (template_dir / "basic.pptx").exists():
            import pytest
            pytest.skip("Basic template not generated yet")

        ds = extract_design_system(template_dir)
        yaml_path = tmp_path / "extracted.yaml"
        ds.to_yaml(yaml_path)

        loaded = DesignSystem.from_yaml(yaml_path)
        assert loaded.fonts.title_font == ds.fonts.title_font
        assert loaded.colors.primary == ds.colors.primary
