"""Tests for document parsers."""

from pathlib import Path

import pytest


class TestTextParser:
    def test_parse_plain_text(self, tmp_path):
        from src.parsers.text_parser import parse_text

        f = tmp_path / "test.txt"
        f.write_text("Hello World\nThis is a test.")
        result = parse_text(f)
        assert "Hello World" in result
        assert "This is a test." in result

    def test_parse_markdown(self, tmp_path):
        from src.parsers.text_parser import parse_text

        f = tmp_path / "test.md"
        f.write_text("# Heading\n\n- Bullet 1\n- Bullet 2")
        result = parse_text(f)
        assert "# Heading" in result
        assert "- Bullet 1" in result

    def test_normalize_blank_lines(self, tmp_path):
        from src.parsers.text_parser import parse_text

        f = tmp_path / "test.txt"
        f.write_text("Line 1\n\n\n\n\nLine 2")
        result = parse_text(f)
        # Should collapse to max 2 blank lines
        assert "\n\n\n\n" not in result
        assert "Line 1" in result
        assert "Line 2" in result


class TestParserDispatch:
    def test_dispatch_txt(self, tmp_path):
        from src.parsers import parse

        f = tmp_path / "test.txt"
        f.write_text("Hello")
        result = parse(f)
        assert result == "Hello"

    def test_dispatch_md(self, tmp_path):
        from src.parsers import parse

        f = tmp_path / "test.md"
        f.write_text("# Hello")
        result = parse(f)
        assert "# Hello" in result

    def test_unsupported_format(self, tmp_path):
        from src.parsers import parse

        f = tmp_path / "test.xyz"
        f.write_text("data")
        with pytest.raises(ValueError, match="Unsupported file format"):
            parse(f)


class TestPptxParser:
    def test_parse_basic_template(self):
        """Test parsing the built-in basic template."""
        from src.parsers.pptx_parser import parse_pptx

        template = Path(__file__).parent.parent / "templates" / "minimal" / "basic.pptx"
        if not template.exists():
            pytest.skip("Basic template not generated yet")

        result = parse_pptx(template)
        assert "Presentation Title" in result
        assert "Slide" in result

    def test_parse_generated_pptx(self, tmp_path):
        """Test parsing a dynamically generated .pptx file."""
        from pptx import Presentation
        from pptx.util import Inches

        from src.parsers.pptx_parser import parse_pptx

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        txBox.text_frame.paragraphs[0].text = "Test Content Here"

        path = tmp_path / "test.pptx"
        prs.save(str(path))

        result = parse_pptx(path)
        assert "Test Content Here" in result


class TestDocxParser:
    def test_parse_simple_docx(self, tmp_path):
        """Test parsing a simple DOCX file."""
        from docx import Document

        from src.parsers.docx_parser import parse_docx

        doc = Document()
        doc.add_heading("Test Heading", level=1)
        doc.add_paragraph("This is body text.")
        doc.add_paragraph("Bullet item", style="List Bullet")

        path = tmp_path / "test.docx"
        doc.save(str(path))

        result = parse_docx(path)
        assert "# Test Heading" in result
        assert "This is body text." in result
