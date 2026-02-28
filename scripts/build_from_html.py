#!/usr/bin/env python3
"""Build a PowerPoint presentation from an HTML slide deck preview.

The HTML file (produced by render_html.py or the agent directly) is the
**visual contract** — every slide is a <div class="slide"> with
absolute-positioned text elements, explicit fonts/colors, and data
attributes that tell this builder exactly how to construct the PPTX.

Two build modes (set via data-build-mode on the slide div):

  1. **clone** — Clone a template slide from a source PPTX, preserving all
     branded visuals.  Replace text ONLY in named shapes (data-shape-name).
     Clear any remaining text shapes that were not targeted.  No new
     textboxes are created.

  2. **compose** — Create a slide from a branded base-template layout
     (master backgrounds are inherited).  Build all text elements as new
     textboxes from scratch.  No cloned content shapes.

Usage:
    python scripts/build_from_html.py workspace/deck_preview.html \\
        -o output.pptx \\
        --base-template "templates/base/Shopify - Example Technical Workshop Slide Bank.pptx"
"""

import argparse
import json
import logging
import re
import sys
from pathlib import Path

# Add project root to path
sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from bs4 import BeautifulSoup, Tag

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt, Emu

from src.pptx_engine.slide_operations import (
    add_blank_slide,
    add_slide_from_layout,
    clear_clone_caches,
    clone_slide_as_is,
    open_base_template,
    create_presentation,
)
from src.pptx_engine.text_operations import (
    add_textbox,
    add_bullet_list,
    add_accent_bar,
    add_image_placeholder,
    estimate_fit_font_size,
    set_autofit_shrink,
)
from src.schemas.html_schema import DPI

logger = logging.getLogger(__name__)
logging.basicConfig(level=logging.INFO, format="%(levelname)s: %(message)s")


# ---------------------------------------------------------------------------
# CSS variable resolution
# ---------------------------------------------------------------------------

def _parse_css_variables(soup: BeautifulSoup) -> dict[str, str]:
    """Extract CSS custom properties from the :root block in <style> tags.

    Returns a dict mapping variable names (e.g., '--color-primary') to their
    resolved literal values (e.g., '#CDF986').
    """
    var_map: dict[str, str] = {}
    for style_tag in soup.find_all("style"):
        css_text = style_tag.string or ""
        # Find :root { ... } blocks
        root_match = re.search(r":root\s*\{([^}]+)\}", css_text, re.DOTALL)
        if not root_match:
            continue
        root_block = root_match.group(1)
        # Parse each --variable: value; declaration
        for m in re.finditer(r"(--[\w-]+)\s*:\s*([^;]+);", root_block):
            name = m.group(1).strip()
            value = m.group(2).strip()
            var_map[name] = value
    return var_map


def _resolve_vars(style: str, var_map: dict[str, str]) -> str:
    """Replace all var(--token) and var(--token, fallback) in a style string
    with their literal values from var_map.
    """
    if "var(" not in style:
        return style

    def _replace(m: re.Match) -> str:
        token = m.group(1).strip()
        fallback = m.group(2)
        if token in var_map:
            return var_map[token]
        if fallback is not None:
            return fallback.strip()
        return m.group(0)  # leave unchanged if not found

    # Match var(--token) and var(--token, fallback)
    resolved = re.sub(
        r"var\(\s*(--[\w-]+)\s*(?:,\s*([^)]+))?\s*\)",
        _replace,
        style,
    )
    # Handle nested var() (one more pass)
    if "var(" in resolved:
        resolved = re.sub(
            r"var\(\s*(--[\w-]+)\s*(?:,\s*([^)]+))?\s*\)",
            _replace,
            resolved,
        )
    return resolved


# ---------------------------------------------------------------------------
# HTML parsing helpers
# ---------------------------------------------------------------------------

def _px_to_inches(px: float) -> float:
    """Convert pixels (at 96 DPI) to inches."""
    return px / DPI


def _parse_css_value(style: str, prop: str, default: float = 0.0) -> float:
    """Extract a numeric CSS property value from an inline style string.

    Handles px and pt units. Returns the number stripped of units.
    """
    pattern = rf"{prop}\s*:\s*([\d.]+)\s*(px|pt|em|rem|%)?"
    match = re.search(pattern, style)
    if not match:
        return default
    return float(match.group(1))


def _parse_css_color(style: str, prop: str = "color", default: str = "#000000") -> str:
    """Extract a hex color from a CSS style string.

    Handles #RRGGBB, #RGB shorthand, rgb(), and rgba() notations.
    """
    # Try 6-digit hex first
    pattern = rf"{prop}\s*:\s*(#[0-9A-Fa-f]{{6}})"
    match = re.search(pattern, style)
    if match:
        return match.group(1)

    # Try 3-digit hex shorthand (#RGB -> #RRGGBB)
    short_pattern = rf"{prop}\s*:\s*#([0-9A-Fa-f])([0-9A-Fa-f])([0-9A-Fa-f])(?![0-9A-Fa-f])"
    short_match = re.search(short_pattern, style)
    if short_match:
        r, g, b = short_match.group(1), short_match.group(2), short_match.group(3)
        return f"#{r}{r}{g}{g}{b}{b}"

    # Try rgb(r, g, b) or rgba(r, g, b, a)
    rgba_pattern = rf"{prop}\s*:\s*rgba?\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)"
    rgba_match = re.search(rgba_pattern, style)
    if rgba_match:
        r = min(255, int(rgba_match.group(1)))
        g = min(255, int(rgba_match.group(2)))
        b = min(255, int(rgba_match.group(3)))
        return f"#{r:02X}{g:02X}{b:02X}"

    return default


def _parse_css_string(style: str, prop: str, default: str = "") -> str:
    """Extract a string CSS property value (e.g., font-family)."""
    pattern = rf"{prop}\s*:\s*'([^']+)'"
    match = re.search(pattern, style)
    if not match:
        # Try without quotes
        pattern = rf"{prop}\s*:\s*([^;]+)"
        match = re.search(pattern, style)
        if not match:
            return default
        return match.group(1).strip().strip("'\"")
    return match.group(1)


def _parse_css_bool(style: str, prop: str, value: str) -> bool:
    """Check if a CSS property has a specific value."""
    pattern = rf"{prop}\s*:\s*{value}"
    return bool(re.search(pattern, style, re.IGNORECASE))


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor."""
    hex_color = hex_color.lstrip("#")
    if len(hex_color) > 6:
        hex_color = hex_color[:6]
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )


# ---------------------------------------------------------------------------
# Slide building from HTML sections
# ---------------------------------------------------------------------------

def _build_slide_from_section(
    prs: Presentation,
    section: Tag,
    var_map: dict[str, str] | None = None,
) -> str:
    """Build a single PPTX slide from an HTML <div class="slide"> section.

    Three paths:
      - "clone" build-mode: clone template slide, replace text in named shapes
      - template_clone bg-type (compose): clone for background, clear text, add new textboxes
      - plain compose: create from branded layout, add all content as new textboxes

    Returns the build mode used ("clone", "hybrid", or "compose").
    """
    vmap = var_map or {}
    slide_num = section.get("data-slide-number", "?")
    slide_type = section.get("data-slide-type", "content")
    build_mode = section.get("data-build-mode", "compose")

    bg_div = section.find("div", class_="slide-bg")
    bg_type = bg_div.get("data-bg-type", "layout") if bg_div else "layout"

    slide = None

    # =================================================================
    # CLONE MODE — paste template slide, replace text in named shapes
    # =================================================================
    if build_mode == "clone" and bg_type == "template_clone" and bg_div:
        template_file = bg_div.get("data-template-file", "")
        slide_index_str = bg_div.get("data-slide-index", "0")
        slide_index = int(slide_index_str) if slide_index_str else 0

        if template_file and Path(template_file).exists():
            try:
                slide = clone_slide_as_is(prs, template_file, slide_index)
                logger.info(
                    f"Slide {slide_num}: CLONE from {template_file} "
                    f"index {slide_index}"
                )
            except Exception as e:
                logger.warning(
                    f"Slide {slide_num}: clone failed ({e}), "
                    f"falling back to compose mode"
                )

        if slide is not None:
            # Replace text ONLY in named shapes — no new textboxes
            mapped_shape_names: set[str] = set()

            for elem_div in section.find_all("div", class_="element"):
                shape_name = elem_div.get("data-shape-name", "")
                if not shape_name:
                    logger.debug(
                        f"Slide {slide_num}: clone mode element "
                        f"has no shape_name, skipping"
                    )
                    continue

                style = _resolve_vars(elem_div.get("style", ""), vmap)
                font_family = _parse_css_string(style, "font-family", "Arial")
                font_size_pt = _parse_css_value(style, "font-size", 18)
                font_color = _parse_css_color(style, "color", "#000000")
                is_bold = _parse_css_bool(style, "font-weight", "bold")
                is_italic = _parse_css_bool(style, "font-style", "italic")

                bullet_items = _extract_bullet_items(elem_div)
                text_content = _extract_text_content(elem_div)

                replaced = _replace_shape_text_by_name(
                    slide, shape_name, text_content, bullet_items,
                    font_family, font_size_pt, font_color,
                    is_bold, is_italic,
                )
                if replaced:
                    mapped_shape_names.add(shape_name)
                else:
                    logger.warning(
                        f"Slide {slide_num}: shape '{shape_name}' not found "
                        f"on cloned slide"
                    )

            # Clear stale text from shapes that were NOT targeted
            _clear_unmapped_shapes(slide, mapped_shape_names)

            # Add speaker notes
            notes_div = section.find("div", class_="speaker-notes")
            if notes_div and notes_div.get_text(strip=True):
                _add_speaker_notes(slide, notes_div.get_text(strip=True))

            return "clone"

    # =================================================================
    # HYBRID MODE — clone template for background, clear text, compose
    # =================================================================
    if bg_type == "template_clone" and bg_div:
        template_file = bg_div.get("data-template-file", "")
        slide_index_str = bg_div.get("data-slide-index", "0")
        slide_index = int(slide_index_str) if slide_index_str else 0
        preserve_images = bg_div.get("data-preserve-images", "").lower() == "true"

        if template_file and Path(template_file).exists():
            try:
                slide = clone_slide_as_is(prs, template_file, slide_index)
                _clear_all_text_shapes(slide, preserve_images=preserve_images)
                logger.info(
                    f"Slide {slide_num}: HYBRID clone+compose from "
                    f"{Path(template_file).name} index {slide_index}"
                )
            except Exception as e:
                logger.warning(
                    f"Slide {slide_num}: hybrid clone failed ({e}), "
                    f"falling back to compose"
                )
                slide = None

    # =================================================================
    # COMPOSE MODE — build from branded layout, all new textboxes
    # =================================================================
    if slide is None:
        slide = add_blank_slide(prs)

        # Apply background color if specified
        if bg_div:
            bg_style = _resolve_vars(bg_div.get("style", ""), vmap)
            color_match = re.search(r"background\s*:\s*(#[0-9A-Fa-f]{6})", bg_style)
            if color_match:
                _set_slide_bg_color(slide, color_match.group(1))

        logger.info(f"Slide {slide_num}: COMPOSE from layout ({slide_type})")
    else:
        # Hybrid mode — verify background survived the clone.
        if bg_div:
            bg_style = _resolve_vars(bg_div.get("style", ""), vmap)
            fallback_color = _parse_css_color(bg_style, "background", "")
            if not _verify_slide_background(slide) and fallback_color:
                logger.warning(
                    f"Slide {slide_num}: cloned background missing, "
                    f"applying fallback color {fallback_color}"
                )
                _set_slide_bg_color(slide, fallback_color)

    # Detect and render accent stripe (border-left on .slide)
    slide_style = _resolve_vars(section.get("style", ""), vmap)
    # Also check the CSS class-based border from the stylesheet
    visual_profile = section.get("data-visual-profile", "")
    _maybe_add_accent_stripe(slide, slide_style, visual_profile, vmap)

    # Create all text elements as new textboxes
    for elem_div in section.find_all("div", class_="element"):
        role = elem_div.get("data-role", "")
        if role == "image_placeholder":
            _add_image_placeholder_to_slide(slide, elem_div, vmap)
        else:
            _add_element_to_slide_compose(slide, elem_div, vmap)

    # Create decoration shapes (non-text visual elements)
    for deco_div in section.find_all("div", class_="decoration"):
        _add_decoration_to_slide(slide, deco_div, vmap)

    # Post-build overlap detection: warn if text shapes overlap
    _warn_overlapping_shapes(slide, slide_num)

    # Enforce text-background contrast on all text shapes
    if bg_div:
        resolved_bg_style = _resolve_vars(bg_div.get("style", ""), vmap)
        bg_color_for_contrast = _parse_css_color(resolved_bg_style, "background", "")
        if bg_color_for_contrast:
            _enforce_contrast(slide, bg_color_for_contrast, vmap)

    # Speaker notes
    notes_div = section.find("div", class_="speaker-notes")
    if notes_div and notes_div.get_text(strip=True):
        _add_speaker_notes(slide, notes_div.get_text(strip=True))

    return "hybrid" if bg_type == "template_clone" else "compose"


def _fitted_font_size(
    text: str,
    width_inches: float,
    height_inches: float,
    requested_pt: float,
    font_name: str = "Arial",
) -> int:
    """Return the largest font size that fits, capped at *requested_pt*.

    Uses the heuristic estimator and clamps to a 10pt floor.
    """
    fit_pt = estimate_fit_font_size(
        text, width_inches, height_inches,
        max_font_pt=requested_pt, min_font_pt=10.0,
        font_name=font_name,
    )
    return max(10, int(min(fit_pt, requested_pt)))


def _add_element_to_slide_compose(
    slide, elem_div: Tag, var_map: dict[str, str] | None = None,
) -> None:
    """Add a text element to a compose-mode slide as a new textbox.

    Compose mode always creates new textboxes — there are no existing
    shapes to target.  CSS variables are resolved before parsing.
    Font sizes are checked against the available space and reduced if
    the text would overflow the shape.
    """
    vmap = var_map or {}
    style = _resolve_vars(elem_div.get("style", ""), vmap)

    # Extract position from CSS
    left_px = _parse_css_value(style, "left")
    top_px = _parse_css_value(style, "top")
    width_px = _parse_css_value(style, "width")
    height_px = _parse_css_value(style, "height")

    left = _px_to_inches(left_px)
    top = _px_to_inches(top_px)
    width = _px_to_inches(width_px)
    height = _px_to_inches(height_px)

    if width <= 0 or height <= 0:
        return  # Skip elements with no dimensions

    # Extract font properties from CSS (now resolved from var tokens)
    font_family = _parse_css_string(style, "font-family", "Arial")
    font_size_pt = _parse_css_value(style, "font-size", 18)
    font_color = _parse_css_color(style, "color", "#000000")
    is_bold = _parse_css_bool(style, "font-weight", "bold")
    is_italic = _parse_css_bool(style, "font-style", "italic")
    alignment = _parse_css_string(style, "text-align", "left")
    line_spacing_raw = _parse_css_value(style, "line-height", 0)
    line_spacing = line_spacing_raw if line_spacing_raw > 0 else None

    # Extract content
    bullet_items = _extract_bullet_items(elem_div)
    text_content = _extract_text_content(elem_div)

    if bullet_items:
        full_text = "\n".join(bullet_items)
        fitted_size = _fitted_font_size(full_text, width, height, font_size_pt, font_family)
        shape = add_bullet_list(
            slide,
            items=bullet_items,
            left=left,
            top=top,
            width=width,
            height=height,
            font_name=font_family,
            font_size=fitted_size,
            font_color=font_color,
            line_spacing=line_spacing or 1.2,
        )
    else:
        fitted_size = _fitted_font_size(text_content, width, height, font_size_pt, font_family)
        shape = add_textbox(
            slide,
            text_content,
            left=left,
            top=top,
            width=width,
            height=height,
            font_name=font_family,
            font_size=fitted_size,
            font_color=font_color,
            bold=is_bold,
            italic=is_italic,
            alignment=alignment,
            line_spacing=line_spacing,
        )

    # Safety net: enable PowerPoint's native shrink-on-overflow
    if shape is not None:
        set_autofit_shrink(shape, min_font_scale_pct=60)


def _add_decoration_to_slide(
    slide, deco_div: Tag, var_map: dict[str, str] | None = None,
) -> None:
    """Add a decoration shape to a compose-mode slide.

    Handles CSS classes: numbered-badge, accent-card, divider-line,
    icon-placeholder, image-frame. Each becomes a PPTX shape at the
    specified CSS position.
    """
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN

    vmap = var_map or {}
    style = _resolve_vars(deco_div.get("style", ""), vmap)
    classes = deco_div.get("class", [])
    if isinstance(classes, str):
        classes = classes.split()

    left_px = _parse_css_value(style, "left")
    top_px = _parse_css_value(style, "top")
    width_px = _parse_css_value(style, "width")
    height_px = _parse_css_value(style, "height")

    left = _px_to_inches(left_px)
    top = _px_to_inches(top_px)
    width = _px_to_inches(width_px)
    height = _px_to_inches(height_px)

    if width <= 0 or height <= 0:
        return

    text_content = deco_div.get_text(strip=True)

    if "numbered-badge" in classes:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(left), Inches(top), Inches(width), Inches(height),
        )
        fill_color = _parse_css_color(style, "background", vmap.get("--color-badge-fill", "#CDF986"))
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(fill_color)
        shape.line.fill.background()
        if text_content:
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
            p.text = text_content
            p.alignment = PP_ALIGN.CENTER
            text_color = _parse_css_color(style, "color", vmap.get("--color-badge-text", "#191E17"))
            for run in p.runs:
                run.font.size = Pt(max(int(width * 72 * 0.4), 10))
                run.font.color.rgb = _hex_to_rgb(text_color)
                run.font.bold = True
        return

    if "divider-line" in classes:
        line_color = vmap.get("--color-divider", "#CCCCCC")
        add_accent_bar(slide, left, top, width, max(height, 0.01), line_color)
        return

    if "accent-card" in classes:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height),
        )
        fill_color = _parse_css_color(style, "background", vmap.get("--color-surface", "#F4F4F4"))
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(fill_color)
        shape.line.fill.background()
        # Add accent bar on left edge
        accent_color = vmap.get("--color-accent-bar", vmap.get("--color-primary", "#CDF986"))
        add_accent_bar(slide, left, top, 0.04, height, accent_color)
        if text_content:
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
            p.text = text_content
            text_color = _parse_css_color(style, "color", "#434343")
            for run in p.runs:
                run.font.size = Pt(11)
                run.font.color.rgb = _hex_to_rgb(text_color)
        return

    if "icon-placeholder" in classes:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(left), Inches(top), Inches(width), Inches(height),
        )
        fill_color = _parse_css_color(style, "background", vmap.get("--color-surface", "#F4F4F4"))
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(fill_color)
        shape.line.fill.background()
        if text_content:
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
            p.text = text_content
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(10)
                run.font.color.rgb = _hex_to_rgb(vmap.get("--color-text-light", "#888888"))
        return

    if "image-frame" in classes:
        description = text_content or "Image"
        add_image_placeholder(
            slide,
            left=left,
            top=top,
            width=width,
            height=height,
            description=description,
            fill_hex=_parse_css_color(style, "background", "#F4F4F4"),
            border_hex=vmap.get("--color-divider", "#CCCCCC"),
            text_color_hex=vmap.get("--color-text-light", "#888888"),
        )
        return

    # Fallback: generic colored rectangle
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    fill_color = _parse_css_color(style, "background", "#F4F4F4")
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(fill_color)
    shape.line.fill.background()


def _add_image_placeholder_to_slide(
    slide, elem_div: Tag, var_map: dict[str, str] | None = None,
) -> None:
    """Add an image placeholder shape to a slide.

    Reads data-image-description, data-image-style, and CSS position to
    create a labeled rectangle placeholder in the PPTX.
    """
    vmap = var_map or {}
    style = _resolve_vars(elem_div.get("style", ""), vmap)

    left_px = _parse_css_value(style, "left")
    top_px = _parse_css_value(style, "top")
    width_px = _parse_css_value(style, "width")
    height_px = _parse_css_value(style, "height")

    left = _px_to_inches(left_px)
    top = _px_to_inches(top_px)
    width = _px_to_inches(width_px)
    height = _px_to_inches(height_px)

    if width <= 0 or height <= 0:
        return

    description = elem_div.get("data-image-description", "Image placeholder")
    image_style = elem_div.get("data-image-style", "")

    fill_hex = _parse_css_color(style, "background", "#F4F4F4")
    border_color = vmap.get("--color-text-light", "#CCCCCC")
    text_color = vmap.get("--color-text-light", "#888888")

    add_image_placeholder(
        slide,
        left=left,
        top=top,
        width=width,
        height=height,
        description=description,
        fill_hex=fill_hex,
        text_color_hex=text_color,
        border_hex=border_color,
        image_style=image_style,
    )


def _extract_bullet_items(elem_div: Tag) -> list[str]:
    """Extract bullet items from <li> elements within the div."""
    items = []
    for li in elem_div.find_all("li"):
        text = li.get_text(strip=True)
        if text:
            items.append(text)
    return items


def _extract_text_content(elem_div: Tag) -> str:
    """Extract plain text content from an element div.

    Handles <br> as newlines and strips HTML tags.
    """
    # Replace <br> tags with newlines before getting text
    for br in elem_div.find_all("br"):
        br.replace_with("\n")
    return elem_div.get_text(strip=False).strip()


def _replace_shape_text_by_name(
    slide,
    shape_name: str,
    text: str,
    bullet_items: list[str] | None,
    font_family: str,
    font_size_pt: float,
    font_color: str,
    bold: bool,
    italic: bool,
) -> bool:
    """Try to find a shape by name on the slide and replace its text.

    Returns True if the shape was found and text was replaced.
    """
    for shape in slide.shapes:
        if shape.name == shape_name and shape.has_text_frame:
            tf = shape.text_frame
            if not tf.paragraphs:
                return False

            # Preserve first-run formatting, apply our content
            first_para = tf.paragraphs[0]

            # Clear extra paragraphs
            ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
            p_elements = list(tf._element)
            for p_el in p_elements:
                if p_el.tag == f"{{{ns_a}}}p" and p_el is not first_para._element:
                    tf._element.remove(p_el)

            if bullet_items:
                # Multi-paragraph bullet content
                first_para.text = bullet_items[0]
                _apply_font_to_para(first_para, font_family, font_size_pt, font_color, bold, italic)
                for item in bullet_items[1:]:
                    p = tf.add_paragraph()
                    p.text = item
                    _apply_font_to_para(p, font_family, font_size_pt, font_color, bold, italic)
            else:
                first_para.text = text
                _apply_font_to_para(first_para, font_family, font_size_pt, font_color, bold, italic)

            return True
    return False


def _apply_font_to_para(
    para, font_family: str, font_size_pt: float,
    font_color: str, bold: bool, italic: bool
) -> None:
    """Apply font formatting to all runs in a paragraph."""
    for run in para.runs:
        run.font.name = font_family
        run.font.size = Pt(font_size_pt)
        try:
            run.font.color.rgb = _hex_to_rgb(font_color)
        except Exception:
            pass
        run.font.bold = bold
        run.font.italic = italic


def _set_slide_bg_color(slide, hex_color: str) -> None:
    """Set a solid background color on a slide."""
    try:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = _hex_to_rgb(hex_color)
    except Exception as e:
        logger.debug(f"Could not set background color: {e}")


def _verify_slide_background(slide) -> bool:
    """Check whether a slide has an effective background.

    Returns True if the slide has a slide-level, layout-level, or
    master-level background with actual fill content. Returns False
    if the background appears to be empty/missing.
    """
    ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"

    def _has_fill(bg_elem) -> bool:
        if bg_elem is None or len(bg_elem) == 0:
            return False
        for bgPr in bg_elem.iter(f"{{{ns_p}}}bgPr"):
            if (bgPr.find(f"{{{ns_a}}}solidFill") is not None
                    or bgPr.find(f"{{{ns_a}}}gradFill") is not None
                    or bgPr.find(f"{{{ns_a}}}blipFill") is not None
                    or bgPr.find(f"{{{ns_a}}}pattFill") is not None):
                return True
        return False

    try:
        # Check slide-level background
        if _has_fill(slide.background._element):
            return True

        # Check layout-level background
        layout = slide.slide_layout
        if layout:
            for bg in layout._element.iter(f"{{{ns_p}}}bg"):
                if _has_fill(bg):
                    return True

        # Check master-level background
        if layout and layout.slide_master:
            for bg in layout.slide_master._element.iter(f"{{{ns_p}}}bg"):
                if _has_fill(bg):
                    return True

        return False
    except Exception:
        return True  # Assume OK if we can't check


def _warn_overlapping_shapes(slide, slide_num) -> None:
    """Log warnings for text shapes that overlap on a slide."""
    text_shapes = []
    for shape in slide.shapes:
        if not shape.has_text_frame or not shape.text_frame.text.strip():
            continue
        if shape.left is None or shape.top is None:
            continue
        if shape.width is None or shape.height is None:
            continue
        left_in = shape.left / 914400
        if left_in > 15:
            continue
        text_shapes.append(shape)

    for i, a in enumerate(text_shapes):
        for b in text_shapes[i + 1:]:
            if (a.left < b.left + b.width and a.left + a.width > b.left and
                    a.top < b.top + b.height and a.top + a.height > b.top):
                overlap_w = (min(a.left + a.width, b.left + b.width)
                             - max(a.left, b.left))
                overlap_h = (min(a.top + a.height, b.top + b.height)
                             - max(a.top, b.top))
                area = (overlap_w / 914400) * (overlap_h / 914400)
                if area > 0.1:
                    logger.warning(
                        f"Slide {slide_num}: shapes '{a.name}' and '{b.name}' "
                        f"overlap ({area:.1f} sq in)"
                    )


def _hex_luminance(hex_color: str) -> float:
    """Calculate relative luminance of a hex color (0.0 = black, 1.0 = white)."""
    hex_color = hex_color.lstrip("#")
    if len(hex_color) < 6:
        return 0.5
    try:
        r = int(hex_color[0:2], 16) / 255
        g = int(hex_color[2:4], 16) / 255
        b = int(hex_color[4:6], 16) / 255
        return 0.299 * r + 0.587 * g + 0.114 * b
    except (ValueError, IndexError):
        return 0.5


def _enforce_contrast(
    slide,
    bg_hex: str,
    var_map: dict[str, str],
) -> None:
    """Check every text shape on a slide for sufficient background contrast.

    If a text run's color is too similar to the background, flip it to the
    opposite design-system token (light text vs dark text).
    """
    bg_lum = _hex_luminance(bg_hex)
    bg_is_dark = bg_lum < 0.4

    light_color = var_map.get("--color-text-light", "#FFFFF6")
    dark_color = var_map.get("--color-text-dark", "#434343")

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if not run.text.strip():
                    continue
                try:
                    if run.font.color and run.font.color.rgb:
                        text_rgb = str(run.font.color.rgb)
                        text_lum = _hex_luminance(text_rgb)
                        text_is_dark = text_lum < 0.4

                        # Same luminance band = poor contrast
                        if bg_is_dark and text_is_dark:
                            run.font.color.rgb = _hex_to_rgb(light_color)
                        elif not bg_is_dark and not text_is_dark:
                            run.font.color.rgb = _hex_to_rgb(dark_color)
                except (AttributeError, TypeError):
                    pass


def _add_speaker_notes(slide, notes_text: str) -> None:
    """Add speaker notes to a slide."""
    try:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes_text
    except Exception as e:
        logger.debug(f"Could not add speaker notes: {e}")


# ---------------------------------------------------------------------------
# Hybrid mode: clear all text shapes on a cloned slide
# ---------------------------------------------------------------------------

def _is_image_shape(shape) -> bool:
    """Check if a shape is an image (picture) or contains images."""
    try:
        # Check for picture elements
        ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
        ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        ns_r = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

        elem = shape._element
        tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag

        # Direct picture shape
        if tag == "pic":
            return True

        # Shape with blipFill (image fill)
        if elem.find(f".//{{{ns_a}}}blipFill") is not None:
            return True
        if elem.find(f".//{{{ns_a}}}blip") is not None:
            return True

        # Group shape containing pictures
        if tag == "grpSp":
            for child in elem.iter():
                child_tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                if child_tag == "pic" or child.tag == f"{{{ns_a}}}blip":
                    return True
    except Exception:
        pass
    return False


def _clear_all_text_shapes(slide, preserve_images: bool = False) -> None:
    """Clear and move off-canvas text shapes on a cloned slide.

    Used in hybrid mode: the slide is cloned for its branded background,
    but original text is removed so compose-mode textboxes can be added.

    When preserve_images is True, image shapes and group shapes containing
    images are kept in place — only pure text shapes are cleared.
    """
    for shape in slide.shapes:
        if preserve_images and _is_image_shape(shape):
            continue

        if shape.has_text_frame:
            tf = shape.text_frame
            ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
            p_elements = list(tf._element)
            first_p = None
            for p_el in p_elements:
                if p_el.tag == f"{{{ns_a}}}p":
                    if first_p is None:
                        first_p = p_el
                    else:
                        tf._element.remove(p_el)
            if tf.paragraphs:
                tf.paragraphs[0].text = ""
            try:
                shape.left = Emu(914400 * 20)
            except Exception:
                pass


# ---------------------------------------------------------------------------
# Accent stripe decoration
# ---------------------------------------------------------------------------

def _maybe_add_accent_stripe(
    slide, slide_style: str, visual_profile: str,
    var_map: dict[str, str],
) -> None:
    """Detect CSS border-left on the slide and create a matching accent bar.

    The HTML uses `border-left: 3px solid var(--color-primary)` to render a
    left-edge accent stripe.  This function reproduces it as a thin
    rectangle shape in PPTX.
    """
    # Check for explicit border-left in inline style
    border_match = re.search(
        r"border-left\s*:\s*([\d.]+)\s*px\s+solid\s+(#[0-9A-Fa-f]{6})",
        slide_style,
    )
    if border_match:
        width_px = float(border_match.group(1))
        color_hex = border_match.group(2)
        add_accent_bar(
            slide,
            left=0.0,
            top=0.0,
            width=_px_to_inches(width_px),
            height=5.625,  # full slide height
            color_hex=color_hex,
        )
        return

    # Check for visual profile that gets an accent stripe via CSS class rules
    # (.slide[data-visual-profile="dark"] { border-left: 3px solid ... })
    if visual_profile == "dark":
        accent_color = var_map.get("--color-primary", "#CDF986")
        add_accent_bar(
            slide,
            left=0.0,
            top=0.0,
            width=_px_to_inches(3),  # 3px default stripe
            height=5.625,
            color_hex=accent_color,
        )


# ---------------------------------------------------------------------------
# Post-processing: clear unmapped shapes on cloned slides
# ---------------------------------------------------------------------------

def _clear_unmapped_shapes(slide, mapped_shape_names: set[str]) -> None:
    """Clear text from shapes that weren't mapped to any HTML element.

    Only applies to cloned slides to remove stale template text.
    """
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name not in mapped_shape_names:
            tf = shape.text_frame
            if tf.paragraphs:
                # Clear all text
                ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
                p_elements = list(tf._element)
                first_p = tf.paragraphs[0]._element
                for p_el in p_elements:
                    if p_el.tag == f"{{{ns_a}}}p" and p_el is not first_p:
                        tf._element.remove(p_el)
                tf.paragraphs[0].text = ""

                # Move off-canvas
                try:
                    shape.left = Emu(914400 * 20)
                except Exception:
                    pass


# ---------------------------------------------------------------------------
# Post-save fixes (python-pptx bug workarounds)
# ---------------------------------------------------------------------------

def _fix_notes_master_id(pptx_path: Path) -> None:
    """Inject notesMasterIdLst into presentation.xml if missing.

    python-pptx creates a notesMaster relationship when speaker notes are
    added but omits the required <p:notesMasterIdLst> element in
    presentation.xml.  Keynote and Google Slides reject files without it.

    This patches the saved ZIP in place — a single XML insertion.
    """
    import zipfile
    import tempfile
    import shutil

    with zipfile.ZipFile(pptx_path, "r") as zf:
        pres_rels = zf.read("ppt/_rels/presentation.xml.rels").decode("utf-8")
        pres_xml = zf.read("ppt/presentation.xml").decode("utf-8")

        # Check: does a notesMaster relationship exist?
        nm_match = re.search(
            r'Id="([^"]+)"[^>]*notesMaster', pres_rels
        )
        if not nm_match:
            return  # No notes in this file

        # Check: is notesMasterIdLst already present?
        if "notesMasterIdLst" in pres_xml:
            return  # Already fixed

        rid = nm_match.group(1)
        notes_element = (
            "<p:notesMasterIdLst>"
            f'<p:notesMasterId r:id="{rid}"/>'
            "</p:notesMasterIdLst>"
        )

        # Insert after </p:sldMasterIdLst>
        patched_xml = pres_xml.replace(
            "</p:sldMasterIdLst>",
            f"</p:sldMasterIdLst>{notes_element}",
        )

        if patched_xml == pres_xml:
            logger.warning("Could not find insertion point for notesMasterIdLst")
            return

    # Rewrite the ZIP with the patched presentation.xml
    tmp_fd, tmp_path = tempfile.mkstemp(suffix=".pptx")
    try:
        with zipfile.ZipFile(pptx_path, "r") as zf_in:
            with zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as zf_out:
                for name in zf_in.namelist():
                    if name == "ppt/presentation.xml":
                        zf_out.writestr(name, patched_xml.encode("utf-8"))
                    else:
                        zf_out.writestr(name, zf_in.read(name))
        shutil.move(tmp_path, pptx_path)
        logger.info(f"Fixed notesMasterIdLst (r:id={rid})")
    except Exception as e:
        logger.warning(f"Could not fix notesMasterIdLst: {e}")
        try:
            Path(tmp_path).unlink(missing_ok=True)
        except Exception:
            pass
    finally:
        try:
            import os
            os.close(tmp_fd)
        except Exception:
            pass


# ---------------------------------------------------------------------------
# Main build pipeline
# ---------------------------------------------------------------------------

def build_from_html(
    html_path: Path,
    output_path: Path,
    base_template: Path | None = None,
) -> Path:
    """Build a PPTX presentation from an HTML slide deck file.

    Args:
        html_path: Path to the HTML deck preview file.
        output_path: Where to save the output PPTX.
        base_template: Optional base template for layouts/masters.

    Returns:
        Path to the saved PPTX.
    """
    html_content = html_path.read_text(encoding="utf-8")
    soup = BeautifulSoup(html_content, "html.parser")

    # --- CSS variable resolution ---
    var_map = _parse_css_variables(soup)
    if var_map:
        logger.info(f"Resolved {len(var_map)} CSS custom properties from :root")

    # Open base template
    default_base = Path("templates/base/Shopify - Example Technical Workshop Slide Bank.pptx")
    base_path = base_template or default_base

    if base_path.exists():
        prs = open_base_template(base_path)
    else:
        logger.warning(f"Base template not found at {base_path}, creating blank")
        prs = create_presentation()

    # Find all slide sections
    slide_divs = soup.find_all("div", class_="slide")
    if not slide_divs:
        logger.error("No <div class='slide'> sections found in HTML")
        sys.exit(1)

    # Sort by data-slide-number
    def slide_sort_key(div):
        num = div.get("data-slide-number", "0")
        try:
            return int(num)
        except ValueError:
            return 0

    slide_divs.sort(key=slide_sort_key)

    counts: dict[str, int] = {"clone": 0, "hybrid": 0, "compose": 0}

    for slide_div in slide_divs:
        mode_used = _build_slide_from_section(prs, slide_div, var_map)
        counts[mode_used] = counts.get(mode_used, 0) + 1

    # Save
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    # Fix python-pptx bug: notesMasterIdLst missing from presentation.xml.
    # python-pptx creates the notesMaster relationship when speaker notes
    # are added but doesn't register it in the XML.  Keynote and Google
    # Slides require this element.
    _fix_notes_master_id(output_path)

    # Free caches
    clear_clone_caches()

    logger.info(
        f"Saved: {output_path} "
        f"({counts['clone']} cloned, {counts['hybrid']} hybrid, "
        f"{counts['compose']} composed, {len(slide_divs)} total)"
    )
    return output_path


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(description="Build PPTX from HTML slide deck")
    parser.add_argument("html_file", type=Path, help="Path to HTML deck file")
    parser.add_argument("-o", "--output", type=Path, default=Path("output/presentation.pptx"),
                        help="Output PPTX path (default: output/presentation.pptx)")
    parser.add_argument("--base-template", type=Path, default=None,
                        help="Base template PPTX (default: templates/base/...)")
    parser.add_argument("--no-repair", action="store_true",
                        help="Skip the post-build repair/compaction step")
    args = parser.parse_args()

    if not args.html_file.exists():
        print(f"Error: HTML file not found: {args.html_file}", file=sys.stderr)
        sys.exit(1)

    result_path = build_from_html(
        html_path=args.html_file,
        output_path=args.output,
        base_template=args.base_template,
    )

    # Post-build compaction: strip fonts, compress images, inject docProps
    if not args.no_repair:
        try:
            from scripts.repair_pptx import repair_pptx
            original_mb = result_path.stat().st_size / (1024 * 1024)
            if original_mb > 10:  # Only compact files > 10 MB
                logger.info(f"Running post-build compaction ({original_mb:.0f} MB)...")
                stats = repair_pptx(result_path, result_path)
                logger.info(
                    f"Compacted: {stats['original_size_mb']:.0f} MB → "
                    f"{stats['final_size_mb']:.0f} MB "
                    f"({stats['fonts_removed']} fonts, "
                    f"{stats['media_compressed']} images compressed)"
                )
        except Exception as e:
            logger.warning(f"Post-build compaction skipped: {e}")

    print(f"Presentation generated: {result_path}")


if __name__ == "__main__":
    main()
