#!/usr/bin/env python3
"""Analyze PPTX template files and extract structural metadata.

This script performs the non-LLM portion of template analysis:
1. Extracts placeholder info, shape counts, colors, fonts from each slide
2. Writes slide descriptions to a file for the Cursor agent to classify
3. After classification, merges the agent's classifications into a template registry

Usage:
    # Step 1: Extract structural metadata and write descriptions
    python scripts/analyze_templates.py extract <template_dir> [-o workspace/template_descriptions.json]

    # Step 2: After the Cursor agent classifies slides, merge into registry
    python scripts/analyze_templates.py merge <descriptions_json> <classifications_json> [-o template_registry.json]
"""

import argparse
import json
import sys
from pathlib import Path

# Add project root to path
sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from pptx import Presentation

from pptx.enum.shapes import PP_PLACEHOLDER

from src.schemas.slide_schema import SlideType
from src.schemas.template_schema import (
    ContentZone,
    DecorationAsset,
    PlaceholderInfo,
    TemplateRegistry,
    TemplateSlide,
    TextContent,
)
from src.utils.file_utils import find_pptx_files


def _safe_inches(emu_value) -> float:
    """Safely convert EMU value to inches."""
    try:
        if emu_value is None:
            return 0.0
        return emu_value / 914400
    except Exception:
        return 0.0


def _is_unsupported(shape) -> bool:
    """Check if a shape is unsupported."""
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        if shape.shape_type in (
            MSO_SHAPE_TYPE.CHART,
            MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,
            MSO_SHAPE_TYPE.LINKED_OLE_OBJECT,
            MSO_SHAPE_TYPE.SMART_ART,
        ):
            return True
    except Exception:
        pass
    return False


def _extract_text_content(slide) -> dict:
    """Extract structured text content from a slide.

    Returns a dict with 'title', 'body', and 'all_text' keys.
    """
    title_parts = []
    body_parts = []
    all_parts = []

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        text = shape.text_frame.text.strip()
        if not text:
            continue

        all_parts.append(text)

        # Classify by placeholder type
        is_title = False
        is_body = False

        if shape.is_placeholder:
            try:
                ph_type = shape.placeholder_format.type
                if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                    is_title = True
                elif ph_type in (PP_PLACEHOLDER.SUBTITLE, PP_PLACEHOLDER.BODY,
                                 PP_PLACEHOLDER.OBJECT):
                    is_body = True
            except Exception:
                pass

        if is_title:
            title_parts.append(text)
        elif is_body:
            body_parts.append(text)
        else:
            # Heuristic: large text near top = title candidate
            try:
                if shape.top is not None and shape.top / 914400 < 1.5:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size and run.font.size.pt >= 28:
                                title_parts.append(text)
                                break
                        else:
                            continue
                        break
                    else:
                        body_parts.append(text)
                else:
                    body_parts.append(text)
            except Exception:
                body_parts.append(text)

    return {
        "title": " | ".join(title_parts),
        "body": " | ".join(body_parts),
        "all_text": " ".join(all_parts),
    }


def _build_slide_description(
    index: int,
    placeholders: list[dict],
    shape_count: int,
    has_images: bool,
    layout_name: str,
    text_content: dict | None = None,
) -> str:
    """Build a text description of a slide layout for classification."""
    parts = [f"Layout: '{layout_name}'" if layout_name else "Layout: (unnamed)"]
    parts.append(f"Total shapes: {shape_count}")

    if has_images:
        parts.append("Contains images")

    if placeholders:
        ph_desc = []
        for ph in placeholders:
            pos = ph["position"]
            size_desc = f"{pos[2]:.1f}x{pos[3]:.1f} inches"
            ph_desc.append(f"  - {ph['name']} (type={ph['type']}, size={size_desc})")
        parts.append("Placeholders:\n" + "\n".join(ph_desc))
    else:
        parts.append("No placeholders (free-form shapes only)")

    # Include extracted text content for richer classification
    if text_content:
        if text_content.get("title"):
            parts.append(f"Title text: \"{text_content['title']}\"")
        if text_content.get("body"):
            # Truncate very long body text
            body = text_content["body"]
            if len(body) > 300:
                body = body[:300] + "..."
            parts.append(f"Body text: \"{body}\"")
        if not text_content.get("title") and not text_content.get("body"):
            all_text = text_content.get("all_text", "")
            if all_text:
                if len(all_text) > 300:
                    all_text = all_text[:300] + "..."
                parts.append(f"Visible text: \"{all_text}\"")

    return "\n".join(parts)


def _extract_content_zones(slide) -> list[dict]:
    """Identify content zones (replaceable text areas) in a slide.

    A content zone is a text shape that holds primary content (title, body,
    subtitle) as opposed to design elements (labels, badges, decorations).
    """
    zones = []

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        text = shape.text_frame.text.strip()
        if not text:
            continue

        # Skip very small shapes (badges, labels, footers)
        try:
            w_inches = _safe_inches(shape.width) if shape.width else 0
            h_inches = _safe_inches(shape.height) if shape.height else 0
            if w_inches < 1.5 or h_inches < 0.3:
                continue
        except Exception:
            continue

        # Determine zone type from placeholder type or heuristics
        zone_type = "body"
        if shape.is_placeholder:
            try:
                ph_type = shape.placeholder_format.type
                if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                    zone_type = "title"
                elif ph_type == PP_PLACEHOLDER.SUBTITLE:
                    zone_type = "subtitle"
                elif ph_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
                    zone_type = "body"
            except Exception:
                pass
        else:
            # Heuristic: large font near top = title, otherwise body
            max_font = 0
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        max_font = max(max_font, run.font.size.pt)
            top_inches = _safe_inches(shape.top) if shape.top else 0
            if max_font >= 28 and top_inches < 2.0:
                zone_type = "title"
            elif max_font >= 40:
                zone_type = "data_point"

        # Calculate capacity
        area = w_inches * h_inches
        max_chars = int(area * 80)  # ~80 chars per square inch

        # Get font size range
        font_sizes = []
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    font_sizes.append(int(run.font.size.pt))
        font_range = (min(font_sizes), max(font_sizes)) if font_sizes else (12, 24)

        zones.append({
            "zone_type": zone_type,
            "shape_name": shape.name,
            "position": (
                _safe_inches(shape.left),
                _safe_inches(shape.top),
                w_inches,
                h_inches,
            ),
            "max_chars": max_chars,
            "font_size_range": font_range,
        })

    return zones


def _detect_background_type(slide) -> str:
    """Detect the background type of a slide."""
    try:
        bg = slide.background
        if bg is None or bg._element is None:
            return "master_inherited"

        bg_elem = bg._element
        if len(bg_elem) == 0:
            return "master_inherited"

        ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"

        # Check for bgPr (explicit background properties)
        for bgPr in bg_elem.iter(f"{{{ns_p}}}bgPr"):
            # Check for image fill
            for _blip_fill in bgPr.iter(f"{{{ns_a}}}blipFill"):
                return "image"
            # Check for gradient
            for _grad in bgPr.iter(f"{{{ns_a}}}gradFill"):
                return "gradient"
            # Check for solid fill
            for _solid in bgPr.iter(f"{{{ns_a}}}solidFill"):
                return "solid"

        # Check for bgRef (reference to theme background)
        for _bgRef in bg_elem.iter(f"{{{ns_p}}}bgRef"):
            return "master_inherited"

        return "master_inherited"
    except Exception:
        return "none"


def _extract_background_color(slide) -> str | None:
    """Extract the dominant background color from a slide.

    Checks in order: slide-level explicit background, layout background,
    master background. Returns a hex RGB string (e.g., '#000000') or None.
    """
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"

    def _color_from_bg_element(bg_elem) -> str | None:
        """Try to extract a color from a background XML element."""
        if bg_elem is None or len(bg_elem) == 0:
            return None
        # Check bgPr for solid fill
        for bgPr in bg_elem.iter(f"{{{ns_p}}}bgPr"):
            for solid in bgPr.iter(f"{{{ns_a}}}solidFill"):
                for srgb in solid.iter(f"{{{ns_a}}}srgbClr"):
                    val = srgb.get("val", "")
                    if val:
                        return f"#{val}"
            # Gradient: return first stop color
            for grad in bgPr.iter(f"{{{ns_a}}}gradFill"):
                for gs in grad.iter(f"{{{ns_a}}}gs"):
                    for srgb in gs.iter(f"{{{ns_a}}}srgbClr"):
                        val = srgb.get("val", "")
                        if val:
                            return f"#{val}"
                    break  # Only first stop
        return None

    try:
        # 1. Slide-level background
        bg = slide.background
        if bg and bg._element is not None:
            color = _color_from_bg_element(bg._element)
            if color:
                return color

        # 2. Layout-level background
        try:
            layout_part = None
            from pptx.opc.constants import RELATIONSHIP_TYPE as RT
            for _key, rel in slide.part.rels.items():
                if rel.reltype == RT.SLIDE_LAYOUT:
                    layout_part = rel.target_part
                    break
            if layout_part:
                layout_xml = layout_part._element
                for layout_bg in layout_xml.iter(f"{{{ns_p}}}bg"):
                    color = _color_from_bg_element(layout_bg)
                    if color:
                        return color
        except Exception:
            pass

        # 3. Master-level background
        try:
            if layout_part:
                master_part = None
                for _key, rel in layout_part.rels.items():
                    if rel.reltype == RT.SLIDE_MASTER:
                        master_part = rel.target_part
                        break
                if master_part:
                    master_xml = master_part._element
                    for master_bg in master_xml.iter(f"{{{ns_p}}}bg"):
                        color = _color_from_bg_element(master_bg)
                        if color:
                            return color
        except Exception:
            pass

    except Exception:
        pass

    return None


def _detect_visual_profile(slide, has_images: bool, has_background: bool, bg_type: str) -> str:
    """Detect the visual profile of a slide."""
    if bg_type == "image" or (has_images and has_background):
        return "branded_image"

    # Check if background is dark by sampling colors
    dark_colors = {"#000000", "#191E17", "#1a1a1a", "#0d0d0d", "#111111", "#222222"}
    try:
        bg = slide.background
        if bg and bg._element is not None:
            ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
            for srgb in bg._element.iter(f"{{{ns_a}}}srgbClr"):
                val = srgb.get("val", "")
                if f"#{val}" in dark_colors or val.lower() in {"000000", "191e17"}:
                    return "dark"
    except Exception:
        pass

    # Check if shapes suggest darkness (lots of light-colored text)
    light_text_count = 0
    dark_text_count = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                try:
                    if run.font.color and run.font.color.rgb:
                        rgb = str(run.font.color.rgb).upper()
                        r, g, b = int(rgb[0:2], 16), int(rgb[2:4], 16), int(rgb[4:6], 16)
                        luminance = (r * 299 + g * 587 + b * 114) / 1000
                        if luminance > 180:
                            light_text_count += 1
                        else:
                            dark_text_count += 1
                except (AttributeError, TypeError, ValueError):
                    pass

    if light_text_count > dark_text_count and light_text_count >= 2:
        return "dark"
    elif has_images:
        return "branded_image"
    elif has_background:
        return "light"
    else:
        return "minimal"


def _classify_images(slide, image_count: int, bg_type: str) -> str:
    """Classify images on a slide as 'none', 'decorative', or 'content'.

    Heuristic:
    - 'none' if no images at all
    - 'decorative' if the slide has a background image or only 1-2 small images
      (likely logos/accents) or if images are near edges (header/footer areas)
    - 'content' if there are 3+ picture shapes or large centered images
      (likely product shots, screenshots, photos that won't make sense when cloned)
    """
    if image_count == 0:
        return "none"

    # If the only "image" is the background, it's decorative
    if bg_type == "image" and image_count <= 1:
        return "decorative"

    # Count picture shapes and assess their sizes / positions
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    large_images = 0
    center_images = 0
    try:
        slide_w = slide.part.package.presentation.slide_width / 914400
        slide_h = slide.part.package.presentation.slide_height / 914400
    except Exception:
        slide_w, slide_h = 10.0, 5.625

    for shape in slide.shapes:
        try:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
        except Exception:
            continue

        w = (shape.width / 914400) if shape.width else 0
        h = (shape.height / 914400) if shape.height else 0
        left = (shape.left / 914400) if shape.left else 0
        top = (shape.top / 914400) if shape.top else 0

        area = w * h
        if area > 3.0:
            large_images += 1
        # Check if image is roughly centered (not a corner logo)
        cx = left + w / 2
        cy = top + h / 2
        if 0.2 * slide_w < cx < 0.8 * slide_w and 0.2 * slide_h < cy < 0.8 * slide_h:
            center_images += 1

    # 3+ images or 2+ large centered images → content images
    if image_count >= 3 or (large_images >= 2 and center_images >= 2):
        return "content"

    return "decorative"


def _count_images(slide) -> int:
    """Count the number of picture shapes on a slide."""
    count = 0
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        for shape in slide.shapes:
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    count += 1
            except Exception:
                pass
    except Exception:
        pass
    return count


def _extract_decoration_assets(slide, content_zone_names: set[str]) -> list[dict]:
    """Catalog decorative elements (non-content shapes) on a slide.

    Identifies shapes that are NOT content zones (text placeholders) and
    classifies them as decorative assets: accent shapes, divider lines,
    images, badges, logos, etc.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    assets = []

    for shape in slide.shapes:
        # Skip shapes that are content zones (already cataloged as replaceable text)
        if shape.name in content_zone_names:
            continue

        left = _safe_inches(shape.left) if shape.left else 0
        top = _safe_inches(shape.top) if shape.top else 0
        w = _safe_inches(shape.width) if shape.width else 0
        h = _safe_inches(shape.height) if shape.height else 0

        # Skip shapes with zero dimensions
        if w < 0.05 and h < 0.05:
            continue

        position = (left, top, w, h)
        area = w * h
        color = None
        asset_type = None
        description = ""
        is_branded = True
        group_id = None

        try:
            shape_type = shape.shape_type
        except Exception:
            continue

        # --- Pictures ---
        try:
            if shape_type == MSO_SHAPE_TYPE.PICTURE:
                # Classify by size and position
                if area > 20.0:
                    asset_type = "background_image"
                    description = "Full-slide or near-full background image"
                elif w < 1.5 and h < 1.5 and (top < 0.5 or top > 4.5):
                    asset_type = "logo"
                    description = f"Small image ({w:.1f}x{h:.1f}in) in header/footer area"
                elif area > 5.0:
                    # Large centered image — likely content-specific
                    asset_type = "photo"
                    description = f"Large image ({w:.1f}x{h:.1f}in)"
                    is_branded = False
                else:
                    asset_type = "illustration"
                    description = f"Medium image ({w:.1f}x{h:.1f}in)"

                assets.append({
                    "asset_type": asset_type,
                    "shape_name": shape.name,
                    "position": position,
                    "description": description,
                    "is_branded": is_branded,
                    "color": color,
                    "group_id": group_id,
                })
                continue
        except Exception:
            pass

        # --- Auto shapes (rectangles, ovals, lines, etc.) ---
        try:
            if shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                # Extract fill color if available
                try:
                    if shape.fill and shape.fill.type is not None:
                        fg = shape.fill.fore_color
                        if fg and fg.rgb:
                            color = f"#{fg.rgb}"
                except Exception:
                    pass

                # Classify by shape and size
                if h < 0.1 and w > 1.0:
                    asset_type = "divider_line"
                    description = f"Thin horizontal shape ({w:.1f}x{h:.2f}in)"
                elif w < 0.1 and h > 1.0:
                    asset_type = "divider_line"
                    description = f"Thin vertical shape ({w:.2f}x{h:.1f}in)"
                elif w < 0.8 and h < 0.8 and abs(w - h) < 0.2:
                    # Small, roughly square — likely a badge or icon container
                    asset_type = "badge"
                    description = f"Small shape ({w:.1f}x{h:.1f}in)"
                    # Check if it contains text (numbered badge)
                    if shape.has_text_frame and shape.text_frame.text.strip():
                        txt = shape.text_frame.text.strip()
                        description = f"Badge with text '{txt}'"
                elif area > 15.0:
                    asset_type = "frame"
                    description = f"Large shape ({w:.1f}x{h:.1f}in) — likely a frame or panel"
                elif h < 0.15:
                    asset_type = "accent_shape"
                    description = f"Accent bar ({w:.1f}x{h:.2f}in)"
                elif w < 0.15:
                    asset_type = "accent_shape"
                    description = f"Vertical accent ({w:.2f}x{h:.1f}in)"
                else:
                    asset_type = "accent_shape"
                    description = f"Decorative shape ({w:.1f}x{h:.1f}in)"

                assets.append({
                    "asset_type": asset_type,
                    "shape_name": shape.name,
                    "position": position,
                    "description": description,
                    "is_branded": is_branded,
                    "color": color,
                    "group_id": group_id,
                })
                continue
        except Exception:
            pass

        # --- Freeform shapes ---
        try:
            if shape_type == MSO_SHAPE_TYPE.FREEFORM:
                asset_type = "illustration"
                description = f"Freeform shape ({w:.1f}x{h:.1f}in)"

                try:
                    if shape.fill and shape.fill.type is not None:
                        fg = shape.fill.fore_color
                        if fg and fg.rgb:
                            color = f"#{fg.rgb}"
                except Exception:
                    pass

                assets.append({
                    "asset_type": asset_type,
                    "shape_name": shape.name,
                    "position": position,
                    "description": description,
                    "is_branded": is_branded,
                    "color": color,
                    "group_id": group_id,
                })
                continue
        except Exception:
            pass

        # --- Lines ---
        try:
            if shape_type in (MSO_SHAPE_TYPE.LINE,):
                asset_type = "divider_line"
                description = f"Line ({w:.1f}x{h:.1f}in)"

                assets.append({
                    "asset_type": asset_type,
                    "shape_name": shape.name,
                    "position": position,
                    "description": description,
                    "is_branded": is_branded,
                    "color": color,
                    "group_id": group_id,
                })
                continue
        except Exception:
            pass

        # --- Group shapes ---
        try:
            if shape_type == MSO_SHAPE_TYPE.GROUP:
                asset_type = "illustration"
                description = f"Grouped shapes ({w:.1f}x{h:.1f}in)"

                assets.append({
                    "asset_type": asset_type,
                    "shape_name": shape.name,
                    "position": position,
                    "description": description,
                    "is_branded": is_branded,
                    "color": color,
                    "group_id": group_id,
                })
                continue
        except Exception:
            pass

        # --- Placeholder shapes that are picture-type (empty picture placeholders) ---
        if shape.is_placeholder:
            try:
                ph_type = shape.placeholder_format.type
                if ph_type == PP_PLACEHOLDER.PICTURE:
                    asset_type = "chart_placeholder"
                    description = f"Picture placeholder ({w:.1f}x{h:.1f}in)"
                    assets.append({
                        "asset_type": asset_type,
                        "shape_name": shape.name,
                        "position": position,
                        "description": description,
                        "is_branded": is_branded,
                        "color": color,
                        "group_id": group_id,
                    })
            except Exception:
                pass

    return assets


def _calculate_content_capacity(content_zones: list[dict]) -> str:
    """Calculate how much text a slide can hold based on its content zones."""
    total_area = 0
    for zone in content_zones:
        pos = zone["position"]
        total_area += pos[2] * pos[3]  # width * height

    if total_area < 3.0:
        return "low"
    elif total_area < 10.0:
        return "medium"
    else:
        return "high"


def _heuristic_classify(description: str) -> dict:
    """Fallback heuristic classification."""
    desc_lower = description.lower()

    if "title" in desc_lower and "subtitle" in desc_lower:
        return {"slide_type": "title", "tags": ["generic"], "complexity": 1,
                "description": "Title slide with subtitle"}
    elif "title" in desc_lower and "body" not in desc_lower:
        return {"slide_type": "section_header", "tags": ["generic"], "complexity": 1,
                "description": "Section header slide"}
    elif "picture" in desc_lower or "image" in desc_lower:
        return {"slide_type": "image_with_text", "tags": ["visual"], "complexity": 3,
                "description": "Slide with image content"}
    else:
        return {"slide_type": "content", "tags": ["generic"], "complexity": 2,
                "description": "General content slide"}


def extract_metadata(template_dir: Path, output_path: Path):
    """Extract structural metadata from all PPTX files in a directory."""
    pptx_files = find_pptx_files(template_dir)

    if not pptx_files:
        print(f"No .pptx files found in {template_dir}", file=sys.stderr)
        sys.exit(1)

    print(f"Found {len(pptx_files)} .pptx files to analyze")

    all_slides = []
    source_files = []

    for pptx_path in pptx_files:
        print(f"  Analyzing: {pptx_path.name}")
        try:
            prs = Presentation(str(pptx_path))

            for idx, slide in enumerate(prs.slides):
                placeholders = []
                colors: set[str] = set()
                fonts: set[str] = set()
                shape_count = 0
                has_images = False
                has_background = False

                for shape in slide.shapes:
                    shape_count += 1

                    if _is_unsupported(shape):
                        continue

                    if shape.is_placeholder:
                        try:
                            ph_type = shape.placeholder_format.type
                            type_name = ph_type.name if ph_type else "other"
                        except Exception:
                            type_name = "other"

                        placeholders.append({
                            "name": shape.name,
                            "type": type_name,
                            "position": (
                                _safe_inches(shape.left),
                                _safe_inches(shape.top),
                                _safe_inches(shape.width),
                                _safe_inches(shape.height),
                            ),
                        })

                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                try:
                                    if run.font.color and run.font.color.rgb:
                                        colors.add(f"#{run.font.color.rgb}")
                                except (AttributeError, TypeError):
                                    pass
                                if run.font.name:
                                    fonts.add(run.font.name)

                    try:
                        from pptx.enum.shapes import MSO_SHAPE_TYPE
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            has_images = True
                    except Exception:
                        pass

                    try:
                        if shape.fill and shape.fill.type is not None:
                            fg = shape.fill.fore_color
                            if fg and fg.rgb:
                                colors.add(f"#{fg.rgb}")
                    except Exception:
                        pass

                try:
                    bg = slide.background
                    if bg and bg._element is not None and len(bg._element) > 0:
                        has_background = True
                except Exception:
                    pass

                try:
                    layout_name = slide.slide_layout.name if slide.slide_layout else ""
                except Exception:
                    layout_name = ""

                # Extract text content for semantic analysis
                text_content = _extract_text_content(slide)

                # Extract content zones (replaceable text areas)
                content_zones = _extract_content_zones(slide)

                # Detect background type and visual profile
                background_type = _detect_background_type(slide)
                visual_profile = _detect_visual_profile(
                    slide, has_images, has_background, background_type
                )
                content_capacity = _calculate_content_capacity(content_zones)

                # Image classification
                image_count = _count_images(slide)
                image_type = _classify_images(slide, image_count, background_type)

                # Background color extraction
                background_color = _extract_background_color(slide)

                # Extract decoration assets (non-content decorative shapes)
                content_zone_names = {cz["shape_name"] for cz in content_zones}
                decoration_assets = _extract_decoration_assets(slide, content_zone_names)

                description = _build_slide_description(
                    idx, placeholders, shape_count, has_images, layout_name,
                    text_content=text_content,
                )

                all_slides.append({
                    "template_file": str(pptx_path),
                    "slide_index": idx,
                    "layout_name": layout_name,
                    "placeholders": placeholders,
                    "color_scheme": sorted(colors),
                    "font_families": sorted(fonts),
                    "shape_count": shape_count,
                    "has_images": has_images,
                    "has_background": has_background,
                    "text_content": text_content,
                    "content_zones": content_zones,
                    "background_type": background_type,
                    "visual_profile": visual_profile,
                    "content_capacity": content_capacity,
                    "image_type": image_type,
                    "image_count": image_count,
                    "background_color": background_color,
                    "decoration_assets": decoration_assets,
                    "description_for_classification": description,
                })

            source_files.append(str(pptx_path))
        except Exception as e:
            print(f"  Warning: Failed to analyze {pptx_path.name}: {e}", file=sys.stderr)

    result = {
        "source_files": source_files,
        "slides": all_slides,
        "total_slides": len(all_slides),
    }

    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(json.dumps(result, indent=2, default=str), encoding="utf-8")

    print(f"\nExtracted metadata for {len(all_slides)} slides from {len(source_files)} files")
    print(f"Written to: {output_path}")
    print(f"\nNext step: The Cursor agent should read this file and classify each slide,")
    print(f"then run: python scripts/analyze_templates.py merge {output_path} workspace/classifications.json")


def merge_classifications(descriptions_path: Path, classifications_path: Path, output_path: Path):
    """Merge agent classifications with structural metadata into a template registry."""
    descriptions_data = json.loads(descriptions_path.read_text(encoding="utf-8"))
    classifications_data = json.loads(classifications_path.read_text(encoding="utf-8"))

    classifications = classifications_data.get("classifications", [])
    slides_meta = descriptions_data.get("slides", [])

    templates = []
    for i, slide_data in enumerate(slides_meta):
        # Get classification (or use heuristic fallback)
        if i < len(classifications):
            cls = classifications[i]
        else:
            cls = _heuristic_classify(slide_data.get("description_for_classification", ""))

        # Convert placeholders to PlaceholderInfo format
        placeholders = []
        for ph in slide_data.get("placeholders", []):
            placeholders.append(PlaceholderInfo(
                name=ph["name"],
                type=ph["type"],
                position=tuple(ph["position"]),
            ))

        # Build TextContent from extracted data
        text_data = slide_data.get("text_content")
        text_content_obj = None
        if text_data:
            text_content_obj = TextContent(
                title=text_data.get("title", ""),
                body=text_data.get("body", ""),
                all_text=text_data.get("all_text", ""),
            )

        # Build ContentZone objects from extracted data
        content_zones_data = slide_data.get("content_zones", [])
        content_zone_objs = []
        for cz in content_zones_data:
            content_zone_objs.append(ContentZone(
                zone_type=cz.get("zone_type", "body"),
                shape_name=cz.get("shape_name", ""),
                position=tuple(cz["position"]),
                max_chars=cz.get("max_chars", 200),
                font_size_range=tuple(cz.get("font_size_range", (10, 44))),
            ))

        # Build DecorationAsset objects from extracted data
        decoration_assets_data = slide_data.get("decoration_assets", [])
        decoration_asset_objs = []
        for da in decoration_assets_data:
            decoration_asset_objs.append(DecorationAsset(
                asset_type=da.get("asset_type", "accent_shape"),
                shape_name=da.get("shape_name", ""),
                position=tuple(da["position"]),
                description=da.get("description", ""),
                is_branded=da.get("is_branded", True),
                color=da.get("color"),
                group_id=da.get("group_id"),
            ))

        template = TemplateSlide(
            template_file=slide_data["template_file"],
            slide_index=slide_data["slide_index"],
            slide_type=SlideType(cls.get("slide_type", "content")),
            layout_name=slide_data.get("layout_name", ""),
            placeholders=placeholders,
            color_scheme=slide_data.get("color_scheme", []),
            font_families=slide_data.get("font_families", []),
            tags=cls.get("tags", []),
            complexity=cls.get("complexity", 2),
            shape_count=slide_data.get("shape_count", 0),
            has_images=slide_data.get("has_images", False),
            has_background=slide_data.get("has_background", False),
            description=cls.get("description", ""),
            text_content=text_content_obj,
            content_keywords=cls.get("content_keywords", []),
            visual_elements=cls.get("visual_elements", []),
            suitable_for=cls.get("suitable_for", []),
            content_zones=content_zone_objs,
            background_type=slide_data.get("background_type", "none"),
            visual_profile=slide_data.get("visual_profile", "minimal"),
            content_capacity=slide_data.get("content_capacity", "medium"),
            image_type=slide_data.get("image_type", "none"),
            image_count=slide_data.get("image_count", 0),
            background_color=slide_data.get("background_color"),
            decoration_assets=decoration_asset_objs,
        )
        templates.append(template)

    # Deduplicate
    deduped = _deduplicate(templates)

    registry = TemplateRegistry(
        templates=deduped,
        source_files=descriptions_data.get("source_files", []),
    )

    registry.save(output_path)
    print(f"Template registry created: {output_path}")
    print(f"  Total slides analyzed: {len(templates)}")
    print(f"  Unique layouts after dedup: {len(deduped)}")


def _layout_similarity(a: TemplateSlide, b: TemplateSlide) -> float:
    """Compute similarity between two slide layouts (0.0 to 1.0)."""
    score = 0.0
    factors = 0

    if a.slide_type == b.slide_type:
        score += 0.4
    factors += 0.4

    ph_diff = abs(len(a.placeholders) - len(b.placeholders))
    if ph_diff == 0:
        score += 0.2
    elif ph_diff == 1:
        score += 0.1
    factors += 0.2

    shape_diff = abs(a.shape_count - b.shape_count)
    if shape_diff <= 1:
        score += 0.15
    elif shape_diff <= 3:
        score += 0.05
    factors += 0.15

    if a.tags and b.tags:
        overlap = len(set(a.tags) & set(b.tags))
        total = len(set(a.tags) | set(b.tags))
        if total > 0:
            score += 0.15 * (overlap / total)
    factors += 0.15

    if a.layout_name and b.layout_name and a.layout_name == b.layout_name:
        score += 0.1
    factors += 0.1

    return score / factors if factors > 0 else 0.0


def _deduplicate(slides: list[TemplateSlide], threshold: float = 0.85) -> list[TemplateSlide]:
    """Remove near-duplicate slide layouts."""
    if not slides:
        return []

    unique: list[TemplateSlide] = []
    for candidate in slides:
        is_dup = False
        for existing in unique:
            if _layout_similarity(candidate, existing) >= threshold:
                is_dup = True
                if candidate.shape_count > existing.shape_count:
                    unique.remove(existing)
                    unique.append(candidate)
                break
        if not is_dup:
            unique.append(candidate)

    return unique


def main():
    parser = argparse.ArgumentParser(description="Analyze PPTX templates")
    subparsers = parser.add_subparsers(dest="command", required=True)

    # Extract subcommand
    extract_parser = subparsers.add_parser("extract", help="Extract structural metadata from templates")
    extract_parser.add_argument("template_dir", type=Path, help="Template directory to analyze")
    extract_parser.add_argument("-o", "--output", type=Path,
                                default=Path("workspace/template_descriptions.json"),
                                help="Output JSON path")

    # Merge subcommand
    merge_parser = subparsers.add_parser("merge", help="Merge classifications into template registry")
    merge_parser.add_argument("descriptions", type=Path, help="Path to template_descriptions.json")
    merge_parser.add_argument("classifications", type=Path, help="Path to classifications.json")
    merge_parser.add_argument("-o", "--output", type=Path,
                              default=Path("template_registry.json"),
                              help="Output registry path")

    args = parser.parse_args()

    if args.command == "extract":
        if not args.template_dir.is_dir():
            print(f"Error: Not a directory: {args.template_dir}", file=sys.stderr)
            sys.exit(1)
        extract_metadata(args.template_dir, args.output)
    elif args.command == "merge":
        if not args.descriptions.exists():
            print(f"Error: File not found: {args.descriptions}", file=sys.stderr)
            sys.exit(1)
        if not args.classifications.exists():
            print(f"Error: File not found: {args.classifications}", file=sys.stderr)
            sys.exit(1)
        merge_classifications(args.descriptions, args.classifications, args.output)


if __name__ == "__main__":
    main()
