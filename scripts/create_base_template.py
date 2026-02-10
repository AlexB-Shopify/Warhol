#!/usr/bin/env python3
"""Create a clean base template by stripping all slides from a source template.

Preserves the Slide Master, slide layouts, themes, and branding.
Produces a .pptx with zero slides but all layouts available.

Usage:
    python scripts/create_base_template.py <source_template> -o <output_path>
"""

import argparse
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from pptx import Presentation


def strip_slides(prs: Presentation) -> None:
    """Remove all slides from a presentation, keeping layouts intact."""
    slide_id_list = prs.slides._sldIdLst
    # Work backwards to avoid index shift
    for i in range(len(slide_id_list) - 1, -1, -1):
        rId = slide_id_list[i].rId
        prs.part.drop_rel(rId)
        del slide_id_list[i]


def main():
    parser = argparse.ArgumentParser(description="Create clean base template")
    parser.add_argument("source", type=Path, help="Source template PPTX")
    parser.add_argument("-o", "--output", type=Path, default=Path("templates/base/shopify_base.pptx"),
                        help="Output path for clean base template")
    args = parser.parse_args()

    if not args.source.exists():
        print(f"Error: Source template not found: {args.source}", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(str(args.source))

    print(f"Source: {args.source}")
    print(f"  Dimensions: {prs.slide_width.inches:.1f}\" x {prs.slide_height.inches:.3f}\"")
    print(f"  Slides: {len(prs.slides)}")
    print(f"  Layouts: {len(prs.slide_layouts)}")

    strip_slides(prs)

    args.output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(args.output))

    # Verify
    check = Presentation(str(args.output))
    print(f"\nOutput: {args.output}")
    print(f"  Slides: {len(check.slides)} (should be 0)")
    print(f"  Layouts: {len(check.slide_layouts)}")
    for i, layout in enumerate(check.slide_layouts):
        phs = [f"{ph.name}({ph.placeholder_format.type})" for ph in layout.placeholders]
        print(f"    {i}: {layout.name} -> {phs}")


if __name__ == "__main__":
    main()
