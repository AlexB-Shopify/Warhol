#!/usr/bin/env python3
"""Generate the built-in minimal/basic.pptx template.

Run this script once to create the basic template:
    python scripts/generate_basic_template.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


def create_basic_template():
    """Generate a minimal template with common slide layouts."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Colors
    primary = RGBColor(0x1A, 0x73, 0xE8)
    text_dark = RGBColor(0x20, 0x21, 0x24)
    text_light = RGBColor(0x5F, 0x63, 0x68)
    white = RGBColor(0xFF, 0xFF, 0xFF)
    bg_light = RGBColor(0xF8, 0xF9, 0xFA)

    # ---- Slide 1: Title Slide ----
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Top accent bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = primary
    bar.line.fill.background()

    # Title placeholder
    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.333), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Presentation Title"
    p.font.size = Pt(44)
    p.font.name = "Arial"
    p.font.color.rgb = text_dark
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Subtitle placeholder
    sub_box = slide.shapes.add_textbox(Inches(2.0), Inches(3.8), Inches(9.333), Inches(1.0))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Subtitle goes here"
    p.font.size = Pt(24)
    p.font.name = "Arial"
    p.font.color.rgb = text_light
    p.alignment = PP_ALIGN.CENTER

    # ---- Slide 2: Section Header ----
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = primary
    bg.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.333), Inches(2.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Section Title"
    p.font.size = Pt(40)
    p.font.name = "Arial"
    p.font.color.rgb = white
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(Inches(2.0), Inches(4.5), Inches(9.333), Inches(1.0))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Section subtitle"
    p.font.size = Pt(24)
    p.font.name = "Arial"
    p.font.color.rgb = white
    p.alignment = PP_ALIGN.CENTER

    # ---- Slide 3: Content with Title ----
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Thin accent bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = primary
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Slide Title"
    p.font.size = Pt(32)
    p.font.name = "Arial"
    p.font.color.rgb = text_dark
    p.font.bold = True

    # Body content area
    body_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(11.333), Inches(5.0))
    tf = body_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Content goes here"
    p.font.size = Pt(18)
    p.font.name = "Arial"
    p.font.color.rgb = text_dark

    # ---- Slide 4: Two Column ----
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = primary
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Two Column Title"
    p.font.size = Pt(32)
    p.font.name = "Arial"
    p.font.color.rgb = text_dark
    p.font.bold = True

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.0))
    tf = left_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Left column content"
    p.font.size = Pt(18)
    p.font.name = "Arial"
    p.font.color.rgb = text_dark

    # Right column
    right_box = slide.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.5), Inches(5.0))
    tf = right_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Right column content"
    p.font.size = Pt(18)
    p.font.name = "Arial"
    p.font.color.rgb = text_dark

    # ---- Slide 5: Quote ----
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Quote mark
    q_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(2.0), Inches(2.0))
    tf = q_box.text_frame
    p = tf.paragraphs[0]
    p.text = "\u201C"
    p.font.size = Pt(96)
    p.font.name = "Arial"
    p.font.color.rgb = primary
    p.font.bold = True

    # Quote text
    quote_box = slide.shapes.add_textbox(Inches(2.0), Inches(2.5), Inches(9.333), Inches(2.5))
    tf = quote_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Quote text goes here"
    p.font.size = Pt(24)
    p.font.name = "Arial"
    p.font.color.rgb = text_dark
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER

    # Attribution
    attr_box = slide.shapes.add_textbox(Inches(2.0), Inches(5.2), Inches(9.333), Inches(0.6))
    tf = attr_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "\u2014 Attribution"
    p.font.size = Pt(16)
    p.font.name = "Arial"
    p.font.color.rgb = text_light
    p.alignment = PP_ALIGN.CENTER

    # ---- Slide 6: Closing ----
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = primary
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.333), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(40)
    p.font.name = "Arial"
    p.font.color.rgb = text_dark
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(Inches(2.0), Inches(4.2), Inches(9.333), Inches(1.0))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Contact information"
    p.font.size = Pt(24)
    p.font.name = "Arial"
    p.font.color.rgb = text_light
    p.alignment = PP_ALIGN.CENTER

    # Save
    output_path = Path(__file__).parent.parent / "templates" / "minimal" / "basic.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"Created {output_path}")


if __name__ == "__main__":
    create_basic_template()
