#!/usr/bin/env python3
"""Run programmatic quality checks on a generated PowerPoint presentation.

Checks for: text overflow, font consistency, readability (small fonts), pacing.

Usage:
    python scripts/quality_check.py <pptx_file> [--design-system design_systems/default.yaml]
"""

import argparse
import json
import sys
from pathlib import Path

# Add project root to path
sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from pptx import Presentation
from pptx.util import Pt

from src.schemas.design_system import DesignSystem


def _hex_luminance(hex_color: str) -> float:
    """Calculate relative luminance of a hex color (0.0 = black, 1.0 = white)."""
    hex_color = hex_color.lstrip("#")
    if len(hex_color) < 6:
        return 0.5  # Unknown
    try:
        r = int(hex_color[0:2], 16) / 255
        g = int(hex_color[2:4], 16) / 255
        b = int(hex_color[4:6], 16) / 255
        return 0.299 * r + 0.587 * g + 0.114 * b
    except (ValueError, IndexError):
        return 0.5


def _is_slide_background_dark(slide) -> bool | None:
    """Determine if a slide's effective background is dark.

    Checks slide-level background first, then layout, then master.
    Returns True (dark), False (light), or None (can't determine).
    """
    try:
        from lxml import etree

        ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
        ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"

        # Check slide-level background
        bg_elem = slide.background._element if slide.background else None
        if bg_elem is not None and len(bg_elem) > 0:
            color = _extract_bg_color(bg_elem, ns_p, ns_a)
            if color is not None:
                return _hex_luminance(color) < 0.4

        # Check layout-level background
        try:
            layout = slide.slide_layout
            if layout:
                layout_xml = layout._element
                for bg in layout_xml.iter(f"{{{ns_p}}}bg"):
                    color = _extract_bg_color(bg, ns_p, ns_a)
                    if color is not None:
                        return _hex_luminance(color) < 0.4
        except Exception:
            pass

        return None  # Can't determine

    except Exception:
        return None


def _extract_bg_color(bg_elem, ns_p: str, ns_a: str) -> str | None:
    """Extract a hex color string from a background element, if it has a solid fill."""
    for bgPr in bg_elem.iter(f"{{{ns_p}}}bgPr"):
        for solid in bgPr.iter(f"{{{ns_a}}}solidFill"):
            for srgb in solid.iter(f"{{{ns_a}}}srgbClr"):
                return srgb.get("val")
    return None


def run_programmatic_checks(prs: Presentation, design_system: DesignSystem | None) -> list[dict]:
    """Run programmatic quality checks on the presentation."""
    issues = []

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    for slide_idx, slide in enumerate(prs.slides, 1):
        fonts_used: set[str] = set()
        colors_used: set[str] = set()

        for shape in slide.shapes:
            # Text overflow: shape extends beyond slide bounds
            if shape.left is not None and shape.width is not None:
                if (shape.left + shape.width) > slide_width * 1.05:
                    issues.append({
                        "slide_number": slide_idx,
                        "severity": "warning",
                        "category": "text_overflow",
                        "description": f"Shape '{shape.name}' extends beyond right edge",
                        "suggestion": "Reduce width or reposition the shape",
                    })

            if shape.top is not None and shape.height is not None:
                if (shape.top + shape.height) > slide_height * 1.05:
                    issues.append({
                        "slide_number": slide_idx,
                        "severity": "warning",
                        "category": "text_overflow",
                        "description": f"Shape '{shape.name}' extends beyond bottom edge",
                        "suggestion": "Reduce height or reposition the shape",
                    })

            # Collect fonts and colors
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run_obj in para.runs:
                        if run_obj.font.name:
                            fonts_used.add(run_obj.font.name)
                        try:
                            if run_obj.font.color and run_obj.font.color.rgb:
                                colors_used.add(str(run_obj.font.color.rgb))
                        except AttributeError:
                            pass  # Skip scheme colors that don't expose .rgb

                        # Readability: check for very small fonts
                        if run_obj.font.size and run_obj.font.size < Pt(10):
                            issues.append({
                                "slide_number": slide_idx,
                                "severity": "warning",
                                "category": "readability",
                                "description": f"Very small font ({run_obj.font.size.pt}pt) in '{shape.name}'",
                                "suggestion": "Increase font size to at least 12pt",
                            })

        # Font consistency check (includes weight variants as expected)
        if design_system and fonts_used:
            expected_fonts = {
                design_system.fonts.title_font,
                design_system.fonts.body_font,
            }
            # Also allow declared weight variants
            for attr in ("emphasis_font", "light_font", "quote_font"):
                val = getattr(design_system.fonts, attr, None)
                if val:
                    expected_fonts.add(val)
            unexpected = fonts_used - expected_fonts
            if unexpected:
                issues.append({
                    "slide_number": slide_idx,
                    "severity": "info",
                    "category": "font_consistency",
                    "description": f"Unexpected fonts: {', '.join(unexpected)}",
                    "suggestion": f"Use {' or '.join(sorted(expected_fonts))} for consistency",
                })

    # Shape-to-shape overlap check
    for slide_idx, slide in enumerate(prs.slides, 1):
        visible_shapes = []
        for shape in slide.shapes:
            if shape.left is None or shape.top is None:
                continue
            if shape.width is None or shape.height is None:
                continue
            # Skip shapes moved off-canvas (cleared shapes at 20" right)
            left_inches = shape.left / 914400
            if left_inches > 15:
                continue
            # Skip very small decorative shapes (badges, dots, lines)
            w_inches = shape.width / 914400
            h_inches = shape.height / 914400
            if w_inches * h_inches < 0.3:
                continue
            # Only check shapes with text content (most likely to cause readability issues)
            if not shape.has_text_frame or not shape.text_frame.text.strip():
                continue

            visible_shapes.append({
                "name": shape.name,
                "left": shape.left,
                "top": shape.top,
                "right": shape.left + shape.width,
                "bottom": shape.top + shape.height,
            })

        # Check all pairs for overlap
        for a_idx in range(len(visible_shapes)):
            for b_idx in range(a_idx + 1, len(visible_shapes)):
                a = visible_shapes[a_idx]
                b = visible_shapes[b_idx]
                # Bounding box intersection test
                if (a["left"] < b["right"] and a["right"] > b["left"] and
                        a["top"] < b["bottom"] and a["bottom"] > b["top"]):
                    # Calculate overlap area
                    overlap_w = min(a["right"], b["right"]) - max(a["left"], b["left"])
                    overlap_h = min(a["bottom"], b["bottom"]) - max(a["top"], b["top"])
                    overlap_area = (overlap_w / 914400) * (overlap_h / 914400)
                    if overlap_area > 0.2:  # Only report significant overlaps (>0.2 sq inches)
                        issues.append({
                            "slide_number": slide_idx,
                            "severity": "warning",
                            "category": "shape_overlap",
                            "description": (
                                f"Text shapes '{a['name']}' and '{b['name']}' overlap "
                                f"({overlap_area:.1f} sq in)"
                            ),
                            "suggestion": "Reduce text, resize shapes, or reposition to avoid overlap",
                        })

    # Background-text contrast check
    for slide_idx, slide in enumerate(prs.slides, 1):
        bg_is_dark = _is_slide_background_dark(slide)
        if bg_is_dark is None:
            continue  # Can't determine background

        dark_text_runs = 0
        light_text_runs = 0
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run_obj in para.runs:
                    if not run_obj.text.strip():
                        continue
                    try:
                        if run_obj.font.color and run_obj.font.color.rgb:
                            rgb = str(run_obj.font.color.rgb)
                            luminance = _hex_luminance(rgb)
                            if luminance < 0.4:
                                dark_text_runs += 1
                            else:
                                light_text_runs += 1
                    except (AttributeError, TypeError):
                        pass

        total_runs = dark_text_runs + light_text_runs
        if total_runs < 2:
            continue

        # Dark bg + mostly dark text = likely invisible
        if bg_is_dark and dark_text_runs > light_text_runs * 2:
            issues.append({
                "slide_number": slide_idx,
                "severity": "warning",
                "category": "background_contrast",
                "description": (
                    f"Dark background with predominantly dark text "
                    f"({dark_text_runs} dark vs {light_text_runs} light runs)"
                ),
                "suggestion": "Check if background color was preserved correctly from the template",
            })
        # Light bg + mostly light text = likely invisible
        elif not bg_is_dark and light_text_runs > dark_text_runs * 2:
            issues.append({
                "slide_number": slide_idx,
                "severity": "warning",
                "category": "background_contrast",
                "description": (
                    f"Light background with predominantly light text "
                    f"({light_text_runs} light vs {dark_text_runs} dark runs)"
                ),
                "suggestion": "Check if background color was preserved correctly from the template",
            })

    # Pacing check
    slide_count = len(prs.slides)
    if slide_count < 3:
        issues.append({
            "slide_number": 0,
            "severity": "info",
            "category": "pacing",
            "description": f"Very short deck ({slide_count} slides)",
            "suggestion": "Consider adding more context or detail",
        })
    elif slide_count > 30:
        issues.append({
            "slide_number": 0,
            "severity": "warning",
            "category": "pacing",
            "description": f"Very long deck ({slide_count} slides)",
            "suggestion": "Consider condensing or splitting into multiple decks",
        })

    return issues


def verify_content(prs: Presentation, deck_schema_path: Path | None) -> list[dict]:
    """Verify that generated slide content matches the deck schema.

    Checks that each slide's title text actually appears in the output.
    This catches the old bug where template content replaced intended content.
    """
    issues = []

    if not deck_schema_path or not deck_schema_path.exists():
        return issues

    deck_data = json.loads(deck_schema_path.read_text(encoding="utf-8"))
    slides_spec = deck_data.get("slides", [])

    if len(prs.slides) != len(slides_spec):
        issues.append({
            "slide_number": 0,
            "severity": "error",
            "category": "content_verification",
            "description": (
                f"Slide count mismatch: PPTX has {len(prs.slides)} slides, "
                f"deck schema has {len(slides_spec)} slides"
            ),
            "suggestion": "Regenerate the presentation",
        })
        return issues

    for i, (slide, spec) in enumerate(zip(prs.slides, slides_spec), 1):
        # Collect all text from the slide
        all_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text += " " + shape.text_frame.text

        all_text_lower = all_text.lower().strip()

        # Check if the title appears in the slide
        expected_title = spec.get("title", "")
        if expected_title:
            # Check for at least the first 30 chars of the title
            title_check = expected_title[:30].lower().strip()
            if title_check and title_check not in all_text_lower:
                # This slide's content doesn't match the schema
                match_type = "unknown"
                # Could be a drop-in slide — check
                actual_preview = all_text[:80].strip()
                issues.append({
                    "slide_number": i,
                    "severity": "warning",
                    "category": "content_verification",
                    "description": (
                        f"Expected title '{expected_title[:50]}' not found. "
                        f"Slide text starts with: '{actual_preview}'"
                    ),
                    "suggestion": "Check if this is a drop-in slide or if content population failed",
                })

        # Check speaker notes
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            expected_notes = spec.get("speaker_notes", "")
            if expected_notes and not notes:
                issues.append({
                    "slide_number": i,
                    "severity": "info",
                    "category": "content_verification",
                    "description": "Speaker notes missing",
                    "suggestion": "Speaker notes from deck schema were not applied",
                })
        except Exception:
            pass  # Notes slide may not exist

    return issues


def verify_file_integrity(pptx_path: Path) -> list[dict]:
    """Verify the PPTX file can be re-opened cleanly (catches encoding issues)."""
    issues = []
    try:
        # Re-open the file to check it's valid
        prs = Presentation(str(pptx_path))
        _ = len(prs.slides)  # Force access
    except Exception as e:
        issues.append({
            "slide_number": 0,
            "severity": "error",
            "category": "file_integrity",
            "description": f"File cannot be re-opened by python-pptx: {e}",
            "suggestion": "The PPTX file may be corrupted. Regenerate without drop-in slides.",
        })

    # Check ZIP structure
    try:
        import zipfile
        with zipfile.ZipFile(str(pptx_path), "r") as zf:
            bad = zf.testzip()
            if bad:
                issues.append({
                    "slide_number": 0,
                    "severity": "error",
                    "category": "file_integrity",
                    "description": f"Corrupt ZIP entry: {bad}",
                    "suggestion": "The PPTX file is corrupted.",
                })
    except Exception as e:
        issues.append({
            "slide_number": 0,
            "severity": "error",
            "category": "file_integrity",
            "description": f"Cannot read as ZIP: {e}",
            "suggestion": "The PPTX file may not be a valid ZIP archive.",
        })

    return issues


def describe_deck(prs: Presentation) -> str:
    """Create a text description of the deck for review."""
    parts = [f"Presentation: {len(prs.slides)} slides\n"]

    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_parts = [f"### Slide {slide_idx}"]
        text_items = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    text_items.append(f"  [{shape.name}]: {text[:200]}")

        if text_items:
            slide_parts.extend(text_items)
        else:
            slide_parts.append("  (no text content)")

        parts.append("\n".join(slide_parts))

    return "\n\n".join(parts)


def main():
    parser = argparse.ArgumentParser(description="Run quality checks on a PPTX file")
    parser.add_argument("pptx_file", type=Path, help="Path to the .pptx file to check")
    parser.add_argument("--design-system", type=Path, default=None,
                        help="Optional design system YAML for consistency checks")
    parser.add_argument("--deck-schema", type=Path, default=None,
                        help="Optional deck_schema.json for content verification")
    parser.add_argument("-o", "--output", type=Path, default=Path("workspace/quality_report.json"),
                        help="Output JSON report path")
    parser.add_argument("--describe", action="store_true",
                        help="Also output a text description of the deck")
    args = parser.parse_args()

    if not args.pptx_file.exists():
        print(f"Error: PPTX file not found: {args.pptx_file}", file=sys.stderr)
        sys.exit(1)

    # Step 1: File integrity check
    integrity_issues = verify_file_integrity(args.pptx_file)
    if any(i["severity"] == "error" for i in integrity_issues):
        print("ERROR: File integrity check failed!")
        for issue in integrity_issues:
            print(f"  [{issue['severity'].upper()}] {issue['description']}")
        # Still continue with other checks if possible

    # Load presentation
    prs = Presentation(str(args.pptx_file))

    # Load design system if provided
    design_system = None
    if args.design_system and args.design_system.exists():
        design_system = DesignSystem.from_yaml(args.design_system)

    # Step 2: Programmatic quality checks
    issues = run_programmatic_checks(prs, design_system)

    # Step 3: Content verification against deck schema
    content_issues = verify_content(prs, args.deck_schema)

    # Merge all issues
    all_issues = integrity_issues + issues + content_issues

    # Build report
    report = {
        "file": str(args.pptx_file),
        "slide_count": len(prs.slides),
        "issue_count": len(all_issues),
        "integrity_issues": len(integrity_issues),
        "quality_issues": len(issues),
        "content_issues": len(content_issues),
        "issues": all_issues,
    }

    # Ensure output directory exists
    args.output.parent.mkdir(parents=True, exist_ok=True)

    # Write report
    args.output.write_text(json.dumps(report, indent=2), encoding="utf-8")
    print(f"Quality report written to: {args.output}")
    print(f"Slide count: {len(prs.slides)}")
    print(f"Issues: {len(all_issues)} total ({len(integrity_issues)} integrity, "
          f"{len(issues)} quality, {len(content_issues)} content)")

    if all_issues:
        for issue in all_issues:
            severity = issue["severity"].upper()
            print(f"  [{severity}] Slide {issue['slide_number']}: {issue['description']}")

    # Optionally describe the deck
    if args.describe:
        description = describe_deck(prs)
        desc_path = args.output.with_suffix(".txt")
        desc_path.write_text(description, encoding="utf-8")
        print(f"\nDeck description written to: {desc_path}")


if __name__ == "__main__":
    main()
