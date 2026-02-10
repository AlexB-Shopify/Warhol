"""Base composer providing shared decoration and layout methods.

All slide-type composers inherit from BaseComposer and implement
the compose() method to build their specific slide layout.
"""

import logging
from abc import ABC, abstractmethod

from lxml import etree

from pptx.util import Inches, Pt

from src.schemas.design_system import DesignSystem
from src.schemas.slide_schema import SlideSpec
from src.pptx_engine.shape_operations import (
    add_accent_bar,
    add_line,
    add_rectangle,
)
from src.pptx_engine.text_operations import (
    add_label,
    add_multi_format_textbox,
    add_textbox,
)

logger = logging.getLogger(__name__)


class BaseComposer(ABC):
    """Abstract base for all slide composers.

    Subclasses implement compose() to build a specific slide type.
    The base class provides shared helpers for backgrounds, dividers,
    accent bars, section markers, footers, and other decorative elements.
    """

    @abstractmethod
    def compose(self, slide, spec: SlideSpec, design: DesignSystem) -> None:
        """Build the slide content and visual elements.

        Args:
            slide: A blank slide (already added to the presentation).
            spec: The slide specification from the deck schema.
            design: The design system configuration.
        """
        ...

    # ------------------------------------------------------------------
    # Background
    # ------------------------------------------------------------------

    @staticmethod
    def set_background(slide, hex_color: str, force: bool = False) -> None:
        """Apply a solid color background to a slide via XML injection.

        When the slide was created from a branded layout (non-blank), the
        layout/master provides rich backgrounds (images, gradients). In that
        case, we skip the solid-color override to preserve the template design.

        Set force=True to always apply the solid color (e.g., on blank slides).
        """
        if not force:
            try:
                layout = slide.slide_layout
                # If the layout has placeholders, it's a branded layout —
                # don't override its background with a flat solid color.
                if layout and len(layout.placeholders) > 0:
                    return
            except Exception:
                pass  # Can't determine layout — proceed with background set

        try:
            ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
            ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"

            cSld = slide.element
            # Remove existing background
            for existing_bg in cSld.findall(f"{{{ns_p}}}bg"):
                cSld.remove(existing_bg)

            hex_val = hex_color.lstrip("#")
            if len(hex_val) > 6:
                hex_val = hex_val[:6]

            bg = etree.SubElement(cSld, f"{{{ns_p}}}bg")
            bgPr = etree.SubElement(bg, f"{{{ns_p}}}bgPr")
            solidFill = etree.SubElement(bgPr, f"{{{ns_a}}}solidFill")
            srgbClr = etree.SubElement(solidFill, f"{{{ns_a}}}srgbClr")
            srgbClr.set("val", hex_val.upper())
            etree.SubElement(bgPr, f"{{{ns_a}}}effectLst")

            # p:bg must be first child of cSld (before spTree)
            cSld.remove(bg)
            cSld.insert(0, bg)
        except Exception as e:
            logger.warning(f"Could not set background color: {e}")

    # ------------------------------------------------------------------
    # Decorative elements
    # ------------------------------------------------------------------

    def add_divider_line(
        self,
        slide,
        design: DesignSystem,
        x1: float,
        y: float,
        x2: float,
    ) -> None:
        """Add a horizontal divider line."""
        try:
            add_line(
                slide,
                start_x=x1,
                start_y=y,
                end_x=x2,
                end_y=y,
                color=design.divider_line_color_resolved,
                width=design.decoration.divider_line_width,
            )
        except Exception as e:
            logger.debug(f"Could not add divider line: {e}")

    def add_vertical_divider(
        self,
        slide,
        design: DesignSystem,
        x: float,
        y1: float,
        y2: float,
    ) -> None:
        """Add a vertical divider line."""
        try:
            add_line(
                slide,
                start_x=x,
                start_y=y1,
                end_x=x,
                end_y=y2,
                color=design.divider_line_color_resolved,
                width=design.decoration.divider_line_width,
            )
        except Exception as e:
            logger.debug(f"Could not add vertical divider: {e}")

    def add_accent_element(
        self,
        slide,
        design: DesignSystem,
        left: float | None = None,
        top: float | None = None,
        width: float | None = None,
    ) -> None:
        """Add the standard accent bar at a given position.

        If positions are not provided, places it near the bottom-left.
        """
        dims = self.get_dims(slide)
        _left = left if left is not None else design.content_area.margin_left
        _top = top if top is not None else dims[1] - 0.55
        _width = width if width is not None else dims[0] * 0.15

        try:
            add_accent_bar(
                slide,
                left=_left,
                top=_top,
                width=_width,
                height=design.decoration.accent_bar_height,
                color=design.accent_bar_color_resolved,
            )
        except Exception as e:
            logger.debug(f"Could not add accent bar: {e}")

    def add_section_marker(
        self,
        slide,
        design: DesignSystem,
        section_number: int | str,
        section_label: str,
        left: float | None = None,
        top: float | None = None,
        color: str | None = None,
    ) -> None:
        """Add a '01 | Topic Name' style section marker.

        Uses multi-format text: number in one style, label in another.
        """
        dims = self.get_dims(slide)
        _left = left if left is not None else design.content_area.margin_left
        _top = top if top is not None else dims[1] - 0.4
        _color = color or design.dark_slide_text

        num_str = f"{section_number:02d}" if isinstance(section_number, int) else str(section_number)

        try:
            add_multi_format_textbox(
                slide,
                runs=[
                    {
                        "text": num_str,
                        "font_name": design.label_font_resolved,
                        "font_size": design.section_marker_size_resolved,
                        "font_color": design.accent_bar_color_resolved,
                        "bold": False,
                    },
                    {
                        "text": "  |  ",
                        "font_name": design.label_font_resolved,
                        "font_size": design.section_marker_size_resolved,
                        "font_color": _color,
                        "bold": False,
                    },
                    {
                        "text": section_label,
                        "font_name": design.label_font_resolved,
                        "font_size": design.section_marker_size_resolved,
                        "font_color": _color,
                        "bold": False,
                    },
                ],
                left=_left,
                top=_top,
                width=3.0,
                height=0.25,
                alignment="left",
            )
        except Exception as e:
            logger.debug(f"Could not add section marker: {e}")

    def add_slide_footer(
        self,
        slide,
        design: DesignSystem,
        text: str = "Confidential",
        color: str | None = None,
    ) -> None:
        """Add a footer label (confidential notice, page number, etc.)."""
        dims = self.get_dims(slide)
        _color = color or design.dark_slide_text

        try:
            add_label(
                slide,
                text=text,
                left=dims[0] - 2.0,
                top=dims[1] - 0.35,
                width=1.7,
                height=0.2,
                font_name=design.label_font_resolved,
                font_size=7,
                font_color=_color,
                alignment="right",
            )
        except Exception as e:
            logger.debug(f"Could not add footer: {e}")

    # ------------------------------------------------------------------
    # Content extraction helpers
    # ------------------------------------------------------------------

    @staticmethod
    def get_body_text(spec: SlideSpec) -> str:
        """Extract body text from content blocks."""
        parts = []
        for block in spec.content_blocks:
            if block.type in ("body", "caption", "data_point", "subtitle"):
                parts.append(block.content)
        return "\n\n".join(parts)

    @staticmethod
    def get_bullets(spec: SlideSpec) -> list[str]:
        """Extract bullet items from content blocks."""
        bullets = []
        for block in spec.content_blocks:
            if block.type == "bullets":
                for line in block.content.split("\n"):
                    line = line.strip()
                    for prefix in ("- ", "* ", "\u2022 ", "\u00b7 "):
                        if line.startswith(prefix):
                            line = line[len(prefix):]
                            break
                    if len(line) > 2 and line[0].isdigit() and line[1] in (".", ")"):
                        line = line[2:].strip()
                    elif len(line) > 3 and line[:2].isdigit() and line[2] in (".", ")"):
                        line = line[3:].strip()
                    if line:
                        bullets.append(line)
        return bullets

    @staticmethod
    def get_quote_text(spec: SlideSpec) -> str:
        """Extract quote text from content blocks."""
        for block in spec.content_blocks:
            if block.type == "quote":
                return block.content
        # Fallback to first content block
        if spec.content_blocks:
            return spec.content_blocks[0].content
        return ""

    @staticmethod
    def get_data_point(spec: SlideSpec) -> tuple[str, str]:
        """Extract data point number and context from content blocks.

        Returns:
            Tuple of (number_text, context_text).
        """
        number = ""
        context = ""
        for block in spec.content_blocks:
            if block.type == "data_point":
                number = block.content
            elif block.type in ("body", "caption"):
                context = block.content
        if not number and spec.title:
            number = spec.title
        return number, context

    # ------------------------------------------------------------------
    # Utility
    # ------------------------------------------------------------------

    @staticmethod
    def get_dims(slide) -> tuple[float, float]:
        """Return (width, height) in inches for the slide's presentation."""
        try:
            prs = slide.part.package.presentation
            return (prs.slide_width.inches, prs.slide_height.inches)
        except Exception:
            return (10.0, 5.625)

    @staticmethod
    def infer_section_number(spec: SlideSpec) -> int:
        """Guess a section number from the slide number (simple heuristic)."""
        # Section headers reset the count; content slides increment
        # For now, just derive from slide number
        return max(1, (spec.slide_number - 1) // 3 + 1)
