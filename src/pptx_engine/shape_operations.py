"""Shape manipulation operations for PowerPoint slides.

Provides primitives for adding rectangles, lines, ovals, badges,
image placeholder areas, accent bars, and card backgrounds.
"""

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


def add_rectangle(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_color: str | None = None,
    border_color: str | None = None,
    border_width: float = 1.0,
    corner_radius: float | None = None,
) -> object:
    """Add a rectangle shape to a slide.

    Args:
        slide: The slide to add the shape to.
        left, top, width, height: Position and size in inches.
        fill_color: Fill color as hex string (e.g., "#1a73e8"). None for no fill.
        border_color: Border color as hex string. None for no border.
        border_width: Border width in points.
        corner_radius: Corner radius for rounded rectangle in inches.

    Returns:
        The created shape.
    """
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE

    shape = slide.shapes.add_shape(
        shape_type,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )

    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(fill_color)
    else:
        shape.fill.background()

    if border_color:
        shape.line.color.rgb = _hex_to_rgb(border_color)
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()

    return shape


def add_line(
    slide,
    start_x: float,
    start_y: float,
    end_x: float,
    end_y: float,
    color: str = "#5f6368",
    width: float = 1.0,
) -> object:
    """Add a line shape to a slide.

    Args:
        slide: The slide to add the line to.
        start_x, start_y: Start position in inches.
        end_x, end_y: End position in inches.
        color: Line color as hex string.
        width: Line width in points.

    Returns:
        The created connector shape.
    """

    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR_TYPE.STRAIGHT
        Inches(start_x),
        Inches(start_y),
        Inches(end_x),
        Inches(end_y),
    )
    connector.line.color.rgb = _hex_to_rgb(color)
    connector.line.width = Pt(width)

    return connector


def add_oval(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_color: str | None = None,
    border_color: str | None = None,
) -> object:
    """Add an oval/circle shape to a slide.

    For a circle, set width == height.
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )

    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(fill_color)
    else:
        shape.fill.background()

    if border_color:
        shape.line.color.rgb = _hex_to_rgb(border_color)
    else:
        shape.line.fill.background()

    return shape


# ---------------------------------------------------------------------------
# Composite shape helpers
# ---------------------------------------------------------------------------


def add_badge_shape(
    slide,
    text: str,
    left: float,
    top: float,
    size: float = 0.4,
    fill_color: str = "#CDF986",
    text_color: str = "#191E17",
    font_name: str = "Inter Tight SemiBold",
    font_size: int = 11,
) -> object:
    """Add a badge (rounded rectangle with centered text) to a slide.

    Used for numbered list badges like '01', '02', step indicators, etc.

    Args:
        slide: Target slide.
        text: Badge label (e.g., '01').
        left, top: Position in inches.
        size: Badge width and height in inches.
        fill_color: Badge background color.
        text_color: Text color inside the badge.
        font_name: Font for the badge label.
        font_size: Font size in points.

    Returns:
        The created shape.
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(size),
        Inches(size),
    )

    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(fill_color)
    shape.line.fill.background()

    # Add centered text
    tf = shape.text_frame
    tf.word_wrap = False
    try:
        tf.auto_size = None
    except Exception:
        pass

    # Vertical centering
    try:
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        shape.text_frame.paragraphs[0].space_before = Pt(0)
        shape.text_frame.paragraphs[0].space_after = Pt(0)
    except Exception:
        pass

    # Set vertical anchor to middle
    try:
        from lxml import etree
        txBody = shape._element.find(
            ".//{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr"
        )
        if txBody is not None:
            txBody.set("anchor", "ctr")
    except Exception:
        pass

    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = _hex_to_rgb(text_color)
        run.font.bold = False

    return shape


def add_image_placeholder(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_color: str = "#1B1B1B",
    label: str | None = None,
    corner_radius: float | None = 0.1,
) -> object:
    """Add an image placeholder area (colored rectangle marking where an image goes).

    Args:
        slide: Target slide.
        left, top, width, height: Position and size in inches.
        fill_color: Background fill for the placeholder area.
        label: Optional label text (e.g., 'Image') shown inside.
        corner_radius: Corner radius in inches. None for sharp corners.

    Returns:
        The created shape.
    """
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )

    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(fill_color)
    shape.line.fill.background()

    if label:
        tf = shape.text_frame
        tf.word_wrap = True
        try:
            tf.auto_size = None
        except Exception:
            pass

        # Center the label vertically and horizontally
        try:
            from lxml import etree
            txBody = shape._element.find(
                ".//{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr"
            )
            if txBody is not None:
                txBody.set("anchor", "ctr")
        except Exception:
            pass

        p = tf.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.size = Pt(9)
            run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    return shape


def add_accent_bar(
    slide,
    left: float,
    top: float,
    width: float,
    height: float = 0.06,
    color: str = "#CDF986",
) -> object:
    """Add a thin accent bar (rectangle) for visual emphasis.

    Args:
        slide: Target slide.
        left, top: Position in inches.
        width: Bar width in inches.
        height: Bar height/thickness in inches.
        color: Bar fill color.

    Returns:
        The created shape.
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(color)
    shape.line.fill.background()
    return shape


def add_card_background(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_color: str = "#F4F4F4",
    corner_radius: float = 0.15,
) -> object:
    """Add a card background (rounded rectangle behind text content).

    Used for card-style layouts, column backgrounds, feature boxes, etc.

    Args:
        slide: Target slide.
        left, top, width, height: Position and size in inches.
        fill_color: Card background color.
        corner_radius: Corner radius in inches.

    Returns:
        The created shape.
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(fill_color)
    shape.line.fill.background()
    return shape


# ---------------------------------------------------------------------------
# Shape modification helpers
# ---------------------------------------------------------------------------


def set_shape_fill(shape, color: str) -> None:
    """Set the fill color of an existing shape."""
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(color)


def set_shape_border(shape, color: str, width: float = 1.0) -> None:
    """Set the border of an existing shape."""
    shape.line.color.rgb = _hex_to_rgb(color)
    shape.line.width = Pt(width)


def remove_shape_fill(shape) -> None:
    """Remove fill from a shape (make transparent)."""
    shape.fill.background()


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor."""
    hex_color = hex_color.lstrip("#")
    if len(hex_color) > 6:
        hex_color = hex_color[:6]
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )
