"""Core slide manipulation operations using python-pptx.

Supports two modes:
1. Layout-based: Add slides from the base template's slide layouts (for generated content)
2. Clone-based: Clone entire slides from source templates as-is (for drop-in slides)

The clone-based mode uses OPC-level part copying to preserve all visual elements:
- Slide layout backgrounds, decorative shapes, and branded graphics
- Master slide inheritance chain
- Image and media relationships
- Theme colors and fonts from the layout
"""

import copy
import logging
import re
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.package import Part, XmlPart
from pptx.opc.packuri import PackURI
from pptx.parts.slide import SlideLayoutPart, SlideMasterPart
from pptx.slide import SlideLayout
from pptx.util import Inches

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Caches — avoid re-opening source files and re-importing layouts
# ---------------------------------------------------------------------------

_source_prs_cache: dict[str, Presentation] = {}
_imported_layout_cache: dict[tuple[str, str], SlideLayout] = {}
_imported_master_cache: dict[str, object] = {}  # source_master_partname → imported master part
_allocated_partnames: set[str] = set()  # tracks names allocated this run


def clear_clone_caches() -> None:
    """Clear cached source presentations and imported layouts.

    Call between build runs or when memory should be freed.
    """
    _source_prs_cache.clear()
    _imported_layout_cache.clear()
    _imported_master_cache.clear()
    _allocated_partnames.clear()


def _get_source_prs(source_path: Path) -> Presentation:
    """Open a source PPTX, caching to avoid repeated I/O."""
    key = str(source_path.resolve())
    if key not in _source_prs_cache:
        _source_prs_cache[key] = Presentation(str(source_path))
    return _source_prs_cache[key]


# ---------------------------------------------------------------------------
# Presentation creation
# ---------------------------------------------------------------------------

def open_base_template(base_path: str | Path) -> Presentation:
    """Open a base template as the starting presentation.

    Loads the template for its slide masters, layouts, and theme (backgrounds,
    fonts, colors). Any existing slides in the template are removed so the
    output starts clean — only the layouts and master design are kept.
    """
    base_path = Path(base_path)
    if not base_path.exists():
        raise FileNotFoundError(f"Base template not found: {base_path}")

    prs = Presentation(str(base_path))

    # Strip existing slides — we only want the layouts / masters / theme
    existing_count = len(prs.slides)
    if existing_count > 0:
        logger.info(
            f"Stripping {existing_count} existing slides from base template "
            f"(keeping layouts and theme)"
        )
        _strip_existing_slides(prs)

    logger.info(
        f"Opened base template: {base_path.name} "
        f"({prs.slide_width.inches:.1f}\"x{prs.slide_height.inches:.3f}\", "
        f"{len(prs.slide_layouts)} layouts, {existing_count} slides stripped)"
    )
    return prs


def _strip_existing_slides(prs: Presentation) -> None:
    """Remove all existing slides from a presentation, preserving layouts/masters.

    This allows a slide bank (full of example slides) to be used as a base
    template — we inherit the rich design system but start with a blank deck.
    """
    slide_list = prs.slides._sldIdLst
    slide_ids = list(slide_list)
    for sld_id in slide_ids:
        r_id = sld_id.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        if r_id:
            try:
                prs.part.drop_rel(r_id)
            except Exception:
                pass  # Relationship may already be gone
        slide_list.remove(sld_id)


def create_presentation(
    width_inches: float = 10.0,
    height_inches: float = 5.625,
) -> Presentation:
    """Create a new blank 16:9 presentation (fallback when no base template)."""
    prs = Presentation()
    prs.slide_width = Inches(width_inches)
    prs.slide_height = Inches(height_inches)
    return prs


# ---------------------------------------------------------------------------
# Layout-based slide creation (for generated content)
# ---------------------------------------------------------------------------

# Map of layout names to their index in the Shopify base template
LAYOUT_MAP = {
    "title": 0,               # TITLE: CENTER_TITLE + SUBTITLE
    "section_header": 7,       # SECTION_TITLE_AND_DESCRIPTION: TITLE + SUBTITLE + BODY
    "content": 2,              # TITLE_AND_BODY: TITLE + BODY
    "bullet_list": 2,          # TITLE_AND_BODY: TITLE + BODY
    "two_column": 3,           # TITLE_AND_TWO_COLUMNS: TITLE + BODY + BODY
    "comparison": 3,           # TITLE_AND_TWO_COLUMNS: TITLE + BODY + BODY
    "quote": 6,                # MAIN_POINT: TITLE only (we add quote as textbox)
    "chart": 4,                # TITLE_ONLY: TITLE (chart area manual)
    "image_full": 10,          # BLANK
    "image_with_text": 2,      # TITLE_AND_BODY: TITLE + BODY
    "timeline": 2,             # TITLE_AND_BODY: TITLE + BODY
    "team": 4,                 # TITLE_ONLY: TITLE
    "closing": 0,              # TITLE: CENTER_TITLE + SUBTITLE
    "data_point": 9,           # BIG_NUMBER: TITLE + BODY
}


def add_slide_from_layout(prs: Presentation, slide_type: str) -> object:
    """Add a new slide using the appropriate layout for the given slide type.

    Args:
        prs: The presentation (should be opened from base template).
        slide_type: The slide type string (e.g., 'title', 'content', 'section_header').

    Returns:
        The newly added slide with layout placeholders ready to populate.
    """
    layout_idx = LAYOUT_MAP.get(slide_type, 2)  # Default to TITLE_AND_BODY

    layouts = prs.slide_layouts
    if layout_idx >= len(layouts):
        layout_idx = min(2, len(layouts) - 1)  # Fallback to content layout
        logger.warning(f"Layout index {layout_idx} out of range, using fallback")

    layout = layouts[layout_idx]
    slide = prs.slides.add_slide(layout)
    logger.debug(f"Added slide with layout '{layout.name}' for type '{slide_type}'")
    return slide


def add_blank_slide(prs: Presentation) -> object:
    """Add a blank slide (no placeholders except slide number)."""
    blank_idx = 10  # BLANK layout
    layouts = prs.slide_layouts
    if blank_idx >= len(layouts):
        blank_idx = len(layouts) - 1
    return prs.slides.add_slide(layouts[blank_idx])


# ---------------------------------------------------------------------------
# Clone-based slide insertion (OPC-level part copying)
# ---------------------------------------------------------------------------

def clone_slide_as_is(
    target_prs: Presentation,
    source_path: str | Path,
    slide_index: int,
) -> object:
    """Clone a slide from a source template into the target, preserving ALL design.

    Uses OPC-level part copying to preserve the full visual hierarchy:
    1. Imports the source slide's layout (with branded backgrounds, decorative
       shapes, logos, etc.) into the target presentation
    2. Creates a new slide from that imported layout — inheriting all design
    3. Copies the source slide's own shapes (text, images, groups)
    4. Copies any slide-level background overrides

    This ensures master backgrounds, layout graphics, and theme colors all
    carry over from the source template.

    Args:
        target_prs: The presentation to add the cloned slide to.
        source_path: Path to the source .pptx template file.
        slide_index: Index of the slide to clone from the source.

    Returns:
        The newly created slide in the target presentation.
    """
    source_path = Path(source_path)
    source_prs = _get_source_prs(source_path)

    if slide_index >= len(source_prs.slides):
        raise IndexError(
            f"Slide index {slide_index} out of range for {source_path.name} "
            f"(has {len(source_prs.slides)} slides)"
        )

    _check_dimensions(target_prs, source_prs, source_path.name)

    source_slide = source_prs.slides[slide_index]

    # --- Step 1: Import the source slide's layout into the target ---
    imported_layout = _import_slide_layout(target_prs, source_slide, source_prs)

    # --- Step 2: Add a new slide using the imported layout ---
    # This gives us a slide that inherits the layout's backgrounds, branded
    # graphics, decorative shapes, and master theme — automatically.
    new_slide = target_prs.slides.add_slide(imported_layout)

    # --- CRITICAL: Get the REAL spTree from the Part's XML element ---
    # python-pptx's slide.shapes._spTree can become a DETACHED element
    # after clone_layout_placeholders runs (the descriptor returns a stale
    # reference). We must work directly with the Part's XML tree to ensure
    # modifications are serialized when saving.
    real_spTree = new_slide.part._element.cSld.spTree

    # --- Step 3: Clear auto-created placeholder shapes ---
    _clear_spTree(real_spTree)

    # --- Step 4: Copy all shapes from the source slide ---
    target_package = target_prs.part.package
    for shape in source_slide.shapes:
        _clone_shape_to_spTree(
            shape, real_spTree, source_slide, new_slide, target_package
        )

    # --- Step 5: Copy slide-level background overrides (if any) ---
    _copy_slide_background(source_slide, new_slide, target_package)

    # --- Step 6: Invalidate the stale shapes proxy cache ---
    # python-pptx's @lazyproperty on Slide.shapes caches a SlideShapes
    # proxy pointing to the ORIGINAL spTree from CT_Slide.new().
    # clone_layout_placeholders may replace that spTree element, leaving
    # the proxy detached. Deleting the cache forces re-creation from the
    # Part's current XML on next access, so text replacement works.
    try:
        if "shapes" in new_slide.__dict__:
            del new_slide.__dict__["shapes"]
    except Exception:
        pass

    return new_slide


def _check_dimensions(target_prs, source_prs, source_name: str) -> None:
    """Warn if source and target have different slide dimensions."""
    tw, th = target_prs.slide_width, target_prs.slide_height
    sw, sh = source_prs.slide_width, source_prs.slide_height

    if tw and sw and (abs(tw - sw) > 100000 or abs(th - sh) > 100000):
        logger.warning(
            f"Dimension mismatch: target={tw.inches:.1f}\"x{th.inches:.1f}\", "
            f"source {source_name}={sw.inches:.1f}\"x{sh.inches:.1f}\". "
            f"Drop-in slide may not fit correctly."
        )


# ---------------------------------------------------------------------------
# Layout importing — the core of the OPC-level clone approach
# ---------------------------------------------------------------------------

def _import_slide_layout(
    target_prs: Presentation,
    source_slide,
    source_prs: Presentation,
) -> SlideLayout:
    """Import the source slide's layout AND its master into the target.

    Copies the layout's XML (branded backgrounds, decorative shapes,
    placeholders) and all its media. Also imports the source's slide
    master so that theme colors (SchemeColor references) resolve to
    the correct values from the source template, not the target's
    default theme.

    Returns a SlideLayout proxy suitable for passing to add_slide().
    Results are cached so the same layout isn't imported multiple times.
    """
    source_path = str(source_prs.core_properties.title or id(source_prs))

    # Get the actual layout part from the source slide
    source_layout_part = _get_layout_part(source_slide)
    cache_key = (source_path, str(source_layout_part.partname))

    if cache_key in _imported_layout_cache:
        return _imported_layout_cache[cache_key]

    target_package = target_prs.part.package

    # --- Import the source's slide master (with theme/colors) ---
    imported_master = _import_slide_master(
        target_prs, source_layout_part, target_package
    )

    # --- Deep copy the layout XML ---
    layout_xml = copy.deepcopy(source_layout_part._element)

    # --- Create a new SlideLayoutPart in the target package ---
    new_partname = _unique_partname(target_package, source_layout_part.partname)

    new_layout_part = SlideLayoutPart(
        new_partname,
        source_layout_part.content_type,
        target_package,
        layout_xml,
    )

    # --- Link to the IMPORTED master (not the target's default) ---
    new_layout_part.relate_to(imported_master, RT.SLIDE_MASTER)

    # --- Import all non-master relationships (images, etc.) ---
    rid_map = _import_part_rels(
        new_layout_part,
        source_layout_part,
        target_package,
        skip_reltypes={RT.SLIDE_MASTER},
    )

    # --- Remap relationship IDs in the layout XML ---
    if rid_map:
        _remap_rids(layout_xml, rid_map)

    # --- Register the layout with the imported master ---
    # Both the relationship AND the sldLayoutIdLst entry are required.
    # Keynote rejects the file without the sldLayoutIdLst entry.
    master_rId = imported_master.relate_to(new_layout_part, RT.SLIDE_LAYOUT)
    _register_layout_in_master(imported_master, master_rId)

    # --- Build a SlideLayout proxy for add_slide() ---
    imported_layout = SlideLayout(layout_xml, new_layout_part)

    _imported_layout_cache[cache_key] = imported_layout
    logger.debug(
        f"Imported layout {source_layout_part.partname} → {new_partname}"
    )
    return imported_layout


def _import_slide_master(target_prs, source_layout_part, target_package):
    """Import the source layout's slide master into the target presentation.

    The master carries the theme (colors, fonts) that SchemeColor references
    in the layout and slide shapes resolve against. Importing it ensures
    colors and fonts render correctly.

    Cached per source master partname to avoid duplicating masters from
    the same source template.
    """
    # Find the source master via the layout's relationship
    source_master_part = None
    for _key, rel in source_layout_part.rels.items():
        if rel.reltype == RT.SLIDE_MASTER:
            source_master_part = rel.target_part
            break

    if source_master_part is None:
        # Fallback: use the target's first master
        logger.warning("Source layout has no master relationship, using target default")
        return _get_first_master_part(target_prs)

    # Check cache — same source master may be used by multiple layouts
    cache_key = str(source_master_part.partname)
    if cache_key in _imported_master_cache:
        return _imported_master_cache[cache_key]

    # --- Deep copy the master XML ---
    master_xml = copy.deepcopy(source_master_part._element)

    # --- Create a new SlideMasterPart ---
    new_partname = _unique_partname(target_package, source_master_part.partname)

    new_master_part = SlideMasterPart(
        new_partname,
        source_master_part.content_type,
        target_package,
        master_xml,
    )

    # --- Import non-layout relationships (theme, images, etc.) ---
    # Skip SLIDE_LAYOUT rels — we'll add our own as layouts are imported.
    rid_map = _import_part_rels(
        new_master_part,
        source_master_part,
        target_package,
        skip_reltypes={RT.SLIDE_LAYOUT},
    )

    if rid_map:
        _remap_rids(master_xml, rid_map)

    # --- Clear the master's sldLayoutIdLst (we'll populate as layouts import) ---
    ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
    layout_id_lst = master_xml.find(f"{{{ns_p}}}sldLayoutIdLst")
    if layout_id_lst is not None:
        for child in list(layout_id_lst):
            layout_id_lst.remove(child)

    # --- Register the master in the target presentation ---
    # Both the relationship AND the sldMasterIdLst entry are required.
    # Keynote rejects the file without the sldMasterIdLst entry.
    master_rId = target_prs.part.relate_to(new_master_part, RT.SLIDE_MASTER)
    _register_master_in_presentation(target_prs, master_rId)

    _imported_master_cache[cache_key] = new_master_part
    logger.debug(
        f"Imported master {source_master_part.partname} → {new_partname}"
    )
    return new_master_part


def _get_layout_part(slide) -> SlideLayoutPart:
    """Get the SlideLayoutPart for a slide via its relationships."""
    for _key, rel in slide.part.rels.items():
        if rel.reltype == RT.SLIDE_LAYOUT:
            return rel.target_part
    raise ValueError("Slide has no layout relationship")


def _get_first_master_part(prs: Presentation):
    """Get the SlideMasterPart from the first slide master."""
    for _key, rel in prs.part.rels.items():
        if rel.reltype == RT.SLIDE_MASTER:
            return rel.target_part
    raise ValueError("Presentation has no slide master")


# Counter for generating unique sldLayoutId values
_next_layout_id = 2147484000


def _register_layout_in_master(master_part, rId: str) -> None:
    """Add a <p:sldLayoutId> entry to the master's sldLayoutIdLst.

    Keynote requires every layout relationship to have a corresponding
    entry in the master XML's sldLayoutIdLst. Without this, Keynote
    rejects the file as having an invalid format.
    """
    global _next_layout_id

    ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
    ns_r = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

    master_xml = master_part._element
    layout_id_lst = master_xml.find(f"{{{ns_p}}}sldLayoutIdLst")

    if layout_id_lst is None:
        # Create the element if it doesn't exist
        layout_id_lst = etree.SubElement(master_xml, f"{{{ns_p}}}sldLayoutIdLst")

    # Check this rId isn't already registered
    for existing in layout_id_lst.findall(f"{{{ns_p}}}sldLayoutId"):
        if existing.get(f"{{{ns_r}}}id") == rId:
            return  # Already registered

    # Add new entry with a unique id
    new_entry = etree.SubElement(layout_id_lst, f"{{{ns_p}}}sldLayoutId")
    new_entry.set("id", str(_next_layout_id))
    new_entry.set(f"{{{ns_r}}}id", rId)
    _next_layout_id += 1


# Counter for generating unique sldMasterId values
_next_master_id = 2147485000


def _register_master_in_presentation(prs: Presentation, rId: str) -> None:
    """Add a <p:sldMasterId> entry to the presentation's sldMasterIdLst.

    Keynote requires every master relationship to have a corresponding
    entry in the presentation XML's sldMasterIdLst.
    """
    global _next_master_id

    ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
    ns_r = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

    pres_xml = prs.part._element
    master_id_lst = pres_xml.find(f"{{{ns_p}}}sldMasterIdLst")

    if master_id_lst is None:
        master_id_lst = etree.SubElement(pres_xml, f"{{{ns_p}}}sldMasterIdLst")

    # Check this rId isn't already registered
    for existing in master_id_lst.findall(f"{{{ns_p}}}sldMasterId"):
        if existing.get(f"{{{ns_r}}}id") == rId:
            return

    new_entry = etree.SubElement(master_id_lst, f"{{{ns_p}}}sldMasterId")
    new_entry.set("id", str(_next_master_id))
    new_entry.set(f"{{{ns_r}}}id", rId)
    _next_master_id += 1


# ---------------------------------------------------------------------------
# Relationship importing and ID remapping
# ---------------------------------------------------------------------------

def _import_part_rels(
    target_part,
    source_part,
    target_package,
    skip_reltypes: set | None = None,
) -> dict[str, str]:
    """Import all relationships from a source part into a target part.

    For each relationship in the source:
    - External rels (hyperlinks): recreated as external rels on the target
    - Internal media rels (images): media data is copied to a new Part
      with a unique partname in the target package
    - Skipped reltypes: ignored (e.g., RT.SLIDE_MASTER when importing layouts)

    Returns a {old_rId: new_rId} mapping for XML remapping.
    """
    skip = skip_reltypes or set()
    rid_map: dict[str, str] = {}

    for old_rid, rel in source_part.rels.items():
        if rel.reltype in skip:
            continue

        try:
            if rel.is_external:
                new_rid = target_part.rels.get_or_add_ext_rel(
                    rel.reltype, rel.target_ref
                )
            else:
                # Import the target part (image, media, etc.) with a unique name
                imported = _import_media_part(rel.target_part, target_package)
                new_rid = target_part.relate_to(imported, rel.reltype)

            if old_rid != new_rid:
                rid_map[old_rid] = new_rid

        except Exception as e:
            logger.debug(f"Could not import rel {old_rid} ({rel.reltype}): {e}")

    return rid_map


def _import_media_part(source_part, target_package) -> Part:
    """Copy a media part (image, etc.) into the target package.

    Creates a new Part with a unique partname to avoid collisions
    with existing media in the target presentation.
    """
    new_partname = _unique_partname(target_package, source_part.partname)
    return Part(
        new_partname,
        source_part.content_type,
        target_package,
        source_part.blob,
    )


def _remap_rids(element, rid_map: dict[str, str]) -> None:
    """Remap relationship IDs throughout an XML element tree.

    Only remaps attributes in the r: (relationships) namespace to avoid
    false matches on other attributes that happen to look like rIds.
    """
    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

    for el in element.iter():
        for attr_name in list(el.attrib.keys()):
            if attr_name.startswith(f"{{{r_ns}}}"):
                old_val = el.get(attr_name)
                if old_val in rid_map:
                    el.set(attr_name, rid_map[old_val])


def _unique_partname(package, original_partname) -> PackURI:
    """Generate a unique partname in the target package.

    If the original name is available, use it. Otherwise, increment
    the numeric suffix until a free name is found.

    Tracks allocated names in a module-level set so that two imports
    in the same build run don't collide even before parts are serialized.
    """
    # Combine package parts + names allocated earlier in this run
    existing = {str(p.partname) for p in package.iter_parts()} | _allocated_partnames
    original = str(original_partname)

    if original not in existing:
        _allocated_partnames.add(original)
        return PackURI(original)

    # Parse /ppt/media/image42.png → prefix=/ppt/media/image, ext=.png
    match = re.match(r"^(.*?)(\d+)(\.\w+)$", original)
    if match:
        prefix, _, ext = match.groups()
        idx = max(
            (
                int(m.group(1))
                for pn in existing
                if (m := re.search(r"(\d+)\.\w+$", pn))
                and pn.startswith(prefix)
            ),
            default=0,
        ) + 1
        while True:
            candidate = f"{prefix}{idx}{ext}"
            if candidate not in existing:
                _allocated_partnames.add(candidate)
                return PackURI(candidate)
            idx += 1
    else:
        # No number in name — append one
        name, ext = original.rsplit(".", 1)
        idx = 1
        while True:
            candidate = f"{name}{idx}.{ext}"
            if candidate not in existing:
                _allocated_partnames.add(candidate)
                return PackURI(candidate)
            idx += 1


# ---------------------------------------------------------------------------
# Shape clearing and cloning
# ---------------------------------------------------------------------------

_SHAPE_TAGS = {"sp", "pic", "graphicFrame", "grpSp", "cxnSp"}


def _clear_spTree(spTree) -> None:
    """Remove all shape elements from a spTree XML element.

    Preserves the spTree element itself and its nvGrpSpPr / grpSpPr
    children (required structural elements). Removes all actual shapes.
    """
    for child in list(spTree):
        local = etree.QName(child).localname
        if local in _SHAPE_TAGS:
            spTree.remove(child)


def _clone_shape_to_spTree(
    shape, target_spTree, source_slide, target_slide, target_package
) -> None:
    """Clone a shape into a specific spTree element.

    Deep-copies the shape XML, imports referenced media/hyperlinks,
    remaps relationship IDs, and appends to the target spTree.

    Uses the target spTree element directly (not slide.shapes._spTree)
    to ensure modifications end up in the Part's XML tree for serialization.
    """
    try:
        new_el = copy.deepcopy(shape._element)
        rid_map: dict[str, str] = {}

        # Collect and import image/media relationships
        _collect_image_rels(
            new_el, source_slide.part, target_slide.part, target_package, rid_map
        )

        # Collect and import hyperlink relationships
        _collect_hyperlink_rels(
            new_el, source_slide.part, target_slide.part, rid_map
        )

        # Remap all relationship IDs in the cloned element
        if rid_map:
            _remap_rids(new_el, rid_map)

        # Append to the REAL spTree (from the Part's XML), not the proxy
        target_spTree.append(new_el)

    except Exception as e:
        logger.warning(f"Could not clone shape '{shape.name}': {e}")


def _collect_image_rels(
    element, source_part, target_part, target_package, rid_map: dict
) -> None:
    """Find ALL internal relationship references in an XML element and import them.

    Scans every element in the tree for r: namespace attributes (r:embed,
    r:link, r:id on non-hyperlink elements, etc.) and imports the
    referenced parts as media. This catches:
    - a:blip (image fills, picture shapes)
    - a:blipFill (shape background images)
    - OLE object references
    - Any other internal media relationship
    """
    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"

    # Tags where r:id means a hyperlink (handled separately)
    _HYPERLINK_TAGS = {f"{{{a_ns}}}hlinkClick", f"{{{a_ns}}}hlinkHover"}

    for el in element.iter():
        # Skip hyperlink elements — those are handled by _collect_hyperlink_rels
        if el.tag in _HYPERLINK_TAGS:
            continue

        for attr_name in list(el.attrib.keys()):
            if not attr_name.startswith(f"{{{r_ns}}}"):
                continue

            r_id = el.get(attr_name)
            if not r_id or r_id in rid_map:
                continue

            try:
                source_media = source_part.related_part(r_id)
                imported = _import_media_part(source_media, target_package)
                # Determine the correct relationship type from the source
                source_rel = source_part.rels.get(r_id)
                rel_type = source_rel.reltype if source_rel else RT.IMAGE
                new_rid = target_part.relate_to(imported, rel_type)
                rid_map[r_id] = new_rid
            except (KeyError, Exception) as e:
                logger.debug(f"Could not import media rel {r_id}: {e}")


def _collect_hyperlink_rels(
    element, source_part, target_part, rid_map: dict
) -> None:
    """Find hyperlink references in an XML element and import them."""
    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

    for tag in (f"{{{a_ns}}}hlinkClick", f"{{{a_ns}}}hlinkHover"):
        for hlink in element.iter(tag):
            r_id = hlink.get(f"{{{r_ns}}}id")
            if r_id and r_id not in rid_map:
                try:
                    source_rel = source_part.rels.get(r_id)
                    if source_rel and source_rel.is_external:
                        new_rid = target_part.rels.get_or_add_ext_rel(
                            RT.HYPERLINK, source_rel.target_ref
                        )
                        rid_map[r_id] = new_rid
                except Exception as e:
                    logger.debug(f"Could not import hyperlink rel {r_id}: {e}")


# ---------------------------------------------------------------------------
# Slide-level background copying
# ---------------------------------------------------------------------------

def _copy_slide_background(source_slide, target_slide, target_package) -> None:
    """Copy an explicit slide-level background from source to target.

    Only copies if the source slide has its own background definition
    (overriding the layout). If the source inherits its background from
    the layout/master, this is a no-op — the imported layout already
    provides the correct background.
    """
    try:
        src_bg = source_slide.background
        if src_bg is None or src_bg._element is None:
            return

        bg_elem = src_bg._element
        if len(bg_elem) == 0:
            # No explicit background — inherits from layout (correct)
            return

        # Deep copy the background XML
        new_bg = copy.deepcopy(bg_elem)

        # Import any image references in the background
        rid_map: dict[str, str] = {}
        _collect_background_image_rels(
            new_bg, source_slide.part, target_slide.part, target_package, rid_map
        )
        if rid_map:
            _remap_rids(new_bg, rid_map)

        # Replace target slide's background children
        tgt_bg = target_slide.background._element
        for child in list(tgt_bg):
            tgt_bg.remove(child)
        for child in new_bg:
            tgt_bg.append(child)

    except Exception as e:
        logger.warning(f"Could not copy slide background: {e}")


def _collect_background_image_rels(
    bg_elem, source_part, target_part, target_package, rid_map: dict
) -> None:
    """Import image parts referenced in a background element."""
    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

    for blip in bg_elem.iter(f"{{{a_ns}}}blip"):
        r_embed = blip.get(f"{{{r_ns}}}embed")
        if r_embed and r_embed not in rid_map:
            try:
                source_image = source_part.related_part(r_embed)
                imported = _import_media_part(source_image, target_package)
                new_rid = target_part.relate_to(imported, RT.IMAGE)
                rid_map[r_embed] = new_rid
            except (KeyError, Exception) as e:
                logger.debug(f"Could not import background image: {e}")


# ---------------------------------------------------------------------------
# Utility
# ---------------------------------------------------------------------------

def get_slide_dimensions(prs: Presentation) -> tuple[float, float]:
    """Return (width, height) in inches for the presentation."""
    return (
        prs.slide_width.inches if prs.slide_width else 10.0,
        prs.slide_height.inches if prs.slide_height else 5.625,
    )
