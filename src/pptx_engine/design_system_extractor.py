"""Extract a DesignSystem from .pptx template files.

Analyzes fonts (with weight variants), colors, paragraph formatting,
content area positioning, slide dimensions, and backgrounds across all
slides in the template bank to derive a coherent design system.
"""

import logging
from collections import Counter
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

from src.schemas.design_system import (
    ColorConfig,
    ContentAreaConfig,
    DesignSystem,
    FontConfig,
    ParagraphConfig,
    SlideDimensions,
    SlideTypeOverrides,
)

logger = logging.getLogger(__name__)

# Colors to ignore -- near-black, near-white, and pure grays are "structural"
_NEUTRAL_COLORS = {
    "#000000", "#ffffff", "#FFFFFF",
    "#202124", "#333333", "#444444", "#555555",
    "#666666", "#777777", "#888888", "#999999",
    "#aaaaaa", "#bbbbbb", "#cccccc", "#dddddd",
    "#eeeeee", "#f0f0f0", "#f5f5f5", "#f8f8f8",
    "#f8f9fa", "#5f6368",
}

_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


# -----------------------------------------------------------------------
# Accumulator — collects all raw observations across files
# -----------------------------------------------------------------------

class _Accum:
    """Mutable accumulator for all design-system-relevant observations."""

    def __init__(self):
        # Fonts by role
        self.title_fonts: Counter[str] = Counter()
        self.body_fonts: Counter[str] = Counter()
        self.subtitle_fonts: Counter[str] = Counter()

        # Font sizes by role
        self.title_sizes: Counter[int] = Counter()
        self.subtitle_sizes: Counter[int] = Counter()
        self.body_sizes: Counter[int] = Counter()

        # Font weight variants (the full font-name including weight)
        self.emphasis_fonts: Counter[str] = Counter()   # Bold / SemiBold
        self.light_fonts: Counter[str] = Counter()      # Light / ExtraLight

        # Colors
        self.accent_colors: Counter[str] = Counter()
        self.text_dark_colors: Counter[str] = Counter()
        self.text_light_colors: Counter[str] = Counter()
        self.bg_colors: Counter[str] = Counter()
        self.theme_colors: Counter[str] = Counter()

        # Paragraph formatting
        self.title_alignments: Counter[str] = Counter()
        self.body_alignments: Counter[str] = Counter()
        self.subtitle_alignments: Counter[str] = Counter()
        self.body_line_spacings: Counter[float] = Counter()
        self.title_line_spacings: Counter[float] = Counter()
        self.space_afters: Counter[int] = Counter()   # in EMU
        self.space_befores: Counter[int] = Counter()   # in EMU

        # Content positioning (placeholder bounds)
        self.title_positions: list[tuple[float, float, float, float]] = []
        self.body_positions: list[tuple[float, float, float, float]] = []

        # Slide dimensions (from first file)
        self.slide_width: float | None = None
        self.slide_height: float | None = None


# -----------------------------------------------------------------------
# Public API
# -----------------------------------------------------------------------

def extract_design_system_from_file(
    pptx_path: str | Path,
    name: str = "Extracted",
) -> DesignSystem:
    """Derive a DesignSystem from a single .pptx file."""
    pptx_path = Path(pptx_path)
    if not pptx_path.exists():
        raise FileNotFoundError(f"File not found: {pptx_path}")
    if pptx_path.suffix.lower() != ".pptx":
        raise ValueError(f"Expected a .pptx file, got: {pptx_path.suffix}")
    return _extract_from_files([pptx_path], name=name)


def extract_design_system(
    template_dir: str | Path,
    name: str = "Extracted",
) -> DesignSystem:
    """Derive a DesignSystem by analyzing .pptx files in a directory."""
    from src.utils.file_utils import find_pptx_files

    template_dir = Path(template_dir)
    pptx_files = find_pptx_files(template_dir)

    if not pptx_files:
        logger.warning(f"No .pptx files found in {template_dir}. Returning defaults.")
        return DesignSystem(name=name)

    return _extract_from_files(pptx_files, name=name)


# -----------------------------------------------------------------------
# Core extraction
# -----------------------------------------------------------------------

def _extract_from_files(
    pptx_files: list[Path],
    name: str = "Extracted",
) -> DesignSystem:
    if not pptx_files:
        return DesignSystem(name=name)

    acc = _Accum()

    for pptx_path in pptx_files:
        try:
            _analyze_file(pptx_path, acc)
        except Exception as e:
            logger.warning(f"Skipping {pptx_path.name}: {e}")

    # Build the design system from accumulated data
    fonts = _derive_fonts(acc)
    colors = _derive_colors(acc)
    paragraph = _derive_paragraph(acc)
    content_area = _derive_content_area(acc)
    dimensions = SlideDimensions(
        width=acc.slide_width or 10.0,
        height=acc.slide_height or 5.625,
    )
    overrides = _derive_overrides(acc)

    ds = DesignSystem(
        name=name,
        fonts=fonts,
        colors=colors,
        paragraph=paragraph,
        content_area=content_area,
        dimensions=dimensions,
        overrides=overrides,
    )
    logger.info(
        f"Extracted design system '{name}': "
        f"title={fonts.title_font} ({fonts.title_size}pt), "
        f"body={fonts.body_font} ({fonts.body_size}pt), "
        f"emphasis={fonts.emphasis_font}, light={fonts.light_font}, "
        f"primary={colors.primary}, bg={colors.background}, "
        f"body_line_spacing={paragraph.body_line_spacing}"
    )
    return ds


def _analyze_file(pptx_path: Path, acc: _Accum) -> None:
    """Analyze a single .pptx file and accumulate observations."""
    prs = Presentation(str(pptx_path))

    # Capture slide dimensions from first file analyzed
    if acc.slide_width is None:
        try:
            acc.slide_width = prs.slide_width.inches
            acc.slide_height = prs.slide_height.inches
        except Exception:
            pass

    # Theme colors from slide master
    _extract_theme_colors(prs, acc.theme_colors)

    for slide in prs.slides:
        _extract_slide_background_color(slide, acc.bg_colors)

        for shape in slide.shapes:
            role = _classify_shape_context(shape)

            # Content area positioning
            if shape.is_placeholder and shape.left is not None:
                try:
                    pos = (
                        shape.left / 914400,
                        shape.top / 914400,
                        shape.width / 914400,
                        shape.height / 914400,
                    )
                    if role == "title":
                        acc.title_positions.append(pos)
                    elif role == "body":
                        acc.body_positions.append(pos)
                except Exception:
                    pass

            if not shape.has_text_frame:
                # Shape fills → accent colors
                _extract_shape_fill(shape, acc)
                continue

            for para in shape.text_frame.paragraphs:
                # --- Paragraph-level properties ---
                _extract_paragraph_formatting(para, role, acc)

                # --- Run-level properties ---
                para_font_name = None
                para_font_size = None
                try:
                    para_font_name = para.font.name
                    para_font_size = para.font.size
                except Exception:
                    pass

                for run in para.runs:
                    font_name = run.font.name or para_font_name
                    font_size = run.font.size or para_font_size

                    if font_name:
                        # Categorize by role AND by weight variant
                        if role == "title":
                            acc.title_fonts[font_name] += 1
                        elif role == "subtitle":
                            acc.subtitle_fonts[font_name] += 1
                        else:
                            acc.body_fonts[font_name] += 1

                        # Detect weight variants
                        _categorize_font_weight(font_name, acc)

                    if font_size:
                        size_pt = int(font_size.pt)
                        if role == "title":
                            acc.title_sizes[size_pt] += 1
                        elif role == "subtitle":
                            acc.subtitle_sizes[size_pt] += 1
                        else:
                            acc.body_sizes[size_pt] += 1

                    # Text colors
                    try:
                        color_rgb = run.font.color.rgb
                        if color_rgb is not None:
                            hex_color = f"#{color_rgb}"
                            if _is_dark_color(hex_color):
                                acc.text_dark_colors[hex_color] += 1
                            else:
                                acc.text_light_colors[hex_color] += 1
                    except (AttributeError, TypeError):
                        pass

                # Paragraph without runs
                if not para.runs and para.text.strip():
                    if para_font_name:
                        if role == "title":
                            acc.title_fonts[para_font_name] += 1
                        else:
                            acc.body_fonts[para_font_name] += 1
                    if para_font_size:
                        size_pt = int(para_font_size.pt)
                        if role == "title":
                            acc.title_sizes[size_pt] += 1
                        elif role == "subtitle":
                            acc.subtitle_sizes[size_pt] += 1
                        else:
                            acc.body_sizes[size_pt] += 1

            # Shape fills → accent colors
            _extract_shape_fill(shape, acc)


# -----------------------------------------------------------------------
# Paragraph formatting extraction
# -----------------------------------------------------------------------

_ALIGNMENT_MAP = {
    0: "left",     # LEFT
    1: "left",     # LEFT
    2: "center",   # CENTER
    3: "right",    # RIGHT
    4: "justify",  # JUSTIFY
}


def _extract_paragraph_formatting(para, role: str, acc: _Accum) -> None:
    """Extract alignment, line spacing, and space before/after from a paragraph."""
    # Alignment
    try:
        if para.alignment is not None:
            align_str = _ALIGNMENT_MAP.get(int(para.alignment), "left")
            if role == "title":
                acc.title_alignments[align_str] += 1
            elif role == "subtitle":
                acc.subtitle_alignments[align_str] += 1
            else:
                acc.body_alignments[align_str] += 1
    except Exception:
        pass

    # Line spacing
    try:
        if para.line_spacing is not None:
            ls = float(para.line_spacing)
            if role == "title":
                acc.title_line_spacings[ls] += 1
            else:
                acc.body_line_spacings[ls] += 1
    except Exception:
        pass

    # Space after / before (in EMU — 914400 EMU = 1 inch, 12700 EMU = 1 pt)
    try:
        if para.space_after is not None:
            acc.space_afters[int(para.space_after)] += 1
    except Exception:
        pass
    try:
        if para.space_before is not None:
            acc.space_befores[int(para.space_before)] += 1
    except Exception:
        pass


# -----------------------------------------------------------------------
# Font weight categorization
# -----------------------------------------------------------------------

_BOLD_KEYWORDS = {"bold", "semibold", "semi bold", "black", "heavy", "extrabold"}
_LIGHT_KEYWORDS = {"light", "extralight", "extra light", "thin", "ultralight"}


def _categorize_font_weight(font_name: str, acc: _Accum) -> None:
    """Categorize a full font name into emphasis/light weight buckets."""
    name_lower = font_name.lower()
    for kw in _BOLD_KEYWORDS:
        if kw in name_lower:
            acc.emphasis_fonts[font_name] += 1
            return
    for kw in _LIGHT_KEYWORDS:
        if kw in name_lower:
            acc.light_fonts[font_name] += 1
            return


# -----------------------------------------------------------------------
# Shape fill extraction
# -----------------------------------------------------------------------

def _extract_shape_fill(shape, acc: _Accum) -> None:
    """Extract fill colors from a shape as accent color candidates."""
    try:
        if shape.fill and shape.fill.type is not None:
            fg = shape.fill.fore_color
            if fg and fg.rgb:
                hex_color = f"#{fg.rgb}"
                if _normalize(hex_color) not in _NEUTRAL_COLORS:
                    acc.accent_colors[hex_color] += 1
    except Exception:
        pass


# -----------------------------------------------------------------------
# Background / theme extraction
# -----------------------------------------------------------------------

def _extract_slide_background_color(slide, bg_colors: Counter) -> None:
    """Extract solid background color from a slide's background fill."""
    try:
        bg = slide.background
        if bg is None or bg._element is None:
            return

        bg_elem = bg._element
        ns_a = _NS["a"]

        for bgPr in bg_elem.iter(f"{{{ns_a}}}bgPr"):
            for solidFill in bgPr.iter(f"{{{ns_a}}}solidFill"):
                hex_color = _extract_color_from_fill_element(solidFill)
                if hex_color:
                    bg_colors[hex_color] += 1
                return

        for solidFill in bg_elem.iter(f"{{{ns_a}}}solidFill"):
            hex_color = _extract_color_from_fill_element(solidFill)
            if hex_color:
                bg_colors[hex_color] += 1

    except Exception as e:
        logger.debug(f"Could not extract slide background color: {e}")


def _extract_theme_colors(prs: Presentation, theme_colors: Counter) -> None:
    """Extract named theme colors from the first slide master's theme XML."""
    try:
        master = prs.slide_masters[0]
        theme_elem = master.element
        ns_a = _NS["a"]

        for clrScheme in theme_elem.iter(f"{{{ns_a}}}clrScheme"):
            for child in clrScheme:
                for srgbClr in child.iter(f"{{{ns_a}}}srgbClr"):
                    val = srgbClr.get("val")
                    if val:
                        theme_colors[f"#{val}"] += 1
                for sysClr in child.iter(f"{{{ns_a}}}sysClr"):
                    last_clr = sysClr.get("lastClr")
                    if last_clr:
                        theme_colors[f"#{last_clr}"] += 1
    except Exception as e:
        logger.debug(f"Could not extract theme colors: {e}")

    try:
        master = prs.slide_masters[0]
        theme_part = master.part.related_parts.get(
            next(
                (rid for rid, rel in master.part.rels.items()
                 if "theme" in rel.reltype.lower()),
                None
            )
        )
        if theme_part:
            theme_xml = etree.fromstring(theme_part.blob)
            ns_a = _NS["a"]
            for clrScheme in theme_xml.iter(f"{{{ns_a}}}clrScheme"):
                for child in clrScheme:
                    for srgbClr in child.iter(f"{{{ns_a}}}srgbClr"):
                        val = srgbClr.get("val")
                        if val:
                            theme_colors[f"#{val}"] += 1
    except Exception as e:
        logger.debug(f"Could not extract theme colors from theme part: {e}")


def _extract_color_from_fill_element(fill_elem) -> str | None:
    """Extract hex color from a fill XML element."""
    ns_a = _NS["a"]
    for srgb in fill_elem.iter(f"{{{ns_a}}}srgbClr"):
        val = srgb.get("val")
        if val:
            return f"#{val}"
    return None


# -----------------------------------------------------------------------
# Shape context classification
# -----------------------------------------------------------------------

def _classify_shape_context(shape) -> str:
    """Classify a shape as 'title', 'subtitle', or 'body'."""
    if shape.is_placeholder:
        try:
            ph_type = shape.placeholder_format.type
            if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                return "title"
            elif ph_type == PP_PLACEHOLDER.SUBTITLE:
                return "subtitle"
        except Exception:
            pass

    name_lower = shape.name.lower() if shape.name else ""
    if "title" in name_lower:
        return "title"
    elif "subtitle" in name_lower:
        return "subtitle"

    try:
        if shape.has_text_frame and shape.top is not None and shape.height is not None:
            top_inches = shape.top / 914400
            height_inches = shape.height / 914400
            if top_inches < 1.5 and height_inches > 0.6:
                if shape.text_frame.paragraphs:
                    for run in shape.text_frame.paragraphs[0].runs:
                        if run.font.size and run.font.size.pt >= 28:
                            return "title"
    except Exception:
        pass

    return "body"


# -----------------------------------------------------------------------
# Derivation — convert raw observations into DesignSystem components
# -----------------------------------------------------------------------

def _derive_fonts(acc: _Accum) -> FontConfig:
    """Pick the most common fonts, sizes, and weight variants."""
    title_font = _most_common_or_default(acc.title_fonts, "Arial")
    body_font = _most_common_or_default(acc.body_fonts, "Arial")

    title_size = _most_common_or_default(acc.title_sizes, 44)
    subtitle_size = _most_common_or_default(acc.subtitle_sizes, 28)
    body_size = _most_common_or_default(acc.body_sizes, 18)
    bullet_size = max(12, body_size - 2)

    # Weight variants
    emphasis_font = _most_common_or_default(acc.emphasis_fonts, None)
    light_font = _most_common_or_default(acc.light_fonts, None)

    # Quote font: prefer light font variant, or body font
    quote_font = light_font

    # Quote size: look for sizes near subtitle range used with italic
    quote_size = subtitle_size

    # Data point size: larger than title
    data_point_size = int(title_size * 1.5) if title_size else 54

    # Caption size
    caption_size = max(10, bullet_size - 2) if bullet_size else 12

    return FontConfig(
        title_font=title_font,
        body_font=body_font,
        title_size=title_size,
        subtitle_size=subtitle_size,
        body_size=body_size,
        bullet_size=bullet_size,
        emphasis_font=emphasis_font,
        light_font=light_font,
        quote_font=quote_font,
        quote_size=quote_size,
        data_point_size=data_point_size,
        caption_size=caption_size,
    )


def _derive_colors(acc: _Accum) -> ColorConfig:
    """Pick the dominant brand colors from frequency data."""
    filtered_accents = Counter({
        k: v for k, v in acc.accent_colors.items()
        if _normalize(k) not in _NEUTRAL_COLORS
    })
    filtered_theme = Counter({
        k: v for k, v in acc.theme_colors.items()
        if _normalize(k) not in _NEUTRAL_COLORS
    })
    combined = filtered_accents + filtered_theme
    top = combined.most_common(3)

    primary = top[0][0] if len(top) > 0 else "#1a73e8"
    secondary = top[1][0] if len(top) > 1 else "#34a853"
    accent = top[2][0] if len(top) > 2 else "#ea4335"

    text_dark = _most_common_or_default(acc.text_dark_colors, "#202124")
    text_light = _most_common_or_default(acc.text_light_colors, "#5f6368")
    background = _most_common_or_default(acc.bg_colors, "#ffffff")

    return ColorConfig(
        primary=primary,
        secondary=secondary,
        accent=accent,
        text_dark=text_dark,
        text_light=text_light,
        background=background,
    )


def _derive_paragraph(acc: _Accum) -> ParagraphConfig:
    """Derive paragraph formatting defaults."""
    title_align = _most_common_or_default(acc.title_alignments, "left")
    body_align = _most_common_or_default(acc.body_alignments, "left")
    subtitle_align = _most_common_or_default(acc.subtitle_alignments, "left")

    body_ls = _most_common_or_default(acc.body_line_spacings, None)
    title_ls = _most_common_or_default(acc.title_line_spacings, None)

    # Convert space after/before from EMU to points (12700 EMU = 1 pt)
    space_after_emu = _most_common_or_default(acc.space_afters, None)
    space_before_emu = _most_common_or_default(acc.space_befores, None)

    space_after_pt = round(space_after_emu / 12700) if space_after_emu else None
    space_before_pt = round(space_before_emu / 12700) if space_before_emu else None

    return ParagraphConfig(
        title_alignment=title_align,
        body_alignment=body_align,
        subtitle_alignment=subtitle_align,
        body_line_spacing=body_ls,
        title_line_spacing=title_ls,
        space_after_body=space_after_pt,
        space_before_body=space_before_pt,
    )


def _derive_content_area(acc: _Accum) -> ContentAreaConfig:
    """Derive content area positioning from placeholder positions."""
    cfg = ContentAreaConfig()

    if acc.title_positions:
        positions = acc.title_positions
        cfg.title_left = round(sum(p[0] for p in positions) / len(positions), 2)
        cfg.title_top = round(sum(p[1] for p in positions) / len(positions), 2)
        cfg.title_width = round(sum(p[2] for p in positions) / len(positions), 2)
        cfg.title_height = round(sum(p[3] for p in positions) / len(positions), 2)

        # Derive margins from title position
        cfg.margin_left = round(min(p[0] for p in positions), 2)
        cfg.margin_top = round(min(p[1] for p in positions), 2)

    if acc.body_positions:
        positions = acc.body_positions
        cfg.body_left = round(sum(p[0] for p in positions) / len(positions), 2)
        cfg.body_top = round(sum(p[1] for p in positions) / len(positions), 2)
        cfg.body_width = round(sum(p[2] for p in positions) / len(positions), 2)
        cfg.body_height = round(sum(p[3] for p in positions) / len(positions), 2)

    if acc.slide_width and (acc.title_positions or acc.body_positions):
        # Right margin = slide width - (left + width) of rightmost content
        all_rights = []
        for p in acc.title_positions + acc.body_positions:
            all_rights.append(p[0] + p[2])
        if all_rights:
            cfg.margin_right = round(acc.slide_width - max(all_rights), 2)

    return cfg


def _derive_overrides(acc: _Accum) -> SlideTypeOverrides:
    """Derive per-slide-type overrides from frequency data."""
    dark_bgs = Counter({
        k: v for k, v in acc.bg_colors.items()
        if _is_dark_color(k)
    })
    section_header_bg = _most_common_or_default(dark_bgs, None)
    dark_text = _most_common_or_default(acc.text_light_colors, None)

    filtered = Counter({
        k: v for k, v in acc.accent_colors.items()
        if _normalize(k) not in _NEUTRAL_COLORS
    })
    data_accent = filtered.most_common(1)[0][0] if filtered else None

    return SlideTypeOverrides(
        section_header_bg=section_header_bg,
        data_point_accent=data_accent,
        dark_slide_text=dark_text,
        quote_bg=section_header_bg,
    )


# -----------------------------------------------------------------------
# Utility
# -----------------------------------------------------------------------

def _most_common_or_default(counter: Counter, default):
    if counter:
        return counter.most_common(1)[0][0]
    return default


def _normalize(hex_color: str) -> str:
    return hex_color.lower()


def _is_dark_color(hex_color: str) -> bool:
    hex_color = hex_color.lstrip("#")
    try:
        r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
        return (0.299 * r + 0.587 * g + 0.114 * b) < 128
    except (ValueError, IndexError):
        return True
