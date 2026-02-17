"""Text manipulation operations for PowerPoint slides.

Provides primitives for text boxes, bullet lists, labels, hero numbers,
multi-format text boxes with independently styled runs, and decoration
shapes (accent bars / stripes).
"""

import logging
from typing import Any

from lxml import etree
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)

# Namespace for DrawingML elements
_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


# ---------------------------------------------------------------------------
# Core text box helpers
# ---------------------------------------------------------------------------


def add_textbox(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_name: str = "Arial",
    font_size: int = 18,
    font_color: str = "#202124",
    bold: bool = False,
    italic: bool = False,
    alignment: str = "left",
    word_wrap: bool = True,
    vertical_anchor: str = "top",
    line_spacing: float | None = None,
) -> object:
    """Add a text box to a slide.

    Args:
        slide: The slide to add the text box to.
        text: The text content.
        left, top, width, height: Position and size in inches.
        font_name: Font family name.
        font_size: Font size in points.
        font_color: Hex color string (e.g., "#202124").
        bold: Whether text should be bold.
        italic: Whether text should be italic.
        alignment: Text alignment ("left", "center", "right").
        word_wrap: Whether to enable word wrapping.
        vertical_anchor: Vertical text position ("top", "middle", "bottom").
        line_spacing: Optional line spacing multiplier.

    Returns:
        The created text box shape.
    """
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap

    # Zero out internal margins for pixel-accurate positioning
    tf.margin_left = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_bottom = Inches(0)

    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass

    # Set vertical anchor via XML
    if vertical_anchor in ("middle", "bottom"):
        try:
            from lxml import etree
            txBody = txBox._element.find(
                ".//{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr"
            )
            if txBody is not None:
                anchor_map = {"top": "t", "middle": "ctr", "bottom": "b"}
                txBody.set("anchor", anchor_map.get(vertical_anchor, "t"))
        except Exception:
            pass

    p = tf.paragraphs[0]
    p.text = text
    _apply_run_format(p, font_name, font_size, font_color, bold, italic)
    p.alignment = _get_alignment(alignment)

    if line_spacing is not None:
        try:
            p.line_spacing = line_spacing
        except Exception:
            pass

    return txBox


def add_bullet_list(
    slide,
    items: list[str],
    left: float,
    top: float,
    width: float,
    height: float,
    font_name: str = "Arial",
    font_size: int = 16,
    font_color: str = "#202124",
    bullet_color: str | None = None,
    line_spacing: float = 1.2,
) -> object:
    """Add a bullet list to a slide.

    Args:
        slide: The slide to add bullets to.
        items: List of bullet point strings.
        left, top, width, height: Position and size in inches.
        font_name: Font family name.
        font_size: Font size in points.
        font_color: Hex color string for text.
        bullet_color: Hex color for bullets (defaults to font_color).
        line_spacing: Line spacing multiplier.

    Returns:
        The created text box shape.
    """
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    # Zero out internal margins for pixel-accurate positioning
    tf.margin_left = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_bottom = Inches(0)

    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = item
        p.level = 0
        _apply_run_format(p, font_name, font_size, font_color, bold=False, italic=False)
        p.alignment = PP_ALIGN.LEFT

        # Explicitly set bullet character via XML (don't rely on layout defaults)
        _set_bullet_char(p, "\u2022", font_color)

        # Set bullet formatting
        p.space_after = Pt(font_size * 0.4)

        # Set line spacing
        try:
            p.line_spacing = line_spacing
        except Exception:
            pass

    return txBox


# ---------------------------------------------------------------------------
# Enhanced text helpers
# ---------------------------------------------------------------------------


def add_label(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float = 0.25,
    font_name: str = "Poppins Medium",
    font_size: int = 9,
    font_color: str = "#434343",
    alignment: str = "left",
    bold: bool = False,
) -> object:
    """Add a small label text (section markers, dates, confidential notices).

    Optimized for 8-10pt supporting text that appears throughout
    the base template as contextual information.

    Args:
        slide: Target slide.
        text: Label text (e.g., '01 | Revenue growth', 'Confidential').
        left, top: Position in inches.
        width: Text area width in inches.
        height: Text area height in inches.
        font_name: Label font (typically Poppins Medium).
        font_size: Font size in points (typically 8-10).
        font_color: Text color.
        alignment: Text alignment.
        bold: Whether text should be bold.

    Returns:
        The created text box shape.
    """
    return add_textbox(
        slide,
        text,
        left=left,
        top=top,
        width=width,
        height=height,
        font_name=font_name,
        font_size=font_size,
        font_color=font_color,
        bold=bold,
        alignment=alignment,
    )


def add_hero_number(
    slide,
    number_text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_name: str = "Inter Tight ExtraLight",
    font_size: int = 88,
    font_color: str = "#CDF986",
    alignment: str = "left",
) -> object:
    """Add a large display/hero number with ExtraLight font.

    Used for data point slides where a single statistic is the
    visual centerpiece (e.g., '47%', '$2.3M', '10x').

    Args:
        slide: Target slide.
        number_text: The display number/value string.
        left, top, width, height: Position and size in inches.
        font_name: Display font (typically ExtraLight weight).
        font_size: Large font size in points (80-100).
        font_color: Number color (typically accent).
        alignment: Text alignment.

    Returns:
        The created text box shape.
    """
    return add_textbox(
        slide,
        number_text,
        left=left,
        top=top,
        width=width,
        height=height,
        font_name=font_name,
        font_size=font_size,
        font_color=font_color,
        bold=False,
        alignment=alignment,
        vertical_anchor="middle",
    )


def add_multi_format_textbox(
    slide,
    runs: list[dict[str, Any]],
    left: float,
    top: float,
    width: float,
    height: float,
    alignment: str = "left",
    line_spacing: float | None = None,
) -> object:
    """Add a text box with multiple independently formatted runs.

    Enables patterns like '01 | Revenue growth' where '01' uses one
    style and 'Revenue growth' uses another.

    Args:
        slide: Target slide.
        runs: List of run specifications, each a dict with:
            - text (str): The text content.
            - font_name (str, optional): Font family.
            - font_size (int, optional): Size in points.
            - font_color (str, optional): Hex color.
            - bold (bool, optional): Bold flag.
            - italic (bool, optional): Italic flag.
        left, top, width, height: Position and size in inches.
        alignment: Text alignment for the paragraph.
        line_spacing: Optional line spacing multiplier.

    Returns:
        The created text box shape.
    """
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass

    p = tf.paragraphs[0]
    p.alignment = _get_alignment(alignment)

    if line_spacing is not None:
        try:
            p.line_spacing = line_spacing
        except Exception:
            pass

    # Clear default text
    p.text = ""

    for i, run_spec in enumerate(runs):
        run = p.add_run()
        run.text = run_spec.get("text", "")

        if "font_name" in run_spec:
            run.font.name = run_spec["font_name"]
        if "font_size" in run_spec:
            run.font.size = Pt(run_spec["font_size"])
        if "font_color" in run_spec:
            run.font.color.rgb = _hex_to_rgb(run_spec["font_color"])
        if "bold" in run_spec:
            run.font.bold = run_spec["bold"]
        if "italic" in run_spec:
            run.font.italic = run_spec["italic"]

    return txBox


# ---------------------------------------------------------------------------
# Placeholder operations
# ---------------------------------------------------------------------------


def set_placeholder_text(
    slide,
    placeholder_name: str,
    text: str,
    font_name: str | None = None,
    font_size: int | None = None,
    font_color: str | None = None,
    bold: bool | None = None,
) -> bool:
    """Set text in a named placeholder on a slide.

    Args:
        slide: The slide containing the placeholder.
        placeholder_name: Name or partial name of the placeholder shape.
        text: Text to set.
        font_name, font_size, font_color, bold: Optional formatting overrides.

    Returns:
        True if placeholder was found and updated, False otherwise.
    """
    name_lower = placeholder_name.lower()

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        shape_name = shape.name.lower()
        if name_lower in shape_name or shape_name in name_lower:
            tf = shape.text_frame
            if tf.paragraphs:
                p = tf.paragraphs[0]
                p.text = text
                if font_name or font_size or font_color or bold is not None:
                    _apply_run_format(
                        p,
                        font_name or "Arial",
                        font_size or 18,
                        font_color or "#202124",
                        bold=bold or False,
                        italic=False,
                    )
            return True

    logger.warning(f"Placeholder '{placeholder_name}' not found on slide")
    return False


def populate_slide_text(
    slide,
    title: str | None = None,
    subtitle: str | None = None,
    body: str | None = None,
    font_config: dict | None = None,
) -> None:
    """Populate standard text placeholders on a slide.

    Searches for title/subtitle/body placeholders by name and type,
    then sets the text content with appropriate formatting.
    """
    fc = font_config or {}

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        name_lower = shape.name.lower()
        is_placeholder = shape.is_placeholder

        if is_placeholder:
            try:
                from pptx.enum.shapes import PP_PLACEHOLDER

                ph_type = shape.placeholder_format.type
                if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                    if title:
                        _set_shape_text(shape, title, fc.get("title_font"), fc.get("title_size"))
                    continue
                elif ph_type == PP_PLACEHOLDER.SUBTITLE:
                    if subtitle:
                        _set_shape_text(
                            shape, subtitle, fc.get("body_font"), fc.get("subtitle_size")
                        )
                    continue
                elif ph_type == PP_PLACEHOLDER.BODY:
                    if body:
                        _set_shape_text(shape, body, fc.get("body_font"), fc.get("body_size"))
                    continue
            except Exception:
                pass

        # Fallback: match by shape name
        if "title" in name_lower and title:
            _set_shape_text(shape, title, fc.get("title_font"), fc.get("title_size"))
        elif "subtitle" in name_lower and subtitle:
            _set_shape_text(shape, subtitle, fc.get("body_font"), fc.get("subtitle_size"))
        elif "body" in name_lower or "content" in name_lower or "text" in name_lower:
            if body:
                _set_shape_text(shape, body, fc.get("body_font"), fc.get("body_size"))


# ---------------------------------------------------------------------------
# Decoration shapes
# ---------------------------------------------------------------------------


def add_accent_bar(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    color_hex: str,
) -> object:
    """Add a thin colored rectangle shape (accent stripe / bar).

    Used to replicate CSS border-left accent stripes in the PPTX.

    Args:
        slide: Target slide.
        left, top, width, height: Position and size in inches.
        color_hex: Fill color as hex string (e.g., '#CDF986').

    Returns:
        The created shape.
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(color_hex)
    # Remove the shape's own outline
    shape.line.fill.background()
    return shape


def add_image_placeholder(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    description: str = "",
    fill_hex: str = "#F4F4F4",
    text_color_hex: str = "#888888",
    border_hex: str = "#CCCCCC",
    image_style: str = "",
) -> object:
    """Add a labeled placeholder rectangle for an image or illustration.

    Creates a rectangle shape with a dashed border, light fill, and centered
    label text describing what image should go there. The presenter (or a
    future image pipeline) can replace this placeholder with actual imagery.

    Args:
        slide: Target slide.
        left, top, width, height: Position and size in inches.
        description: Label text describing the intended image.
        fill_hex: Background fill color (default light gray).
        text_color_hex: Label text color (default medium gray).
        border_hex: Dashed border color (default gray).
        image_style: Image style hint (e.g., 'diagram', 'photo').

    Returns:
        The created shape.
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )

    # Light fill
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(fill_hex)

    # Dashed border
    shape.line.color.rgb = _hex_to_rgb(border_hex)
    shape.line.width = Pt(1.5)
    shape.line.dash_style = 4  # MSO_LINE_DASH_STYLE.DASH (enum value 4)

    # Label text centered in the shape
    label = description
    if image_style and image_style not in description.lower():
        label = f"[{image_style}] {description}"

    tf = shape.text_frame
    tf.word_wrap = True
    if tf.paragraphs:
        p = tf.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.size = Pt(10)
            run.font.color.rgb = _hex_to_rgb(text_color_hex)
            run.font.name = "Arial"
            run.font.italic = True

    return shape


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------


def _set_bullet_char(paragraph, char: str = "\u2022", color_hex: str | None = None) -> None:
    """Explicitly set a bullet character on a paragraph using XML.

    This ensures bullets appear regardless of slide layout defaults.
    """
    pPr = paragraph._p.get_or_add_pPr()

    # Remove any existing buNone (explicit "no bullet") element
    for existing in pPr.findall(f"{{{_NS_A}}}buNone"):
        pPr.remove(existing)
    # Remove any existing buChar to avoid duplicates
    for existing in pPr.findall(f"{{{_NS_A}}}buChar"):
        pPr.remove(existing)
    # Remove any existing buAutoNum
    for existing in pPr.findall(f"{{{_NS_A}}}buAutoNum"):
        pPr.remove(existing)

    # Set bullet color if specified
    if color_hex:
        # Remove existing buClr
        for existing in pPr.findall(f"{{{_NS_A}}}buClr"):
            pPr.remove(existing)
        buClr = etree.SubElement(pPr, f"{{{_NS_A}}}buClr")
        srgbClr = etree.SubElement(buClr, f"{{{_NS_A}}}srgbClr")
        srgbClr.set("val", color_hex.lstrip("#")[:6])

    # Set bullet size relative to text
    for existing in pPr.findall(f"{{{_NS_A}}}buSzPct"):
        pPr.remove(existing)
    buSzPct = etree.SubElement(pPr, f"{{{_NS_A}}}buSzPct")
    buSzPct.set("val", "100000")  # 100% of text size

    # Set the actual bullet character
    buChar = etree.SubElement(pPr, f"{{{_NS_A}}}buChar")
    buChar.set("char", char)


def _set_shape_text(
    shape,
    text: str,
    font_name: str | None = None,
    font_size: int | None = None,
) -> None:
    """Set text on a shape, preserving existing formatting where possible."""
    tf = shape.text_frame
    if tf.paragraphs:
        # Preserve formatting from first run if it exists
        p = tf.paragraphs[0]
        existing_font_name = None
        existing_font_size = None
        if p.runs:
            run = p.runs[0]
            existing_font_name = run.font.name
            if run.font.size:
                existing_font_size = int(run.font.size.pt)

        p.text = text
        if p.runs:
            run = p.runs[0]
            if font_name:
                run.font.name = font_name
            elif existing_font_name:
                run.font.name = existing_font_name
            if font_size:
                run.font.size = Pt(font_size)
            elif existing_font_size:
                run.font.size = Pt(existing_font_size)


def _apply_run_format(
    paragraph,
    font_name: str,
    font_size: int,
    font_color: str,
    bold: bool,
    italic: bool,
) -> None:
    """Apply font formatting to all runs in a paragraph."""
    color_rgb = _hex_to_rgb(font_color)
    for run in paragraph.runs:
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color_rgb


def _get_alignment(alignment: str) -> int:
    """Convert alignment string to PP_ALIGN constant."""
    align_map = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }
    return align_map.get(alignment.lower(), PP_ALIGN.LEFT)


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor."""
    hex_color = hex_color.lstrip("#")
    if len(hex_color) > 6:
        hex_color = hex_color[:6]
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )


# ---------------------------------------------------------------------------
# Font sizing utilities
# ---------------------------------------------------------------------------

def estimate_fit_font_size(
    text: str,
    shape_width_inches: float,
    shape_height_inches: float,
    max_font_pt: float = 44.0,
    min_font_pt: float = 10.0,
) -> float:
    """Estimate the largest font size (pt) at which text fits a shape.

    Uses a simple character-per-line heuristic (not pixel-perfect, but
    prevents gross overflows). Assumes roughly 1.0 characters per point
    of width, and 1.3x font size per line height.

    Args:
        text: The text to fit.
        shape_width_inches: Available width in inches.
        shape_height_inches: Available height in inches.
        max_font_pt: Maximum font size to consider.
        min_font_pt: Minimum font size (floor).

    Returns:
        Recommended font size in points.
    """
    if not text or shape_width_inches <= 0 or shape_height_inches <= 0:
        return max_font_pt

    # Try sizes from max down to min
    for size_pt in range(int(max_font_pt), int(min_font_pt) - 1, -1):
        # Approximate chars per line at this font size
        # At 12pt, roughly 9 chars per inch of width; scales inversely
        chars_per_inch = max(1.0, 9.0 * (12.0 / size_pt))
        chars_per_line = int(shape_width_inches * chars_per_inch)
        if chars_per_line < 1:
            continue

        # Line height ≈ 1.3x font size in inches
        line_height_inches = size_pt / 72.0 * 1.3
        max_lines = max(1, int(shape_height_inches / line_height_inches))

        # Estimate lines needed (word-wrap simulation)
        lines_needed = 0
        for paragraph in text.split("\n"):
            if not paragraph.strip():
                lines_needed += 1
                continue
            lines_needed += max(1, -(-len(paragraph) // chars_per_line))  # ceil division

        if lines_needed <= max_lines:
            return float(size_pt)

    return min_font_pt
