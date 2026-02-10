"""Image manipulation operations for PowerPoint slides."""

import logging
from pathlib import Path

from pptx.util import Inches

logger = logging.getLogger(__name__)


def add_image(
    slide,
    image_path: str | Path,
    left: float,
    top: float,
    width: float | None = None,
    height: float | None = None,
) -> object | None:
    """Add an image to a slide.

    Args:
        slide: The slide to add the image to.
        image_path: Path to the image file.
        left, top: Position in inches.
        width: Width in inches (None for original).
        height: Height in inches (None for original).

    Returns:
        The created picture shape, or None if the image couldn't be added.
    """
    image_path = Path(image_path)
    if not image_path.exists():
        logger.warning(f"Image not found: {image_path}")
        return None

    kwargs = {
        "image_file": str(image_path),
        "left": Inches(left),
        "top": Inches(top),
    }

    if width is not None:
        kwargs["width"] = Inches(width)
    if height is not None:
        kwargs["height"] = Inches(height)

    try:
        return slide.shapes.add_picture(**kwargs)
    except Exception as e:
        logger.warning(f"Could not add image {image_path}: {e}")
        return None


def add_image_centered(
    slide,
    image_path: str | Path,
    slide_width: float = 13.333,
    slide_height: float = 7.5,
    max_width: float | None = None,
    max_height: float | None = None,
) -> object | None:
    """Add an image centered on the slide.

    The image is scaled to fit within the specified maximum dimensions
    while maintaining aspect ratio.
    """
    image_path = Path(image_path)
    if not image_path.exists():
        logger.warning(f"Image not found: {image_path}")
        return None

    # Use max dimensions or slide dimensions with margin
    max_w = max_width or (slide_width * 0.8)
    max_h = max_height or (slide_height * 0.7)

    try:
        # Get original image dimensions for aspect ratio
        from PIL import Image

        with Image.open(image_path) as img:
            orig_w, orig_h = img.size
    except ImportError:
        # Without PIL, just use max dimensions
        width = max_w
        height = max_h
        left = (slide_width - width) / 2
        top = (slide_height - height) / 2
        return add_image(slide, image_path, left, top, width, height)
    except Exception as e:
        logger.warning(f"Could not read image dimensions: {e}")
        width = max_w
        height = max_h
        left = (slide_width - width) / 2
        top = (slide_height - height) / 2
        return add_image(slide, image_path, left, top, width, height)

    # Scale to fit within max dimensions
    aspect = orig_w / orig_h
    if aspect > (max_w / max_h):
        # Width-constrained
        width = max_w
        height = width / aspect
    else:
        # Height-constrained
        height = max_h
        width = height * aspect

    left = (slide_width - width) / 2
    top = (slide_height - height) / 2

    return add_image(slide, image_path, left, top, width, height)


def replace_placeholder_image(
    slide,
    placeholder_name: str,
    image_path: str | Path,
) -> bool:
    """Replace a picture placeholder with an image.

    Args:
        slide: The slide containing the placeholder.
        placeholder_name: Name or partial name of the placeholder.
        image_path: Path to the replacement image.

    Returns:
        True if placeholder was found and image inserted, False otherwise.
    """
    image_path = Path(image_path)
    if not image_path.exists():
        logger.warning(f"Image not found: {image_path}")
        return False

    name_lower = placeholder_name.lower()

    for shape in slide.shapes:
        if name_lower in shape.name.lower():
            # Get the placeholder position and size
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height

            try:
                slide.shapes.add_picture(
                    str(image_path), left, top, width, height
                )
                return True
            except Exception as e:
                logger.warning(f"Could not replace placeholder image: {e}")
                return False

    logger.warning(f"Placeholder '{placeholder_name}' not found on slide")
    return False
