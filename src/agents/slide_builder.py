"""Slide Builder Agent — Clone-and-Replace Architecture.

Primary mode: Clone a matched template slide (preserving all visual elements —
backgrounds, images, decorations, branded graphics) and replace only the text
content zones with deck-schema content.

Fallback mode: When no template match exists, use a layout-based slide from
the base template (inheriting master backgrounds) and compose content via the
type-specific composer system.
"""

import logging
from pathlib import Path

from pptx import Presentation

from src.schemas.design_system import DesignSystem
from src.schemas.slide_schema import DeckSchema, SlideSpec, SlideType
from src.pptx_engine.composers import get_composer
from src.pptx_engine.slide_operations import (
    add_blank_slide,
    add_slide_from_layout,
    clone_slide_as_is,
    open_base_template,
    create_presentation,
)
from src.pptx_engine.text_operations import estimate_fit_font_size

logger = logging.getLogger(__name__)


class SlideBuilderAgent:
    """Build a .pptx using clone-and-replace as primary mode.

    For every slide:
    1. If a template match exists → clone the template slide, replace text
    2. If no match → create from branded layout + composer (fallback)
    """

    DEFAULT_BASE = Path("templates/base/Shopify - Example Technical Workshop Slide Bank.pptx")

    def build(
        self,
        deck_schema: DeckSchema,
        design_system: DesignSystem,
        output_path: str | Path,
        matches: list[dict] | None = None,
        template_registry: object | None = None,
        base_template: str | Path | None = None,
    ) -> Path:
        output_path = Path(output_path)

        base_path = Path(base_template) if base_template else self.DEFAULT_BASE
        if base_path.exists():
            prs = open_base_template(base_path)
        else:
            logger.warning(f"Base template not found at {base_path}, creating blank")
            prs = create_presentation()

        match_lookup = {}
        if matches:
            for m in matches:
                match_lookup[m["slide_number"]] = m

        cloned = 0
        composed = 0
        for slide_spec in deck_schema.slides:
            match_info = match_lookup.get(slide_spec.slide_number)

            # --- Primary mode: clone template + replace content ---
            if match_info and match_info.get("match_type") == "use_as_is" and template_registry:
                try:
                    template = template_registry.templates[match_info["template_index"]]
                    slide = clone_slide_as_is(
                        prs, template.template_file, template.slide_index
                    )
                    # Replace text using content zones if available, else heuristic
                    content_zones = getattr(template, "content_zones", None)
                    if content_zones:
                        self._populate_with_zones(slide, slide_spec, content_zones)
                    else:
                        self._populate_cloned_slide(slide, slide_spec, design_system)

                    if slide_spec.speaker_notes:
                        self._add_speaker_notes(slide, slide_spec.speaker_notes)
                    cloned += 1
                    logger.info(
                        f"Slide {slide_spec.slide_number}: cloned from "
                        f"{template.template_file} index {template.slide_index}"
                    )
                    continue
                except Exception as e:
                    logger.warning(
                        f"Slide {slide_spec.slide_number}: clone failed ({e}), "
                        f"falling back to layout-based compose"
                    )

            # --- Fallback mode: branded layout + composer ---
            self._build_composed_slide(prs, slide_spec, design_system)
            composed += 1

        prs.save(str(output_path))
        logger.info(
            f"Saved: {output_path} ({cloned} cloned, {composed} composed, "
            f"{len(deck_schema.slides)} total)"
        )
        return output_path

    # ------------------------------------------------------------------
    # Composed slide generation (fallback)
    # ------------------------------------------------------------------

    def _build_composed_slide(
        self,
        prs: Presentation,
        spec: SlideSpec,
        design: DesignSystem,
    ) -> None:
        """Build a slide using the composition system.

        Creates a slide from the appropriate branded layout (inheriting
        master backgrounds and placeholders) and delegates to the
        type-specific composer for content placement.
        """
        slide = add_slide_from_layout(prs, spec.slide_type.value)
        composer = get_composer(spec.slide_type)
        composer.compose(slide, spec, design)

        if spec.speaker_notes:
            self._add_speaker_notes(slide, spec.speaker_notes)

    # ------------------------------------------------------------------
    # Content zone-based text replacement
    # ------------------------------------------------------------------

    def _populate_with_zones(
        self,
        slide,
        spec: SlideSpec,
        content_zones: list,
    ) -> None:
        """Replace text in cloned slide using content zone map.

        Content zones precisely identify which shapes hold replaceable text
        and what type of content they expect (title, body, subtitle, etc.).
        After mapping, ALL unmapped text shapes have their text cleared.
        """
        title_text = spec.title or ""
        body_text = self._get_combined_body(spec)
        subtitle_text = spec.subtitle or ""
        data_text = self._get_data_point_text(spec)

        # Build a shape name → shape lookup
        shape_lookup = {}
        for shape in slide.shapes:
            shape_lookup[shape.name] = shape

        mapped_names: set[str] = set()

        for zone in content_zones:
            zone_type = zone.zone_type if hasattr(zone, "zone_type") else zone.get("zone_type", "body")
            shape_name = zone.shape_name if hasattr(zone, "shape_name") else zone.get("shape_name", "")

            shape = shape_lookup.get(shape_name)
            if not shape or not shape.has_text_frame:
                continue

            # Extract font size range from zone metadata
            font_range = None
            if hasattr(zone, "font_size_range"):
                font_range = zone.font_size_range
            elif isinstance(zone, dict):
                font_range = zone.get("font_size_range")

            size_kwargs = {}
            if font_range and len(font_range) == 2:
                size_kwargs["min_font_pt"] = font_range[0]
                size_kwargs["max_font_pt"] = font_range[1]

            # Extract max_chars for truncation
            max_chars = None
            if hasattr(zone, "max_chars"):
                max_chars = zone.max_chars
            elif isinstance(zone, dict):
                max_chars = zone.get("max_chars")

            if zone_type == "title" and title_text:
                text = self._truncate_to_fit(title_text, max_chars)
                self._replace_shape_text(shape, text, **size_kwargs)
                mapped_names.add(shape_name)
            elif zone_type == "subtitle" and subtitle_text:
                text = self._truncate_to_fit(subtitle_text, max_chars)
                self._replace_shape_text(shape, text, **size_kwargs)
                mapped_names.add(shape_name)
            elif zone_type == "data_point" and data_text:
                self._replace_shape_text(shape, data_text, **size_kwargs)
                mapped_names.add(shape_name)
            elif zone_type in ("body", "bullet_area", "caption") and body_text:
                text = self._truncate_to_fit(body_text, max_chars)
                self._replace_shape_text(shape, text, **size_kwargs)
                mapped_names.add(shape_name)

        # Clear ALL unmapped text shapes to remove stale template text
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name not in mapped_names:
                self._clear_shape_text(shape)

    # ------------------------------------------------------------------
    # Heuristic text replacement for cloned slides
    # ------------------------------------------------------------------

    # Minimum shape area (sq inches) to be considered a content shape.
    # Smaller shapes are treated as decorative labels/badges.
    _CONTENT_AREA_THRESHOLD = 1.5

    def _populate_cloned_slide(
        self,
        slide,
        spec: SlideSpec,
        design: DesignSystem,
    ) -> None:
        """Replace ALL text on a cloned slide with deck schema content.

        Strategy:
        1. Collect every text-bearing shape with size/font metadata
        2. Separate into content shapes (large) and decorative (small)
        3. Map deck content to content shapes by role (title, body, etc.)
        4. Clear text from ALL unmapped shapes — no stale template text
        """
        title_text = spec.title or ""
        subtitle_text = spec.subtitle or ""
        body_text = self._get_combined_body(spec)
        data_text = self._get_data_point_text(spec)

        # --- Collect all text shapes with metadata ---
        all_text_shapes: list[dict] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            try:
                w = shape.width / 914400 if shape.width else 0.0
                h = shape.height / 914400 if shape.height else 0.0
                top = shape.top / 914400 if shape.top else 0.0
            except Exception:
                w, h, top = 0.0, 0.0, 0.0

            max_font = 0.0
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        max_font = max(max_font, run.font.size.pt)

            all_text_shapes.append({
                "shape": shape,
                "font_size": max_font,
                "top": top,
                "width": w,
                "area": w * h,
                "text": shape.text_frame.text.strip(),
            })

        if not all_text_shapes:
            return

        # --- Classify: content shapes vs decorative ---
        content_shapes = [
            s for s in all_text_shapes
            if s["area"] >= self._CONTENT_AREA_THRESHOLD
        ]
        # If nothing qualifies as "content", treat the largest shapes as content
        if not content_shapes:
            all_text_shapes.sort(key=lambda s: s["area"], reverse=True)
            content_shapes = all_text_shapes[:3]

        # Sort content shapes: largest font first (title candidate), then by
        # vertical position (top shapes before bottom)
        content_shapes.sort(key=lambda s: (-s["font_size"], s["top"]))

        # --- Map deck content to shapes ---
        mapped_ids: set[int] = set()

        # Design system guardrails for font size clamping
        title_max_pt = design.fonts.title_size or 44
        body_max_pt = design.fonts.body_size or 18
        subtitle_max_pt = design.fonts.subtitle_size or 28

        # Title → first content shape (largest font / topmost)
        if title_text and content_shapes:
            target = content_shapes[0]
            self._replace_shape_text(
                target["shape"], title_text,
                max_font_pt=title_max_pt, min_font_pt=14,
            )
            mapped_ids.add(id(target["shape"]))

        remaining = [s for s in content_shapes if id(s["shape"]) not in mapped_ids]

        # Subtitle → next available content shape (if subtitle exists)
        if subtitle_text and remaining:
            target = remaining[0]
            self._replace_shape_text(
                target["shape"], subtitle_text,
                max_font_pt=subtitle_max_pt, min_font_pt=12,
            )
            mapped_ids.add(id(target["shape"]))
            remaining = [s for s in remaining if id(s["shape"]) not in mapped_ids]

        # Data point → next available if we have one
        if data_text and remaining:
            target = remaining[0]
            # Data points can be large — allow up to design system's data_point size
            dp_max = design.data_point_size_resolved or 60
            self._replace_shape_text(
                target["shape"], data_text,
                max_font_pt=dp_max, min_font_pt=20,
            )
            mapped_ids.add(id(target["shape"]))
            remaining = [s for s in remaining if id(s["shape"]) not in mapped_ids]

        # Body/bullets → distribute across remaining content shapes
        if body_text and remaining:
            # If multiple body shapes, split body text across them
            body_lines = body_text.split("\n\n")
            for i, target in enumerate(remaining):
                if i < len(body_lines):
                    self._replace_shape_text(
                        target["shape"], body_lines[i],
                        max_font_pt=body_max_pt, min_font_pt=10,
                    )
                else:
                    # Ran out of content blocks — clear this shape
                    self._clear_shape_text(target["shape"])
                mapped_ids.add(id(target["shape"]))

        # --- Clear ALL unmapped text shapes ---
        for s in all_text_shapes:
            if id(s["shape"]) not in mapped_ids:
                self._clear_shape_text(s["shape"])

    # ------------------------------------------------------------------
    # Shape text operations
    # ------------------------------------------------------------------

    @staticmethod
    def _replace_shape_text(
        shape,
        new_text: str,
        *,
        min_font_pt: float | None = None,
        max_font_pt: float | None = None,
    ) -> None:
        """Replace ALL text in a shape, preserving first-run formatting.

        Clears every paragraph (not just the first), then sets the new
        text on paragraph 1. Handles theme colors gracefully.

        Optional font size guardrails:
        - min_font_pt / max_font_pt: clamp the preserved font size to a range
        - Auto-fit: if the new text is significantly longer than the shape
          can hold at the current font size, reduce the size to fit
        """
        tf = shape.text_frame
        if not tf.paragraphs:
            return

        # Preserve formatting from the first run of the first paragraph
        first_para = tf.paragraphs[0]
        font_props = {}
        if first_para.runs:
            run = first_para.runs[0]
            color_rgb = None
            try:
                if run.font.color and run.font.color.type is not None:
                    color_rgb = run.font.color.rgb
            except AttributeError:
                pass  # SchemeColor — let theme handle it

            font_props = {
                "name": run.font.name,
                "size": run.font.size,
                "bold": run.font.bold,
                "italic": run.font.italic,
                "color": color_rgb,
            }

        # --- Clamp font size to requested range ---
        if font_props.get("size"):
            from pptx.util import Pt
            current_pt = font_props["size"].pt
            clamped_pt = current_pt
            if max_font_pt is not None and clamped_pt > max_font_pt:
                clamped_pt = max_font_pt
            if min_font_pt is not None and clamped_pt < min_font_pt:
                clamped_pt = min_font_pt

            # Auto-fit: estimate if text fits at current size, reduce if needed
            try:
                shape_w = (shape.width / 914400) if shape.width else 0
                shape_h = (shape.height / 914400) if shape.height else 0
                if shape_w > 0 and shape_h > 0:
                    fit_pt = estimate_fit_font_size(
                        new_text, shape_w, shape_h,
                        max_font_pt=clamped_pt,
                        min_font_pt=min_font_pt or 10.0,
                    )
                    clamped_pt = min(clamped_pt, fit_pt)
            except Exception:
                pass  # If estimation fails, keep the clamped size

            if clamped_pt != current_pt:
                font_props["size"] = Pt(clamped_pt)

        # --- Clear ALL paragraphs beyond the first ---
        p_elements = list(tf._element)
        ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        for p_el in p_elements:
            if p_el.tag == f"{{{ns_a}}}p" and p_el is not first_para._element:
                tf._element.remove(p_el)

        # --- Set new text on the (now only) paragraph ---
        first_para.text = new_text

        # --- Reapply formatting ---
        if font_props and first_para.runs:
            for run in first_para.runs:
                if font_props.get("name"):
                    run.font.name = font_props["name"]
                if font_props.get("size"):
                    run.font.size = font_props["size"]
                if font_props.get("bold") is not None:
                    run.font.bold = font_props["bold"]
                if font_props.get("italic") is not None:
                    run.font.italic = font_props["italic"]
                # Skip re-applying explicit RGB colors — let the
                # theme/layout color inheritance work naturally.
                # Forcing a hardcoded color from the source template
                # causes white-on-white or black-on-black when the
                # layout background differs from the source.

    @staticmethod
    def _clear_shape_text(shape) -> None:
        """Remove all visible text and hide the shape to prevent visual clutter.

        Used to clear stale template text from shapes that don't receive
        deck schema content. Also moves the shape off-canvas so it doesn't
        create empty boxes or overlap with content shapes.
        """
        tf = shape.text_frame
        if not tf.paragraphs:
            return

        # Remove all paragraphs beyond the first
        ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        p_elements = list(tf._element)
        first_p = tf.paragraphs[0]._element
        for p_el in p_elements:
            if p_el.tag == f"{{{ns_a}}}p" and p_el is not first_p:
                tf._element.remove(p_el)

        # Clear the remaining paragraph
        tf.paragraphs[0].text = ""

        # Move the shape off-canvas to prevent empty box artifacts
        try:
            from pptx.util import Emu
            shape.left = Emu(914400 * 20)  # 20 inches off-screen right
        except Exception:
            pass

    # ------------------------------------------------------------------
    # Text truncation
    # ------------------------------------------------------------------

    @staticmethod
    def _truncate_to_fit(text: str, max_chars: int | None) -> str:
        """Truncate text to fit within a zone's max_chars capacity.

        Tries to break at word boundaries. Appends an ellipsis when
        truncation occurs.
        """
        if max_chars is None or max_chars <= 0 or len(text) <= max_chars:
            return text

        # Truncate at word boundary
        truncated = text[:max_chars]
        last_space = truncated.rfind(" ")
        if last_space > max_chars * 0.6:
            truncated = truncated[:last_space]

        return truncated.rstrip() + "..."

    # ------------------------------------------------------------------
    # Content extraction helpers
    # ------------------------------------------------------------------

    @staticmethod
    def _get_combined_body(spec: SlideSpec) -> str:
        """Extract body/bullet/caption text as a single string.

        Content blocks are joined with double newlines so they can be
        split apart when distributing across multiple shapes.
        """
        parts = []
        for block in spec.content_blocks:
            if block.type in ("body", "caption", "bullets"):
                parts.append(block.content)
        return "\n\n".join(parts)

    @staticmethod
    def _get_data_point_text(spec: SlideSpec) -> str:
        """Extract the data point text from content blocks."""
        for block in spec.content_blocks:
            if block.type == "data_point":
                return block.content
        return ""

    # ------------------------------------------------------------------
    # Speaker notes
    # ------------------------------------------------------------------

    @staticmethod
    def _add_speaker_notes(slide, notes: str) -> None:
        """Add speaker notes to a slide."""
        try:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes
        except Exception as e:
            logger.debug(f"Could not add speaker notes: {e}")
