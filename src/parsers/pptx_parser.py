"""Parser for PPTX files using python-pptx.

Extracts text content from existing presentations for re-designing.
"""

from pathlib import Path


def parse_pptx(path: Path) -> str:
    """Extract text from an existing PowerPoint presentation.

    Each slide's content is extracted with slide number markers.
    Text from shapes, placeholders, tables, and notes is included.
    """
    from pptx import Presentation

    prs = Presentation(str(path))
    sections: list[str] = []

    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_parts: list[str] = [f"\n--- Slide {slide_idx} ---\n"]

        for shape in slide.shapes:
            # Extract text from text frames
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue

                    # Classify based on placeholder type or shape name
                    name_lower = shape.name.lower()
                    if shape.is_placeholder:
                        ph_type = shape.placeholder_format.type
                        # Type 1 = CENTER_TITLE or TITLE, Type 2 = BODY/SUBTITLE
                        from pptx.enum.shapes import PP_PLACEHOLDER

                        if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                            slide_parts.append(f"# {text}")
                        elif ph_type == PP_PLACEHOLDER.SUBTITLE:
                            slide_parts.append(f"## {text}")
                        else:
                            slide_parts.append(text)
                    elif "title" in name_lower:
                        slide_parts.append(f"# {text}")
                    elif "subtitle" in name_lower:
                        slide_parts.append(f"## {text}")
                    else:
                        # Check if paragraph has bullet formatting
                        if para.level and para.level > 0:
                            indent = "  " * (para.level - 1)
                            slide_parts.append(f"{indent}- {text}")
                        else:
                            slide_parts.append(text)

            # Extract text from tables
            if shape.has_table:
                table = shape.table
                for row_idx, row in enumerate(table.rows):
                    cells = [cell.text.strip() for cell in row.cells]
                    slide_parts.append("| " + " | ".join(cells) + " |")
                    if row_idx == 0:
                        slide_parts.append("|" + "|".join(["---"] * len(cells)) + "|")

        # Extract speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_parts.append(f"\n> Speaker notes: {notes}")

        # Only include slides that have content
        if len(slide_parts) > 1:
            sections.append("\n".join(slide_parts))

    result = "\n".join(sections).strip()
    if not result:
        raise ValueError(f"No text content could be extracted from {path}")
    return result
